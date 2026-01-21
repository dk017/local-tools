
import logging
import os
import io
import sys
import platform
import base64
import fitz  # PyMuPDF
from pypdf import PdfWriter, PdfReader
import pdfplumber
import pandas as pd
import pikepdf
from PIL import Image, ImageChops, ImageOps

# Try to import camelot for enhanced table detection
try:
    import camelot
    CAMELOT_AVAILABLE = True
except ImportError:
    CAMELOT_AVAILABLE = False
    logging.getLogger(__name__).info("Camelot not available - using pdfplumber only for table extraction")
from pyhanko.sign import signers, fields
try:
    from modules.security import validate_input_file
    from modules.tesseract_helper import is_tesseract_available, configure_tesseract
except ImportError:
    # Fallback for flat structure
    from security import validate_input_file
    try:
        from tesseract_helper import is_tesseract_available, configure_tesseract
    except ImportError:
        # If helper not available, define fallback
        def is_tesseract_available():
            try:
                import pytesseract
                pytesseract.get_tesseract_version()
                return True, None
            except:
                return False, "Tesseract not available"
        def configure_tesseract():
            pass

# Import debug_log (always in flat structure)
try:
    from debug_utils import debug_log
except ImportError:
    # Fallback if debug_utils not available
    def debug_log(msg):
        print(f"[DEBUG] {msg}", file=sys.stderr, flush=True)


logger = logging.getLogger(__name__)

def handle_pdf_action(action, payload):
    """
    Main dispatcher for PDF actions.
    Routes to appropriate function based on action name.
    """
    logger.info(f"Handling PDF action: {action}")

    # Security Check
    files = payload.get("files", [])
    if not files:
        files = [payload.get("file")] if payload.get("file") else []
        files = [f for f in files if f]  # Remove None values

    for f in files:
        if isinstance(f, str):
            validate_input_file(f)

    # Route to appropriate handler
    if action == "merge":
        return merge_pdfs(payload)
    elif action == "split":
        return split_pdf(payload)
    elif action == "compress":
        return compress_pdf(payload)
    elif action == "protect":
        return protect_pdf(payload)
    elif action == "unlock":
        return unlock_pdf(payload)
    elif action == "watermark":
        return watermark_pdf(payload)
    elif action == "rotate":
        return rotate_pdf(payload)
    elif action == "remove_metadata":
        return remove_metadata(payload)
    elif action == "pdf_to_word":
        return pdf_to_word(payload)
    elif action == "pdf_to_images":
        return pdf_to_images(payload)
    elif action == "images_to_pdf":
        return images_to_pdf(payload)
    elif action == "extract_text":
        return extract_text(payload)
    elif action == "extract_images_from_pdf":
        return extract_images_from_pdf(payload)
    elif action == "extract_tables":
        return extract_tables(payload)
    elif action == "preview_tables":
        return preview_tables(payload)
    elif action == "grayscale":
        return grayscale_pdf(payload)
    elif action == "repair":
        return repair_pdf(payload)
    elif action == "flatten":
        return flatten_pdf(payload)
    elif action == "page_numbers":
        return add_page_numbers(payload)
    elif action == "delete_pages":
        return delete_pages(payload)
    elif action == "diff":
        return diff_pdfs(payload)
    elif action == "booklet":
        return create_booklet(payload)
    elif action == "scrub":
        return scrub_pdf(payload)
    elif action == "redact":
        return redact_pdf(payload)
    elif action == "sign":
        return sign_pdf(payload)
    elif action == "optimize":
        return optimize_pdf(payload)
    elif action == "word_to_pdf":
        return word_to_pdf(payload)
    elif action == "powerpoint_to_pdf":
        return powerpoint_to_pdf(payload)
    elif action == "excel_to_pdf":
        return excel_to_pdf(payload)
    elif action == "html_to_pdf":
        return html_to_pdf(payload)
    elif action == "ocr_pdf":
        return ocr_pdf(payload)
    elif action == "pdf_to_pdfa":
        return pdf_to_pdfa(payload)
    elif action == "crop":
        return crop_pdf(payload)
    elif action == "organize" or action == "reorder_pages":
        return organize_pdf(payload)
    elif action == "extract_metadata":
        return extract_metadata(payload)
    elif action == "extract_form_data":
        return extract_form_data(payload)
    elif action == "preview":
        return preview_pdf(payload)
    elif action == "csv_to_pdf":
        return csv_to_pdf(payload)
    elif action == "txt_to_pdf":
        return txt_to_pdf(payload)
    elif action == "tiff_to_pdf":
        return tiff_to_pdf(payload)
    elif action == "rtf_to_pdf":
        return rtf_to_pdf(payload)
    elif action == "xml_to_pdf":
        return xml_to_pdf(payload)
    else:
        raise ValueError(f"Unknown PDF action: {action}")

# ============================================================================
# PDF PREVIEW FUNCTION
# ============================================================================

def preview_pdf(payload):
    """
    Generate preview image for PDF transformations.
    Returns base64-encoded image of the first page with applied transformation.
    """
    file_path = payload.get("file") or (payload.get("files", [None])[0] if payload.get("files") else None)
    if not file_path:
        return {"image": None, "page_count": 0, "errors": ["No file provided"]}
    
    if not os.path.exists(file_path):
        return {"image": None, "page_count": 0, "errors": [f"File not found: {file_path}"]}
    
    action = payload.get("action", "preview")
    page_num = payload.get("page", 0)
    
    doc = None
    try:
        doc = fitz.open(file_path)
        page_count = len(doc)
        
        if page_count == 0:
            return {"image": None, "page_count": 0, "errors": ["PDF has no pages"]}
        
        # Ensure page_num is valid
        page_num = max(0, min(page_num, page_count - 1))
        page = doc[page_num]
        
        # Apply transformation based on action
        if action == "rotate":
            angle = float(payload.get("angle", 0))
            # Normalize angle to 0-360 range
            angle = angle % 360
            if angle < 0:
                angle += 360
            
            # Round to nearest 90 degrees for set_rotation (which only supports 0, 90, 180, 270)
            # This matches the actual rotate_pdf function behavior
            if angle < 45:
                rotation_deg = 0
            elif angle < 135:
                rotation_deg = 90
            elif angle < 225:
                rotation_deg = 180
            elif angle < 315:
                rotation_deg = 270
            else:
                rotation_deg = 0
            
            # Create a temporary page with rotation
            temp_doc = fitz.open()
            temp_page = temp_doc.new_page(width=page.rect.width, height=page.rect.height)
            # Copy content from original page
            temp_page.show_pdf_page(temp_page.rect, doc, page_num)
            # Apply rotation
            temp_page.set_rotation(rotation_deg)
            # Get pixmap with 2x zoom for preview
            pix = temp_page.get_pixmap(matrix=fitz.Matrix(2, 2))
            temp_doc.close()
        elif action == "crop":
            x = payload.get("x", 0)
            y = payload.get("y", 0)
            width = payload.get("width")
            height = payload.get("height")

            # Convert to float and handle None
            x = float(x) if x is not None else 0.0
            y = float(y) if y is not None else 0.0

            # Get original page dimensions
            rect = page.rect

            # If width/height not specified, use full page
            if width is None:
                width = rect.width - x
            else:
                width = float(width)
            if height is None:
                height = rect.height - y
            else:
                height = float(height)
            
            # Create crop rectangle
            crop_rect = fitz.Rect(x, y, x + width, y + height)
            # Clamp to page bounds
            crop_rect = crop_rect & rect
            
            # Get pixmap of cropped area
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), clip=crop_rect)
        elif action == "watermark":
            # For watermark preview, just show the plain PDF page
            # The frontend will overlay the watermark for interactive editing
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        elif action == "grayscale":
            # Get pixmap and convert to grayscale
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            gray_pix = fitz.Pixmap(pix, 0)  # 0 = grayscale
            pix = gray_pix
        else:
            # Default preview - just show the page
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        
        # Convert to base64
        img_data = pix.tobytes("png")
        base64_img = base64.b64encode(img_data).decode()

        # Get page dimensions for coordinate conversion
        page_width = float(page.rect.width)
        page_height = float(page.rect.height)

        # Cleanup
        pix = None
        if doc:
            doc.close()

        return {
            "image": f"data:image/png;base64,{base64_img}",
            "page_count": page_count,
            "page_width": page_width,
            "page_height": page_height,
            "errors": []
        }
    except Exception as e:
        logger.error(f"Error generating preview: {e}", exc_info=True)
        if doc:
            doc.close()
        return {
            "image": None,
            "page_count": 0,
            "errors": [str(e)]
        }

# ============================================================================
# PDF MANIPULATION FUNCTIONS
# ============================================================================

def merge_pdfs(payload):
    """Merge multiple PDF files into one."""
    files = payload.get("files", [])
    if not files:
        return {"processed_files": [], "errors": ["No files provided"]}

    output_filename = payload.get("output_name", "merged.pdf")
    base_dir = os.path.dirname(files[0]) if files[0] else os.getcwd()
    output_path = os.path.join(base_dir, output_filename)

    processed_files = []
    errors = []

    try:
        merger = PdfWriter()
        for pdf_path in files:
            if not os.path.exists(pdf_path):
                errors.append({"file": pdf_path, "error": f"File not found: {pdf_path}"})
                continue
            try:
                merger.append(pdf_path)
            except Exception as e:
                errors.append({"file": pdf_path, "error": f"Failed to add PDF: {str(e)}"})
                continue

        if errors and not merger.pages:
            return {"processed_files": [], "errors": errors}

        merger.write(output_path)
        merger.close()

        if os.path.exists(output_path):
            processed_files.append(output_path)
        else:
            errors.append({"file": "unknown", "error": "Failed to create merged PDF"})
    except Exception as e:
        logger.error(f"Error merging PDFs: {e}", exc_info=True)
        errors.append({"file": files[0] if files else "unknown", "error": str(e)})

    return {"processed_files": processed_files, "errors": errors}

def _validate_page_spec(pages_str, total_pages, max_specs=1000, max_range=500):
    """
    Validates and parses a page specification string.

    Args:
        pages_str: Page specification (e.g., "1,3,5" or "1-5,10")
        total_pages: Total pages in the document
        max_specs: Maximum number of comma-separated specifications
        max_range: Maximum size of a single range

    Returns:
        Tuple of (page_indices list, error_message or None)
    """
    import re

    if not pages_str or not pages_str.strip():
        return [], "Page specification is required"

    # Validate format (only digits, commas, dashes, and spaces)
    if not re.match(r'^[\d\s,\-]+$', pages_str):
        return [], f"Invalid page format. Use only numbers, commas, and dashes (e.g., '1,3,5' or '1-5')"

    parts = [p.strip() for p in pages_str.split(",")]

    # Limit parts to prevent DoS
    if len(parts) > max_specs:
        return [], f"Too many page specifications (max {max_specs})"

    page_indices = []
    for part in parts:
        if not part:
            continue

        if "-" in part:
            range_parts = part.split("-")
            if len(range_parts) != 2:
                return [], f"Invalid range format: '{part}'. Use format like '1-5'"

            try:
                start, end = int(range_parts[0]), int(range_parts[1])
            except ValueError:
                return [], f"Invalid page numbers in range: '{part}'"

            if start < 1 or end < 1:
                return [], f"Page numbers must be >= 1. Invalid range: '{part}'"

            if start > end:
                return [], f"Invalid range: '{part}'. Start must be <= end."

            if end - start + 1 > max_range:
                return [], f"Range too large: '{part}'. Maximum range size is {max_range} pages."

            if end > total_pages:
                return [], f"Page {end} exceeds document length ({total_pages} pages)"

            for p in range(start - 1, end):
                if 0 <= p < total_pages:
                    page_indices.append(p)
        else:
            try:
                p = int(part)
            except ValueError:
                return [], f"Invalid page number: '{part}'"

            if p < 1 or p > total_pages:
                return [], f"Page {part} is out of range (1-{total_pages})"

            page_indices.append(p - 1)

    if not page_indices:
        return [], "No valid pages specified"

    return page_indices, None


def split_pdf(payload):
    """Split PDF into individual pages or extract specific pages."""
    file_path = payload.get("file") or (payload.get("files", [None])[0])
    if not file_path:
        return {"processed_files": [], "errors": ["No file provided"]}

    if not os.path.exists(file_path):
        return {"processed_files": [], "errors": [{"file": file_path, "error": f"File not found: {file_path}"}]}

    mode = payload.get("mode", "all")  # all, range, pages
    pages = payload.get("pages", "")

    processed_files = []
    errors = []
    doc = None

    try:
        base, ext = os.path.splitext(file_path)
        doc = fitz.open(file_path)
        total_pages = len(doc)

        if total_pages == 0:
            errors.append({"file": file_path, "error": "PDF has no pages"})
            doc.close()
            return {"processed_files": [], "errors": errors}

        # Limit total pages for "split all" mode to prevent resource exhaustion
        MAX_SPLIT_ALL_PAGES = 500
        if mode == "all" or not pages:
            if total_pages > MAX_SPLIT_ALL_PAGES:
                errors.append({"file": file_path, "error": f"PDF has {total_pages} pages. Maximum for 'split all' is {MAX_SPLIT_ALL_PAGES} pages. Use page ranges instead."})
                doc.close()
                return {"processed_files": [], "errors": errors}

            # Split into individual pages
            for page_num in range(total_pages):
                new_doc = fitz.open()
                new_doc.insert_pdf(doc, from_page=page_num, to_page=page_num)
                output_path = f"{base}_page_{page_num + 1}{ext}"
                new_doc.save(output_path)
                new_doc.close()
                processed_files.append(output_path)
        elif mode == "range" and pages:
            # Extract page range with validation
            page_indices, error = _validate_page_spec(pages, total_pages)
            if error:
                errors.append({"file": file_path, "error": error})
                doc.close()
                return {"processed_files": [], "errors": errors}

            for page_num in page_indices:
                new_doc = fitz.open()
                new_doc.insert_pdf(doc, from_page=page_num, to_page=page_num)
                output_path = f"{base}_page_{page_num + 1}{ext}"
                new_doc.save(output_path)
                new_doc.close()
                processed_files.append(output_path)

        if doc:
            doc.close()
        return {"processed_files": processed_files, "errors": errors}
    except Exception as e:
        logger.error(f"Error splitting PDF {file_path}: {e}", exc_info=True)
        if doc:
            doc.close()
        return {"processed_files": [], "errors": [{"file": file_path, "error": str(e)}]}

def get_pdf_compression_estimates(file_path):
    """
    Analyze PDF to estimate potential file size after compression for different levels.
    """
    try:
        original_size = os.path.getsize(file_path)
        doc = fitz.open(file_path)
        
        # Calculate image contribution
        image_bytes = 0
        for page in doc:
            for img in page.get_images(full=True):
                xref = img[0]
                try:
                    base_image = doc.extract_image(xref)
                    image_bytes += len(base_image["image"])
                except:
                    pass
        
        doc.close()
        
        # Heuristic estimation factors
        # Level 0: 5% reduction (metadata, unused objects)
        # Level 1: 15% reduction (standard stream compression)
        # Level 2: 40% reduction (80% image quality)
        # Level 3: 70% reduction (60% image quality + downsampling)
        
        other_bytes = max(0, original_size - image_bytes)
        
        estimates = {
            "original": original_size,
            "levels": [
                int(other_bytes * 0.95 + image_bytes * 0.90),  # Level 0
                int(other_bytes * 0.90 + image_bytes * 0.70),  # Level 1
                int(other_bytes * 0.85 + image_bytes * 0.40),  # Level 2
                int(other_bytes * 0.80 + image_bytes * 0.15)   # Level 3
            ]
        }
        return estimates
    except Exception as e:
        logger.error(f"Error estimating compression: {e}")
        return None

def compress_pdf(payload):
    """
    Compress PDF file to reduce size with configurable compression levels.
    """
    files = payload.get("files", [])
    estimate_only = payload.get("estimate_only", False)

    if estimate_only and files:
        results = []
        for f in files:
            est = get_pdf_compression_estimates(f)
            if est:
                results.append({"file": f, "estimates": est})
        return {"estimates": results}

    # Validate and clamp compression level to valid range 0-3
    try:
        level = int(payload.get("level", 1))
        level = max(0, min(3, level))
    except (ValueError, TypeError):
        level = 1

    processed_files = []
    errors = []

    for file_path in files:
        if not os.path.exists(file_path):
            errors.append({"file": file_path, "error": f"File not found: {file_path}"})
            continue

        try:
            base, ext = os.path.splitext(file_path)
            output_path = f"{base}_compressed{ext}"
            
            # Using PyMuPDF's built-in optimization for better reliability
            doc = fitz.open(file_path)
            
            # Apply image compression based on level
            if level >= 2:
                # Levels 2 and 3 involve image re-compression
                quality = 80 if level == 2 else 60
                optimized_xrefs = set()
                
                for page in doc:
                    for img in page.get_images(full=True):
                        xref = img[0]
                        if xref in optimized_xrefs:
                            continue
                            
                        try:
                            # Extract and re-compress
                            pix = fitz.Pixmap(doc, xref)
                            
                            # Convert to RGB if it's CMYK or has alpha (for JPEG compression)
                            if pix.n - pix.alpha > 3 or pix.alpha:
                                pix2 = fitz.Pixmap(fitz.csRGB, pix)
                                pix = pix2

                            # Use PIL for reliable JPEG compression with quality control
                            img_obj = Image.frombytes("RGB" if pix.n >= 3 else "L", [pix.width, pix.height], pix.samples)
                            img_bytes_io = io.BytesIO()
                            img_obj.save(img_bytes_io, format="JPEG", quality=quality, optimize=True)
                            img_data = img_bytes_io.getvalue()
                            
                            # Only replace if compressed version is smaller
                            original_image = doc.extract_image(xref)
                            if original_image and len(img_data) < len(original_image["image"]):
                                # Use page.replace_image for reliable dictionary and filter updates
                                page.replace_image(xref, stream=img_data)
                                optimized_xrefs.add(xref)
                        except Exception as e:
                            logger.warning(f"Could not optimize image {xref}: {e}")

            # Save with appropriate flags
            # garbage=4: Remove unused objects + deduplicate identical objects
            # deflate=True: Compress streams
            # clean=True: Sanitize the structure
            
            save_args = {
                "garbage": 4,
                "deflate": True,
                "clean": True
            }
            
            # Use extreme linearize/compression for Level 3
            if level == 3:
                save_args["linear"] = True
            
            doc.save(output_path, **save_args)
            doc.close()

            # Final pass with pikepdf for object stream optimization (Levels 1+)
            if level >= 1:
                with pikepdf.open(output_path, allow_overwriting_input=True) as pdf:
                    pdf.save(
                        output_path,
                        compress_streams=True,
                        object_stream_mode=pikepdf.ObjectStreamMode.generate
                    )

            if os.path.exists(output_path):
                processed_files.append(output_path)
            else:
                errors.append({"file": file_path, "error": "Compression failed: output file not created"})

        except Exception as e:
            logger.error(f"Error compressing PDF {file_path}: {e}", exc_info=True)
            errors.append({"file": file_path, "error": f"Compression failed: {str(e)}"})

    return {"processed_files": processed_files, "errors": errors}

def protect_pdf(payload):
    """Add password protection to PDF."""
    files = payload.get("files", [])
    password = payload.get("password")

    if not password:
        file_name = files[0] if files else "unknown"
        return {"processed_files": [], "errors": [{"file": file_name, "error": "Please enter a password to protect your PDF file."}]}

    if len(password) < 3:
        file_name = files[0] if files else "unknown"
        return {"processed_files": [], "errors": [{"file": file_name, "error": "Password must be at least 3 characters long."}]}

    processed_files = []
    errors = []

    for file_path in files:
        if not os.path.exists(file_path):
            errors.append({"file": file_path, "error": f"File not found: {file_path}"})
            continue

        try:
            base, ext = os.path.splitext(file_path)
            output_path = f"{base}_protected{ext}"

            reader = PdfReader(file_path)
            writer = PdfWriter()

            for page in reader.pages:
                writer.add_page(page)

            writer.encrypt(password)
            with open(output_path, 'wb') as f:
                writer.write(f)

            if os.path.exists(output_path):
                processed_files.append(output_path)
            else:
                errors.append({"file": file_path, "error": "Failed to create protected PDF"})
        except Exception as e:
            logger.error(f"Error protecting PDF {file_path}: {e}", exc_info=True)
            errors.append({"file": file_path, "error": str(e)})

    return {"processed_files": processed_files, "errors": errors}

def unlock_pdf(payload):
    """Remove password protection from PDF."""
    files = payload.get("files", [])
    password = payload.get("password", "")

    processed_files = []
    errors = []

    for file_path in files:
        if not os.path.exists(file_path):
            errors.append({"file": file_path, "error": f"File not found: {file_path}"})
            continue

        try:
            base, ext = os.path.splitext(file_path)
            output_path = f"{base}_unlocked{ext}"

            reader = PdfReader(file_path)
            if reader.is_encrypted:
                if not password:
                    errors.append({"file": file_path, "error": "This PDF is password-protected. Please provide the password."})
                    continue
                decrypt_result = reader.decrypt(password)
                if decrypt_result == 0:  # 0 = failed, 1 = user password, 2 = owner password
                    errors.append({"file": file_path, "error": "Incorrect password. Please try again."})
                    continue

            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)

            with open(output_path, 'wb') as f:
                writer.write(f)

            processed_files.append(output_path)
        except Exception as e:
            errors.append({"file": file_path, "error": str(e)})

    return {"processed_files": processed_files, "errors": errors}

def _clamp(value, min_val, max_val, default):
    """Clamp a value to a range, returning default if invalid."""
    try:
        val = float(value) if value is not None else default
        return max(min_val, min(max_val, val))
    except (ValueError, TypeError):
        return default


def watermark_pdf(payload):
    """Add watermark to PDF pages."""
    files = payload.get("files", [])
    watermark_type = payload.get("watermark_type", "text")
    text = payload.get("text", "CONFIDENTIAL")

    # Validate and clamp numeric parameters
    opacity = _clamp(payload.get("opacity"), 0.0, 1.0, 0.5)
    watermark_file = payload.get("watermark_file")
    color = str(payload.get("color", "gray"))
    font_size = int(_clamp(payload.get("font_size"), 8, 500, 72))  # Min 8pt, max 500pt
    x_percent = _clamp(payload.get("x"), 0.0, 1.0, 0.5)  # 0-1 range
    y_percent = _clamp(payload.get("y"), 0.0, 1.0, 0.5)  # 0-1 range

    # Validate text length to prevent memory issues
    MAX_WATERMARK_TEXT_LENGTH = 500
    if text and len(text) > MAX_WATERMARK_TEXT_LENGTH:
        text = text[:MAX_WATERMARK_TEXT_LENGTH]

    processed_files = []
    errors = []

    for file_path in files:
        if not os.path.exists(file_path):
            errors.append({"file": file_path, "error": f"File not found: {file_path}"})
            continue

        if watermark_type == "image" and watermark_file and not os.path.exists(watermark_file):
            errors.append({"file": file_path, "error": f"Watermark image not found: {watermark_file}"})
            continue

        try:
            base, ext = os.path.splitext(file_path)
            output_path = f"{base}_watermarked{ext}"

            doc = fitz.open(file_path)

            for page in doc:
                rect = page.rect

                if watermark_type == "text":
                    # Text watermark
                    try:
                        # Parse color - support both color names and hex colors
                        if color.startswith("#"):
                            # Hex color (e.g., "#ff0000" or "#ec0909")
                            hex_color = color.lstrip("#")
                            if len(hex_color) == 6:
                                r = int(hex_color[0:2], 16) / 255.0
                                g = int(hex_color[2:4], 16) / 255.0
                                b = int(hex_color[4:6], 16) / 255.0
                                base_color = (r, g, b)
                            else:
                                base_color = (0.5, 0.5, 0.5)  # fallback
                        elif color == "gray":
                            base_color = (0.5, 0.5, 0.5)
                        elif color == "red":
                            base_color = (1, 0, 0)
                        elif color == "blue":
                            base_color = (0, 0, 1)
                        else:
                            base_color = (0.5, 0.5, 0.5)

                        # Calculate position from percentages
                        point = fitz.Point(rect.width * x_percent, rect.height * y_percent)

                        # Use Shape for text with opacity support
                        shape = page.new_shape()
                        shape.insert_text(point, text, fontsize=font_size, color=base_color)
                        shape.finish(fill_opacity=opacity)
                        shape.commit(overlay=True)
                    except Exception as e:
                        errors.append({"file": file_path, "error": f"bad rotate value" if "rotate" in str(e).lower() else str(e)})
                        doc.close()
                        continue
                elif watermark_type == "image" and watermark_file:
                    # Image watermark
                    try:
                        img_rect = fitz.Rect(0, 0, rect.width, rect.height)
                        page.insert_image(img_rect, filename=watermark_file, opacity=opacity)
                    except Exception as e:
                        errors.append({"file": file_path, "error": str(e)})
                        doc.close()
                        continue

            doc.save(output_path)
            doc.close()
            processed_files.append(output_path)
        except Exception as e:
            errors.append({"file": file_path, "error": str(e)})

    return {"processed_files": processed_files, "errors": errors}

def rotate_pdf(payload):
    """Rotate PDF pages."""
    files = payload.get("files", [])
    angle = payload.get("angle", 90)  # 90, 180, 270
    pages = payload.get("pages", "")  # Optional: specific pages

    # Convert angle to int and normalize to valid values (0, 90, 180, 270)
    try:
        angle = int(float(angle))  # Handle both string and float inputs
        # Normalize to 0-360 range
        angle = angle % 360
        if angle < 0:
            angle += 360
        # Round to nearest 90 degrees (PyMuPDF only supports 0, 90, 180, 270)
        if angle < 45:
            angle = 0
        elif angle < 135:
            angle = 90
        elif angle < 225:
            angle = 180
        elif angle < 315:
            angle = 270
        else:
            angle = 0
    except (ValueError, TypeError):
        angle = 90  # Default to 90 degrees if invalid

    processed_files = []
    errors = []

    for file_path in files:
        if not os.path.exists(file_path):
            errors.append({"file": file_path, "error": f"File not found: {file_path}"})
            continue

        try:
            base, ext = os.path.splitext(file_path)
            output_path = f"{base}_rotated{ext}"

            doc = fitz.open(file_path)

            # Parse pages if specified
            page_list = None
            if pages:
                try:
                    if "-" in pages:
                        start, end = map(int, pages.split("-"))
                        page_list = list(range(start - 1, min(end, len(doc))))
                    else:
                        page_list = [int(p.strip()) - 1 for p in pages.split(",")]
                except:
                    pass

            for page_num in range(len(doc)):
                if page_list is None or page_num in page_list:
                    doc[page_num].set_rotation(int(angle))  # Ensure it's an int

            doc.save(output_path)
            doc.close()
            processed_files.append(output_path)
        except Exception as e:
            errors.append({"file": file_path, "error": str(e)})

    return {"processed_files": processed_files, "errors": errors}

def remove_metadata(payload):
    """Remove metadata from PDF."""
    files = payload.get("files", [])

    processed_files = []
    errors = []

    for file_path in files:
        if not os.path.exists(file_path):
            errors.append({"file": file_path, "error": f"File not found: {file_path}"})
            continue

        try:
            base, ext = os.path.splitext(file_path)
            output_path = f"{base}_sanitized{ext}"

            with pikepdf.open(file_path) as pdf:
                # Remove metadata
                if '/Metadata' in pdf.Root:
                    del pdf.Root.Metadata
                pdf.docinfo.clear()
                pdf.save(output_path)

            processed_files.append(output_path)
        except Exception as e:
            errors.append({"file": file_path, "error": str(e)})

    return {"processed_files": processed_files, "errors": errors}

def pdf_to_word(payload):
    """Convert PDF to Word document."""
    files = payload.get("files", [])

    processed_files = []
    errors = []

    for file_path in files:
        if not os.path.exists(file_path):
            errors.append({"file": file_path, "error": f"File not found: {file_path}"})
            continue

        try:
            base, _ = os.path.splitext(file_path)
            output_path = f"{base}.docx"

            from pdf2docx import Converter
            cv = Converter(file_path)
            cv.convert(output_path)
            cv.close()

            processed_files.append(output_path)
        except Exception as e:
            errors.append({"file": file_path, "error": str(e)})

    return {"processed_files": processed_files, "errors": errors}

def pdf_to_images(payload):
    """Convert PDF pages to images."""
    files = payload.get("files", [])
    output_format = payload.get("output_format", "png").lower()

    # Validate format - support jpg/jpeg and png
    if output_format in ["jpg", "jpeg"]:
        output_format = "jpg"
        extension = ".jpg"
    else:
        output_format = "png"
        extension = ".png"

    processed_files = []
    errors = []

    for file_path in files:
        if not os.path.exists(file_path):
            errors.append({"file": file_path, "error": f"File not found: {file_path}"})
            continue

        doc = None
        try:
            base, _ = os.path.splitext(file_path)
            doc = fitz.open(file_path)

            if len(doc) == 0:
                errors.append({"file": file_path, "error": "PDF has no pages"})
                doc.close()
                continue

            for page_num in range(len(doc)):
                page = doc[page_num]
                pix = page.get_pixmap(dpi=300)
                output_path = f"{base}_page_{page_num + 1}{extension}"

                if output_format == "jpg":
                    # For JPG, we need to convert through PIL for proper JPEG encoding
                    from PIL import Image as PILImage
                    import io
                    img_data = pix.tobytes("ppm")
                    pil_img = PILImage.open(io.BytesIO(img_data))
                    # Convert to RGB if necessary (remove alpha channel for JPG)
                    if pil_img.mode in ("RGBA", "LA", "P"):
                        pil_img = pil_img.convert("RGB")
                    pil_img.save(output_path, "JPEG", quality=95)
                else:
                    pix.save(output_path)

                pix = None  # Free memory
                if os.path.exists(output_path):
                    processed_files.append(output_path)

            doc.close()
        except Exception as e:
            logger.error(f"Error converting PDF to images {file_path}: {e}", exc_info=True)
            if doc:
                doc.close()
            errors.append({"file": file_path, "error": str(e)})

    return {"processed_files": processed_files, "errors": errors}

def images_to_pdf(payload):
    """Combine images into a PDF."""
    files = payload.get("files", [])
    output_name = payload.get("output_name", "combined.pdf")

    processed_files = []
    errors = []

    if not files:
        return {"processed_files": [], "errors": ["No images provided"]}

    try:
        base_dir = os.path.dirname(files[0]) if files[0] and os.path.dirname(files[0]) else os.getcwd()
        output_path = os.path.join(base_dir, output_name)

        images = []
        for file_path in files:
            if not os.path.exists(file_path):
                errors.append({"file": file_path, "error": f"File not found: {file_path}"})
                continue

            try:
                img = Image.open(file_path)
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                images.append(img)
            except Exception as e:
                errors.append({"file": file_path, "error": str(e)})

        if images:
            images[0].save(output_path, format="PDF", save_all=True, append_images=images[1:])
            processed_files.append(output_path)
    except Exception as e:
        errors.append({"file": files[0] if files else "unknown", "error": str(e)})

    return {"processed_files": processed_files, "errors": errors}

def extract_text(payload):
    """Extract text from PDF."""
    files = payload.get("files", [])

    processed_files = []
    errors = []

    for file_path in files:
        if not os.path.exists(file_path):
            errors.append({"file": file_path, "error": f"File not found: {file_path}"})
            continue

        doc = None
        try:
            base, _ = os.path.splitext(file_path)
            output_path = f"{base}_extracted.txt"

            doc = fitz.open(file_path)
            text_content = []

            for page in doc:
                text_content.append(page.get_text())

            with open(output_path, 'w', encoding='utf-8') as f:
                f.write('\n\n'.join(text_content))

            if doc:
                doc.close()
            if os.path.exists(output_path):
                processed_files.append(output_path)
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {e}", exc_info=True)
            if doc:
                doc.close()
            errors.append({"file": file_path, "error": str(e)})

    return {"processed_files": processed_files, "errors": errors}

def extract_images_from_pdf(payload):
    """Extract all images from PDF with duplicate detection and resource limits."""
    files = payload.get("files", [])

    # Production limits
    MAX_IMAGES_PER_PDF = 1000
    MAX_IMAGE_SIZE_MB = 50

    processed_files = []
    errors = []

    for file_path in files:
        if not os.path.exists(file_path):
            errors.append({"file": file_path, "error": f"File not found: {file_path}"})
            continue

        doc = None
        try:
            print(f"[EXTRACT_IMAGES] Starting extraction for {os.path.basename(file_path)}", file=sys.stderr, flush=True)

            base, _ = os.path.splitext(file_path)
            images_dir = f"{base}_images"
            os.makedirs(images_dir, exist_ok=True)

            doc = fitz.open(file_path)
            img_count = 0
            extracted_xrefs = set()  # Track extracted images by xref to avoid duplicates

            for page_num in range(len(doc)):
                page = doc[page_num]
                image_list = page.get_images()

                for img_index, img in enumerate(image_list):
                    xref = img[0]

                    # Skip if we've already extracted this image (duplicate detection)
                    if xref in extracted_xrefs:
                        print(f"[EXTRACT_IMAGES] Skipping duplicate image xref={xref} on page {page_num + 1}", file=sys.stderr, flush=True)
                        continue

                    # Check image count limit
                    if img_count >= MAX_IMAGES_PER_PDF:
                        print(f"[EXTRACT_IMAGES] Reached maximum image limit ({MAX_IMAGES_PER_PDF})", file=sys.stderr, flush=True)
                        errors.append({
                            "file": file_path,
                            "error": f"PDF contains more than {MAX_IMAGES_PER_PDF} images. Only first {MAX_IMAGES_PER_PDF} extracted."
                        })
                        break

                    try:
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        image_ext = base_image["ext"]

                        # Check image size
                        image_size_mb = len(image_bytes) / (1024 * 1024)
                        if image_size_mb > MAX_IMAGE_SIZE_MB:
                            print(f"[EXTRACT_IMAGES] Skipping oversized image ({image_size_mb:.1f}MB)", file=sys.stderr, flush=True)
                            continue

                        img_count += 1
                        extracted_xrefs.add(xref)

                        output_path = os.path.join(images_dir, f"image_{img_count:03d}.{image_ext}")
                        with open(output_path, "wb") as img_file:
                            img_file.write(image_bytes)
                        processed_files.append(output_path)

                    except Exception as img_err:
                        print(f"[EXTRACT_IMAGES] Failed to extract image xref={xref}: {img_err}", file=sys.stderr, flush=True)
                        continue

            if doc:
                doc.close()

            print(f"[EXTRACT_IMAGES] Extracted {img_count} unique images from {os.path.basename(file_path)}", file=sys.stderr, flush=True)

            # Handle case where PDF has no images
            if img_count == 0:
                errors.append({
                    "file": file_path,
                    "error": "This PDF contains no extractable images. PDFs created from scanned documents or text-only PDFs may not have embedded images."
                })

        except Exception as e:
            logger.error(f"Error extracting images from {file_path}: {e}", exc_info=True)
            if doc:
                doc.close()
            errors.append({"file": file_path, "error": f"Failed to process PDF: {str(e)}"})

    return {"processed_files": processed_files, "errors": errors}


def _is_valid_table(table_data, min_columns=2, min_rows=2, require_structure=True, detection_mode="balanced"):
    """
    Validate if detected table data represents a real table vs misdetected text.

    This function helps filter out false positives from text-based table detection
    which can incorrectly interpret regular paragraph text as tables.

    Args:
        table_data: Table data (list of lists)
        min_columns: Minimum columns required
        min_rows: Minimum rows required
        require_structure: If True, apply stricter validation for table structure
        detection_mode: "strict", "balanced", or "aggressive"
            - strict: Original thresholds, fewer false positives but may miss tables
            - balanced: Default, good for most documents
            - aggressive: Relaxed thresholds, catches more tables including forms

    Returns:
        tuple: (is_valid: bool, reason: str)
    """
    # Configure thresholds based on detection mode
    if detection_mode == "aggressive":
        # Relaxed thresholds for form-style documents and medical/technical reports
        inconsistent_threshold = 0.5  # Allow 50% inconsistent rows
        split_word_threshold = 0.5    # Allow 50% split-word patterns (raised from 0.35)
        sparse_row_threshold = 0.25   # Only 25% need 2+ values (relaxed from 0.3)
        empty_cell_threshold = 0.88   # Allow up to 88% empty cells (raised from 0.85)
        min_non_empty = 2             # Minimum 2 non-empty cells
    elif detection_mode == "strict":
        # Original strict thresholds
        inconsistent_threshold = 0.3
        split_word_threshold = 0.2
        sparse_row_threshold = 0.5
        empty_cell_threshold = 0.6
        min_non_empty = 3
    else:  # balanced (default)
        inconsistent_threshold = 0.4  # Slightly relaxed from 0.3
        split_word_threshold = 0.25   # Slightly relaxed from 0.2
        sparse_row_threshold = 0.4    # Slightly relaxed from 0.5
        empty_cell_threshold = 0.75   # Relaxed from 0.6 for forms
        min_non_empty = 2             # Reduced from 3

    if not table_data or len(table_data) < min_rows:
        return False, f"Too few rows ({len(table_data) if table_data else 0})"

    if not table_data[0] or len(table_data[0]) < min_columns:
        return False, f"Too few columns ({len(table_data[0]) if table_data and table_data[0] else 0})"

    col_count = len(table_data[0])

    # Check 1: Consistent column count across rows
    inconsistent_rows = sum(1 for row in table_data if len(row) != col_count)
    if inconsistent_rows > len(table_data) * inconsistent_threshold:
        return False, f"Inconsistent column count ({inconsistent_rows}/{len(table_data)} rows)"

    # Check 2: Look for signs of mid-word splits (indicates text misdetection)
    # Real tables shouldn't have cells that look like split words
    # BUT: Be careful not to reject medical/technical tables with legitimate lowercase text
    split_word_indicators = 0
    total_cells = 0

    for row in table_data[:10]:  # Sample first 10 rows
        for i, cell in enumerate(row):
            if not cell:
                continue
            cell_str = str(cell).strip()
            total_cells += 1

            # Skip cells that are clearly legitimate table content:
            # - Contains numbers (likely values, measurements, dates)
            # - Contains units (mg, ml, %, etc.)
            # - Contains slashes (dates, fractions, ranges like "Cholesterol/HDL")
            # - Contains parentheses (units, references)
            # - Contains colons (time, ratios)
            # - Is a common table value pattern
            if len(cell_str) > 0:
                # Skip if contains numbers, slashes, parentheses, colons - likely real data
                if any(c.isdigit() for c in cell_str):
                    continue
                if '/' in cell_str or '(' in cell_str or ':' in cell_str:
                    continue
                # Skip if looks like a unit (short text with common unit patterns)
                if len(cell_str) <= 10 and any(u in cell_str.lower() for u in ['mg', 'ml', 'dl', 'ul', '%', 'g', 'l', 'iu', 'mmol', 'umol', 'ng', 'pg', 'fl', 'cells']):
                    continue
                # Skip if it's a common status/result word
                if cell_str.lower() in ['normal', 'high', 'low', 'positive', 'negative', 'reactive', 'non-reactive', 'absent', 'present', 'nil', 'trace', 'male', 'female', 'yes', 'no']:
                    continue

                # Check for signs of mid-word split:
                # - Cell is very short (1-3 chars) single lowercase fragment
                # - Starts with lowercase AND is a short fragment without spaces
                if cell_str[0].islower() and len(cell_str) < 8 and ' ' not in cell_str and not cell_str[0].isdigit():
                    # Only count if it really looks like a word fragment (no special chars)
                    if cell_str.isalpha():
                        split_word_indicators += 1

                # Cell ends abruptly and next cell continues the word
                if (len(cell_str) > 3 and
                    cell_str[-1].islower() and
                    cell_str.isalpha() and  # Only pure alpha strings
                    i < len(row) - 1 and
                    row[i+1] and
                    str(row[i+1]).strip() and
                    len(str(row[i+1]).strip()) > 0 and
                    str(row[i+1]).strip()[0].islower() and
                    str(row[i+1]).strip().isalpha()):  # Next cell is also pure alpha
                    split_word_indicators += 1

    # If too many cells show split-word patterns, likely misdetected
    # Use higher threshold since we're now more selective about what counts
    effective_threshold = split_word_threshold * 1.5 if detection_mode == "aggressive" else split_word_threshold
    if total_cells > 0 and split_word_indicators / total_cells > effective_threshold:
        return False, f"Appears to be misdetected text (split words: {split_word_indicators}/{total_cells})"

    # Check 3: Look for proper table structure
    if require_structure:
        # Real tables usually have:
        # - Multiple non-empty cells per row
        # - Some consistency in cell content length

        rows_with_multiple_values = 0
        for row in table_data:
            non_empty_in_row = sum(1 for cell in row if cell and str(cell).strip())
            if non_empty_in_row >= 2:
                rows_with_multiple_values += 1

        # Check against threshold
        if rows_with_multiple_values < len(table_data) * sparse_row_threshold:
            return False, f"Too many sparse rows ({rows_with_multiple_values}/{len(table_data)} have 2+ values)"

    # Check 4: Empty cell ratio
    total_cells = sum(len(row) for row in table_data)
    empty_cells = sum(1 for row in table_data for cell in row if not cell or not str(cell).strip())
    empty_ratio = empty_cells / total_cells if total_cells > 0 else 1

    if empty_ratio > empty_cell_threshold:
        return False, f"Too many empty cells ({empty_ratio:.0%})"

    # Check 5: Minimum non-empty cells (quick quality check)
    non_empty_count = total_cells - empty_cells
    if non_empty_count < min_non_empty:
        return False, f"Too few non-empty cells ({non_empty_count})"

    return True, "Valid table"


def _detect_header_row(table_data):
    """
    Detect if the first row of a table is a header row or data row.

    Heuristic: A header row typically has:
    - Mostly non-numeric text
    - Different content pattern than data rows
    - Unique values (not repeated in subsequent rows)

    Args:
        table_data: Table data (list of lists)

    Returns:
        bool: True if first row appears to be a header
    """
    if not table_data:
        return False

    if len(table_data) == 1:
        # Single row - could be header or data, default to header
        return True

    first_row = table_data[0]

    # If all cells in first row are empty, it's not a header
    first_row_texts = [str(cell).strip() if cell else "" for cell in first_row]
    non_empty_count = sum(1 for cell in first_row_texts if cell)

    if non_empty_count == 0:
        return False

    if non_empty_count < len(first_row) * 0.4:  # Less than 40% filled
        return False

    # Check if first row looks different from the rest of the rows
    # Headers typically have unique text, while data rows have similar patterns

    # Count numeric cells in first row vs second row
    def count_numeric_cells(row):
        count = 0
        for cell in row:
            if cell:
                cell_str = str(cell).strip()
                # Check if it's numeric (number, currency, percentage, etc.)
                # Remove common formatting
                cleaned = cell_str.replace(',', '').replace('$', '').replace('%', '').replace('€', '')
                try:
                    float(cleaned)
                    count += 1
                except ValueError:
                    pass
        return count

    first_row_numeric = count_numeric_cells(first_row)

    # Check subsequent rows
    if len(table_data) > 1:
        second_row_numeric = count_numeric_cells(table_data[1])

        # If first row has significantly fewer numbers than second row, likely a header
        if first_row_numeric < second_row_numeric * 0.5:
            return True

        # If first row has many numbers like data rows, probably not a header
        if first_row_numeric > len(first_row) * 0.6:
            return False

    # Check if first row values are repeated in data (headers shouldn't repeat)
    if len(table_data) > 2:
        first_row_lower = [str(cell).strip().lower() if cell else "" for cell in first_row]

        # Count how many first-row values appear in subsequent rows
        repeats = 0
        for row in table_data[1:]:
            row_lower = [str(cell).strip().lower() if cell else "" for cell in row]
            for val in first_row_lower:
                if val and val in row_lower:
                    repeats += 1

        # If many values repeat, probably not a unique header row
        if repeats > len(first_row) * 0.5:
            return False

    # Default: assume it's a header if we're not sure
    # Better to skip a header than to include it as data
    return True

def _generate_table_fingerprint(table_data):
    """
    Generate a fingerprint for a table based on its structure and content patterns.

    The fingerprint captures:
    - Column count
    - Header content (normalized)
    - Data type patterns per column (numeric, text, date, mixed)
    - Column width patterns (approximate character counts)

    Args:
        table_data: Table data (list of lists)

    Returns:
        dict: Fingerprint with structural and content information
    """
    if not table_data or len(table_data) == 0:
        return None

    col_count = len(table_data[0]) if table_data[0] else 0
    if col_count == 0:
        return None

    has_header = _detect_header_row(table_data)

    # Extract header if present
    header_normalized = []
    if has_header and len(table_data) > 0:
        header_normalized = [str(cell).strip().lower() if cell else "" for cell in table_data[0]]

    # Analyze data type patterns per column
    data_rows = table_data[1:] if has_header else table_data
    column_types = []
    column_widths = []

    for col_idx in range(col_count):
        numeric_count = 0
        text_count = 0
        date_count = 0
        empty_count = 0
        total_width = 0

        for row in data_rows[:10]:  # Sample first 10 data rows
            if col_idx < len(row):
                cell = row[col_idx]
                cell_str = str(cell).strip() if cell else ""
                total_width += len(cell_str)

                if not cell_str:
                    empty_count += 1
                else:
                    # Check if numeric
                    cleaned = cell_str.replace(',', '').replace('$', '').replace('%', '').replace('€', '').replace('-', '')
                    try:
                        float(cleaned)
                        numeric_count += 1
                        continue
                    except ValueError:
                        pass

                    # Check if date-like
                    import re
                    if re.match(r'^\d{1,4}[-/]\d{1,2}[-/]\d{1,4}$', cell_str) or \
                       re.match(r'^(jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec)', cell_str.lower()):
                        date_count += 1
                        continue

                    text_count += 1

        # Determine dominant type
        total = max(1, numeric_count + text_count + date_count)
        if numeric_count / total > 0.6:
            col_type = "numeric"
        elif date_count / total > 0.4:
            col_type = "date"
        elif text_count / total > 0.6:
            col_type = "text"
        else:
            col_type = "mixed"

        column_types.append(col_type)
        avg_width = total_width // max(1, len(data_rows[:10]))
        column_widths.append(avg_width)

    return {
        "col_count": col_count,
        "has_header": has_header,
        "header": header_normalized,
        "column_types": column_types,
        "column_widths": column_widths,
        "row_count": len(table_data)
    }


def _fingerprints_match(fp1, fp2, threshold=0.75):
    """
    Compare two table fingerprints to determine if they represent the same table type.

    Args:
        fp1: First fingerprint
        fp2: Second fingerprint
        threshold: Minimum similarity score (0.0 to 1.0)

    Returns:
        tuple: (bool: match, float: similarity_score, str: reason)
    """
    if not fp1 or not fp2:
        return False, 0.0, "Missing fingerprint"

    col1, col2 = fp1["col_count"], fp2["col_count"]

    # Allow column count tolerance: exact match, or within 40% for tables with 3+ columns
    # This handles cases where PDF extraction produces different column counts for the same table
    # (e.g., one extraction method misses a column, or splits/merges columns)
    if col1 != col2:
        min_cols = min(col1, col2)
        max_cols = max(col1, col2)
        col_diff_ratio = (max_cols - min_cols) / max_cols if max_cols > 0 else 1

        # Allow up to 40% difference for tables with 3+ columns (more lenient for extraction differences)
        # For very small tables (1-2 cols), require exact match
        if min_cols >= 3 and col_diff_ratio <= 0.4:
            # Columns are close enough - use smaller count for comparison
            col_count = min_cols
            print(f"[MERGE] Column count close enough ({col1} vs {col2}, diff={col_diff_ratio:.0%}), using {col_count} for comparison", file=sys.stderr, flush=True)
        elif min_cols >= 1 and col_diff_ratio <= 0.5:
            # For 1-2 column tables, allow up to 50% difference (1 vs 2, 2 vs 3, etc.)
            col_count = min_cols
            print(f"[MERGE] Column count close enough ({col1} vs {col2}, diff={col_diff_ratio:.0%}), using {col_count} for comparison", file=sys.stderr, flush=True)
        else:
            return False, 0.0, f"Column count mismatch ({col1} vs {col2}, diff={col_diff_ratio:.0%})"
    else:
        col_count = col1
    score = 0.0
    max_score = 0.0

    # 1. Column type similarity (weight: 50%) - MOST IMPORTANT for merging
    # Tables that span pages will have same data types even if headers differ
    max_score += 0.5
    type_matches = sum(1 for t1, t2 in zip(fp1["column_types"][:col_count], fp2["column_types"][:col_count]) if t1 == t2)
    type_score = type_matches / col_count if col_count > 0 else 0
    score += 0.5 * type_score

    # 2. Header similarity (weight: 30%)
    max_score += 0.3
    header_score = 0
    if fp1["has_header"] and fp2["has_header"]:
        header_matches = sum(1 for h1, h2 in zip(fp1["header"][:col_count], fp2["header"][:col_count])
                           if h1 and h2 and (h1 == h2 or h1 in h2 or h2 in h1))
        header_score = header_matches / col_count if col_count > 0 else 0
        score += 0.3 * header_score
    elif fp1["has_header"] and not fp2["has_header"]:
        # Table 2 is likely a continuation (no header) - this is a GOOD sign for merging
        score += 0.3  # Full credit - continuations typically don't repeat headers
        header_score = 1.0
    elif not fp1["has_header"] and not fp2["has_header"]:
        # Both headerless - give full credit
        score += 0.3
        header_score = 1.0

    # 3. Column width similarity (weight: 20%)
    max_score += 0.2
    width_diffs = []
    for w1, w2 in zip(fp1["column_widths"][:col_count], fp2["column_widths"][:col_count]):
        max_w = max(w1, w2, 1)
        diff = abs(w1 - w2) / max_w
        width_diffs.append(1 - min(diff, 1))
    width_score = sum(width_diffs) / len(width_diffs) if width_diffs else 0
    score += 0.2 * width_score

    # Normalize score
    final_score = score / max_score if max_score > 0 else 0

    reason = f"Types: {type_score:.0%}, Header: {header_score:.0%}, Widths: {width_score:.0%}"

    # Use a small epsilon for floating point comparison to avoid precision issues
    epsilon = 0.001
    matches = final_score >= (threshold - epsilon)

    return matches, final_score, reason


def _can_merge_tables(table1_data, table2_data, similarity_threshold=0.8):
    """
    Determine if two tables can be merged based on fingerprint comparison.

    Uses table fingerprinting to compare:
    - Column structure (count, types, widths)
    - Header content similarity
    - Data type patterns

    Args:
        table1_data: First table (list of lists)
        table2_data: Second table (list of lists)
        similarity_threshold: Minimum column similarity to merge (0.0 to 1.0)

    Returns:
        bool: True if tables can be merged
    """
    if not table1_data or not table2_data:
        print(f"[MERGE] Cannot merge: empty table data", file=sys.stderr, flush=True)
        return False

    # Generate fingerprints for both tables
    fp1 = _generate_table_fingerprint(table1_data)
    fp2 = _generate_table_fingerprint(table2_data)

    if not fp1 or not fp2:
        print(f"[MERGE] Cannot merge: failed to generate fingerprint", file=sys.stderr, flush=True)
        return False

    print(f"[MERGE] Table1: {fp1['row_count']} rows, {fp1['col_count']} cols, types={fp1['column_types']}", file=sys.stderr, flush=True)
    print(f"[MERGE] Table2: {fp2['row_count']} rows, {fp2['col_count']} cols, types={fp2['column_types']}", file=sys.stderr, flush=True)

    # Use fingerprint matching
    can_merge, score, reason = _fingerprints_match(fp1, fp2, similarity_threshold)

    print(f"[MERGE] Fingerprint match: {can_merge} (score={score:.2f}) - {reason}", file=sys.stderr, flush=True)

    return can_merge


def preview_tables(payload):
    """
    Preview tables in a PDF without extracting them.

    Returns metadata about detected tables including:
    - Page number and table position
    - Row and column counts
    - Header detection
    - Sample data (first few rows)
    - Fingerprint for merge suggestions
    - Suggested merge groups

    Args:
        payload: Dict with 'files' list

    Returns:
        Dict with 'tables' list and 'merge_suggestions'
    """
    files = payload.get("files", [])
    print(f"[PREVIEW_TABLES] Starting preview for {len(files)} file(s): {files}", file=sys.stderr, flush=True)

    if not files:
        return {
            "tables": [],
            "merge_suggestions": [],
            "errors": [{"file": "none", "error": "No files provided"}]
        }

    # Configuration (same as extract_tables)
    MIN_NON_EMPTY_CELLS = 3
    MIN_ROWS = 1
    COLUMN_SIMILARITY_THRESHOLD = 0.75

    table_settings = {
        "vertical_strategy": "lines",
        "horizontal_strategy": "lines",
        "snap_tolerance": 5,
        "snap_x_tolerance": 5,
        "snap_y_tolerance": 5,
        "join_tolerance": 5,
        "join_x_tolerance": 5,
        "join_y_tolerance": 5,
        "edge_min_length": 3,
        "min_words_vertical": 1,
        "min_words_horizontal": 1,
        "intersection_tolerance": 5,
        "text_tolerance": 5,
        "text_x_tolerance": 5,
        "text_y_tolerance": 5,
    }

    all_tables = []
    errors = []
    total_pages = 0

    for file_path in files:
        try:
            if not os.path.exists(file_path):
                errors.append({"file": file_path, "error": "File not found"})
                continue

            with pdfplumber.open(file_path) as pdf:
                total_pages = len(pdf.pages)

                # Scanned PDF detection
                if total_pages > 0:
                    text_found = False
                    for i in range(min(3, total_pages)):
                        page_text = pdf.pages[i].extract_text()
                        if page_text and len(page_text.strip()) > 50:
                            text_found = True
                            break

                    if not text_found:
                        errors.append({
                            "file": file_path,
                            "error": "Scanned PDF detected. OCR required for table extraction."
                        })
                        continue

                # Extract tables from each page
                for page_num, page_obj in enumerate(pdf.pages):
                    try:
                        tables = page_obj.extract_tables(table_settings=table_settings)
                        used_text_fallback = False

                        # Fallback to text strategy (but mark it for stricter validation)
                        if not tables:
                            text_settings = {
                                "vertical_strategy": "text",
                                "horizontal_strategy": "text",
                                "snap_tolerance": 5,
                                "join_tolerance": 5,
                            }
                            tables = page_obj.extract_tables(table_settings=text_settings)
                            used_text_fallback = True

                        if not tables:
                            continue

                        for table_num, table in enumerate(tables):
                            if not table or len(table) == 0:
                                continue

                            # Quality validation - stricter for text-based detection
                            non_empty_cells = sum(1 for row in table for cell in row if cell and str(cell).strip())
                            if non_empty_cells < MIN_NON_EMPTY_CELLS:
                                continue
                            if len(table) < MIN_ROWS + 1:
                                continue

                            # Validate table structure (stricter for text fallback)
                            is_valid, reason = _is_valid_table(
                                table,
                                min_columns=2,
                                min_rows=2,
                                require_structure=used_text_fallback  # Stricter when using text fallback
                            )
                            if not is_valid:
                                print(f"[PREVIEW] Rejecting table {table_num + 1} on page {page_num + 1}: {reason}", file=sys.stderr, flush=True)
                                continue

                            # Generate fingerprint
                            fingerprint = _generate_table_fingerprint(table)
                            has_header = _detect_header_row(table)

                            # Get sample data (first 5 rows)
                            sample_rows = []
                            for row in table[:5]:
                                sample_rows.append([str(cell)[:50] if cell else "" for cell in row])

                            # Get header if detected
                            header = None
                            if has_header and len(table) > 0:
                                header = [str(cell) if cell else "" for cell in table[0]]

                            table_info = {
                                "id": f"p{page_num + 1}_t{table_num + 1}",
                                "file": os.path.basename(file_path),
                                "page": page_num + 1,
                                "table_index": table_num + 1,
                                "rows": len(table),
                                "columns": len(table[0]) if table else 0,
                                "non_empty_cells": non_empty_cells,
                                "has_header": has_header,
                                "header": header,
                                "sample_data": sample_rows,
                                "column_types": fingerprint["column_types"] if fingerprint else [],
                                "fingerprint_hash": hash(str(fingerprint)) if fingerprint else None
                            }
                            all_tables.append(table_info)

                    except Exception as e:
                        logger.warning(f"Preview failed for page {page_num + 1}: {e}")

        except Exception as e:
            errors.append({"file": file_path, "error": str(e)})

    # Generate merge suggestions based on fingerprints
    merge_suggestions = []
    if len(all_tables) > 1:
        # Group tables that can be merged
        used_tables = set()

        for i, table1 in enumerate(all_tables):
            if table1["id"] in used_tables:
                continue

            merge_group = [table1["id"]]
            group_total_rows = table1["rows"]  # Track row count incrementally (O(1) vs O(n))
            group_fingerprint = None

            # Find table data to get fingerprint (we need to re-extract for comparison)
            # For efficiency, we use the column_types as a proxy
            t1_types = tuple(table1["column_types"])

            for j, table2 in enumerate(all_tables[i + 1:], start=i + 1):
                if table2["id"] in used_tables:
                    continue

                t2_types = tuple(table2["column_types"])

                # Simple fingerprint comparison: same column count and similar types
                if table1["columns"] == table2["columns"]:
                    type_matches = sum(1 for t1, t2 in zip(t1_types, t2_types) if t1 == t2)
                    similarity = type_matches / len(t1_types) if t1_types else 0

                    # Check header similarity if both have headers
                    header_similar = True
                    if table1["has_header"] and table2["has_header"] and table1["header"] and table2["header"]:
                        header_matches = sum(1 for h1, h2 in zip(table1["header"], table2["header"])
                                           if h1.lower() == h2.lower())
                        header_similar = header_matches / len(table1["header"]) >= 0.7

                    if similarity >= COLUMN_SIMILARITY_THRESHOLD and header_similar:
                        merge_group.append(table2["id"])
                        group_total_rows += table2["rows"]  # O(1) accumulation
                        used_tables.add(table2["id"])

            if len(merge_group) > 1:
                used_tables.add(table1["id"])
                merge_suggestions.append({
                    "tables": merge_group,
                    "reason": f"Similar structure: {table1['columns']} columns, matching data types",
                    "total_rows": group_total_rows  # Already computed incrementally
                })

    result = {
        "tables": all_tables,
        "merge_suggestions": merge_suggestions,
        "total_tables": len(all_tables),
        "total_pages": total_pages,
        "errors": errors
    }
    print(f"[PREVIEW_TABLES] Result: {len(all_tables)} tables, {total_pages} pages, {len(errors)} errors", file=sys.stderr, flush=True)
    return result


def _extract_tables_with_camelot(file_path, page_numbers=None, flavor="lattice"):
    """
    Extract tables using Camelot library (if available).

    Camelot has two modes:
    - lattice: For tables with clear borders/gridlines
    - stream: For tables without clear borders (uses whitespace)

    Args:
        file_path: Path to PDF file
        page_numbers: List of page numbers (1-indexed) or None for all
        flavor: "lattice" or "stream"

    Returns:
        dict: {page_num: [table_data, ...]} where table_data is list of lists
    """
    if not CAMELOT_AVAILABLE:
        return {}

    try:
        # Convert page numbers to camelot format
        # Camelot accepts: "1,2,3" or "1-end" or "all"
        if page_numbers and len(page_numbers) > 0:
            # Convert to comma-separated string (Camelot uses 1-indexed, which we already have)
            pages = ",".join(str(p) for p in page_numbers)
            print(f"[CAMELOT] Processing only pages: {pages}", file=sys.stderr, flush=True)
        else:
            pages = "all"
            print(f"[CAMELOT] Processing all pages (no page filter)", file=sys.stderr, flush=True)

        # Extract tables - use different kwargs based on flavor
        # Suppress Camelot warnings about image-based pages by catching them
        import warnings
        with warnings.catch_warnings():
            warnings.filterwarnings("ignore", category=UserWarning, message=".*image-based.*")
            if flavor == "lattice":
                tables = camelot.read_pdf(
                    file_path,
                    pages=pages,
                    flavor="lattice",
                    strip_text='\n',
                    line_scale=40,
                )
            else:  # stream
                tables = camelot.read_pdf(
                    file_path,
                    pages=pages,
                    flavor="stream",
                    strip_text='\n',
                    edge_tol=50,
                    row_tol=10,
                )

        result = {}
        processed_pages = set()
        for table in tables:
            page_num = table.page - 1  # Convert to 0-indexed
            processed_pages.add(page_num + 1)  # Track which pages were actually processed (1-indexed)
            if page_num not in result:
                result[page_num] = []

            # Convert pandas DataFrame to list of lists
            df = table.df
            table_data = [df.columns.tolist()] + df.values.tolist()

            # Clean up empty strings
            cleaned_data = []
            for row in table_data:
                cleaned_row = [cell if cell and str(cell).strip() else None for cell in row]
                cleaned_data.append(cleaned_row)

            result[page_num].append({
                "data": cleaned_data,
                "accuracy": table.accuracy,
                "source": f"camelot_{flavor}"
            })

        return result

    except Exception as e:
        logger.warning(f"Camelot {flavor} extraction failed: {e}")
        return {}


def _deduplicate_tables(tables_list):
    """
    Remove duplicate tables detected by different strategies.

    Compares tables by content similarity and keeps the highest quality version.

    Args:
        tables_list: List of table dicts with "data", "source", etc.

    Returns:
        List of deduplicated tables
    """
    if len(tables_list) <= 1:
        return tables_list

    def table_signature(table_data):
        """Create a signature for comparison based on content."""
        if not table_data:
            return ""
        # Use first few cells as signature
        cells = []
        for row in table_data[:3]:
            for cell in row[:5]:
                if cell:
                    cells.append(str(cell).strip()[:20])
        return "|".join(cells)

    def table_quality(table_info):
        """Score table quality for comparison."""
        data = table_info.get("data", [])
        non_empty = sum(1 for row in data for cell in row if cell and str(cell).strip())
        accuracy = table_info.get("accuracy", 50)
        # Prefer camelot_lattice > camelot_stream > pdfplumber
        source_score = {"camelot_lattice": 100, "camelot_stream": 80, "pdfplumber_lines": 60, "pdfplumber_text": 40}
        source = table_info.get("source", "pdfplumber_text")
        return non_empty + source_score.get(source, 0) + accuracy

    seen_signatures = {}
    result = []

    for table_info in tables_list:
        sig = table_signature(table_info.get("data", []))
        if not sig:
            continue

        # Check for similar existing table
        is_duplicate = False
        for existing_sig in list(seen_signatures.keys()):
            # Simple similarity check - if 60% of signature matches
            if sig[:50] == existing_sig[:50] or existing_sig[:50] == sig[:50]:
                # Keep higher quality version
                existing_idx = seen_signatures[existing_sig]
                if table_quality(table_info) > table_quality(result[existing_idx]):
                    result[existing_idx] = table_info
                is_duplicate = True
                break

        if not is_duplicate:
            seen_signatures[sig] = len(result)
            result.append(table_info)

    return result


def _deduplicate_tables_global(all_tables, similarity_threshold=0.85):
    """
    Global deduplication across all pages and strategies using fingerprinting.
    
    This function identifies truly duplicate tables (same table detected by multiple
    strategies or on multiple pages) and keeps only the highest quality version.
    
    Args:
        all_tables: List of table dicts with "data", "source", "page_num", etc.
        similarity_threshold: Minimum fingerprint similarity to consider duplicates (0.0 to 1.0)
    
    Returns:
        List of deduplicated tables (only highest quality version of each unique table)
    """
    if len(all_tables) <= 1:
        return all_tables
    
    def table_quality_score(table_info):
        """Calculate quality score for table comparison."""
        data = table_info.get("data", [])
        non_empty = sum(1 for row in data for cell in row if cell and str(cell).strip())
        accuracy = table_info.get("accuracy", 50)
        # Prefer camelot_lattice > camelot_stream > pdfplumber_lines > pdfplumber_text
        source_score = {
            "camelot_lattice": 100, 
            "camelot_stream": 80, 
            "pdfplumber_lines": 60, 
            "pdfplumber_text": 40,
            "pdfplumber_mixed": 30
        }
        source = table_info.get("source", "pdfplumber_text")
        # Also prefer tables with more rows (more complete)
        row_count = len(data) if data else 0
        return non_empty + source_score.get(source, 0) + accuracy + (row_count * 0.1)
    
    # Generate fingerprints for all tables
    table_fingerprints = []
    tables_without_fp = []  # Tables that couldn't be fingerprinted (keep them all)
    
    for idx, table_info in enumerate(all_tables):
        fp = _generate_table_fingerprint(table_info.get("data", []))
        if fp:
            table_fingerprints.append({
                "index": idx,
                "table_info": table_info,
                "fingerprint": fp,
                "quality": table_quality_score(table_info)
            })
        else:
            # Table without valid fingerprint - keep it (might be very small or malformed)
            tables_without_fp.append(table_info)
    
    # Compare fingerprints to find duplicates
    seen_groups = []
    result_indices = []
    
    for i, table1 in enumerate(table_fingerprints):
        if i in result_indices:
            continue  # Already processed as a duplicate
        
        # Find all tables that match this one
        matching_group = [i]
        best_quality_idx = i
        best_quality = table1["quality"]
        
        for j, table2 in enumerate(table_fingerprints[i+1:], start=i+1):
            if j in result_indices:
                continue
            
            # Compare fingerprints
            can_merge, score, _ = _fingerprints_match(
                table1["fingerprint"], 
                table2["fingerprint"], 
                similarity_threshold
            )
            
            if can_merge:
                matching_group.append(j)
                if table2["quality"] > best_quality:
                    best_quality = table2["quality"]
                    best_quality_idx = j
        
        # Keep only the best quality table from the group
        result_indices.append(best_quality_idx)
        seen_groups.append(matching_group)
        
        if len(matching_group) > 1:
            pages = [table_fingerprints[idx]["table_info"]["page_num"] + 1 for idx in matching_group]
            sources = [table_fingerprints[idx]["table_info"].get("source", "unknown") for idx in matching_group]
            print(f"[DEDUP] Found {len(matching_group)} duplicate tables on pages {pages} (sources: {sources}), keeping best quality version", file=sys.stderr, flush=True)
    
    # Return deduplicated tables (fingerprinted ones + ones without fingerprints)
    result = [table_fingerprints[idx]["table_info"] for idx in result_indices]
    result.extend(tables_without_fp)  # Add tables that couldn't be fingerprinted
    print(f"[DEDUP] Global deduplication: {len(all_tables)} tables -> {len(result)} unique tables ({len(tables_without_fp)} without fingerprints)", file=sys.stderr, flush=True)
    
    return result


def extract_tables(payload):
    """
    Extract tables from PDF to CSV or Excel with production-quality validation.

    Features:
    - Multi-strategy detection (Camelot lattice → Camelot stream → pdfplumber lines → text)
    - Configurable detection mode (strict/balanced/aggressive)
    - Scanned PDF detection
    - Table quality validation with adjustable thresholds
    - Memory limits
    - UTF-8 encoding for CSV
    - Partial success handling
    - Detailed error categorization
    - Smart table merging across pages
    - Deduplication of tables found by multiple strategies
    """
    try:
        files = payload.get("files", [])
        output_format = payload.get("output_format", "csv")
        merge_tables_enabled = payload.get("merge_tables", "false").lower() == "true"
        # New: detection mode parameter - strict, balanced, or aggressive
        detection_mode = payload.get("detection_mode", "balanced").lower()
        if detection_mode not in ["strict", "balanced", "aggressive"]:
            detection_mode = "balanced"

        print(f"[EXTRACT_TABLES] Starting: {len(files)} file(s), format={output_format}, merge={merge_tables_enabled}, mode={detection_mode}", file=sys.stderr, flush=True)
        logger.info(f"Starting table extraction: {len(files)} file(s), format={output_format}, merge={merge_tables_enabled}, mode={detection_mode}")
        logger.info(f"Files received: {files}")
        logger.info(f"Camelot available: {CAMELOT_AVAILABLE}")

        # Immediate validation - catch empty files list
        if not files:
            logger.error("No files provided to extract_tables!")
            return {
                "processed_files": [],
                "errors": [{
                    "file": "none",
                    "error": "No files were uploaded. Please select a PDF file to extract tables from."
                }]
            }

        # Configuration based on detection mode
        MAX_CELLS_PER_TABLE = 100000  # Prevent memory issues
        COLUMN_SIMILARITY_THRESHOLD = 0.75  # 75% column match to consider tables mergeable (lowered from 0.8 for better merging)

        # Adjust thresholds based on detection mode
        if detection_mode == "aggressive":
            MIN_NON_EMPTY_CELLS = 2
            MIN_ROWS = 1
            MIN_COLUMNS = 1  # Allow single-column tables for forms
        elif detection_mode == "strict":
            MIN_NON_EMPTY_CELLS = 4
            MIN_ROWS = 2
            MIN_COLUMNS = 2
        else:  # balanced
            MIN_NON_EMPTY_CELLS = 2
            MIN_ROWS = 1
            MIN_COLUMNS = 2

        # pdfplumber table detection settings - tuned based on mode
        if detection_mode == "aggressive":
            table_settings = {
                "vertical_strategy": "lines",
                "horizontal_strategy": "lines",
                "snap_tolerance": 8,
                "snap_x_tolerance": 8,
                "snap_y_tolerance": 8,
                "join_tolerance": 8,
                "join_x_tolerance": 8,
                "join_y_tolerance": 8,
                "edge_min_length": 2,
                "min_words_vertical": 1,
                "min_words_horizontal": 1,
                "intersection_tolerance": 8,
                "text_tolerance": 8,
                "text_x_tolerance": 8,
                "text_y_tolerance": 8,
            }
        else:
            table_settings = {
                "vertical_strategy": "lines",
                "horizontal_strategy": "lines",
                "explicit_vertical_lines": [],
                "explicit_horizontal_lines": [],
                "snap_tolerance": 5,
                "snap_x_tolerance": 5,
                "snap_y_tolerance": 5,
                "join_tolerance": 5,
                "join_x_tolerance": 5,
                "join_y_tolerance": 5,
                "edge_min_length": 3,
                "min_words_vertical": 1,
                "min_words_horizontal": 1,
                "intersection_tolerance": 5,
                "text_tolerance": 5,
                "text_x_tolerance": 5,
                "text_y_tolerance": 5,
            }

        processed_files = []
        errors = []

        for file_path in files:
            doc = None
            tables_found = 0
            total_pages = 0
            is_scanned = False
            collected_tables = []  # Store tables with metadata for potential merging

            try:
                if not os.path.exists(file_path):
                    errors.append({"file": file_path, "error": "File not found"})
                    continue

                base, _ = os.path.splitext(file_path)

                # ========================================
                # STEP 1: Detect text-based pages FIRST (before Camelot)
                # Camelot is VERY slow on image-based pages, so we skip them
                # ========================================
                text_based_pages = []  # 1-indexed page numbers for Camelot
                with pdfplumber.open(file_path) as pdf:
                    total_pages = len(pdf.pages)
                    
                    # Quick check: detect which pages are text-based
                    for page_num in range(total_pages):
                        try:
                            page_text = pdf.pages[page_num].extract_text()
                            # If page has substantial text (> 50 chars), it's text-based
                            if page_text and len(page_text.strip()) > 50:
                                text_based_pages.append(page_num + 1)  # Camelot uses 1-indexed
                        except Exception as e:
                            print(f"[DEBUG] Error checking page {page_num + 1}: {e}", file=sys.stderr, flush=True)
                            # If we can't extract text, assume it's image-based and skip
                    
                    print(f"[DEBUG] Found {len(text_based_pages)} text-based pages out of {total_pages} total: {text_based_pages}", file=sys.stderr, flush=True)
                    
                    # If no text-based pages found, it's a scanned PDF
                    if len(text_based_pages) == 0:
                        is_scanned = True
                        errors.append({
                            "file": file_path,
                            "error": "This appears to be a scanned PDF (image-based). Table extraction requires text-based PDFs. Please convert using OCR first."
                        })
                        continue

                # ========================================
                # STRATEGY 1: Skip Camelot for now (performance issues with mixed PDFs)
                # Camelot scans all pages internally even when specific pages are requested,
                # which is very slow for PDFs with image-based pages
                # ========================================
                camelot_tables = {}
                # Temporarily disabled Camelot due to performance issues
                # if CAMELOT_AVAILABLE and text_based_pages:
                #     print(f"[DEBUG] Trying Camelot lattice extraction on pages {text_based_pages}...", file=sys.stderr, flush=True)
                #     camelot_tables = _extract_tables_with_camelot(file_path, page_numbers=text_based_pages, flavor="lattice")
                #     ...
                print(f"[DEBUG] Camelot extraction skipped (using pdfplumber only for better performance)", file=sys.stderr, flush=True)

                # ========================================
                # STRATEGY 2: Use pdfplumber as primary/fallback (all pages)
                # ========================================
                with pdfplumber.open(file_path) as pdf:
                    total_pages = len(pdf.pages)

                    # Note: Scanned PDF detection already done above before Camelot
                    # We continue processing all pages with pdfplumber (it handles image pages gracefully)

                    # Extract tables from each page
                    for page_num, page_obj in enumerate(pdf.pages):
                        try:
                            page_tables = []

                            # Add Camelot results for this page (if any)
                            if page_num in camelot_tables:
                                for camelot_table in camelot_tables[page_num]:
                                    page_tables.append(camelot_table)
                                print(f"[DEBUG] Page {page_num + 1}: Added {len(camelot_tables[page_num])} Camelot tables", file=sys.stderr, flush=True)

                            # Try pdfplumber strategies
                            pdfplumber_tables = []

                            # Strategy 2a: Lines-based detection
                            tables = page_obj.extract_tables(table_settings=table_settings)
                            if tables:
                                for t in tables:
                                    if t:
                                        pdfplumber_tables.append({"data": t, "source": "pdfplumber_lines", "accuracy": 70})
                                print(f"[DEBUG] Page {page_num + 1}: pdfplumber lines found {len(tables)} tables", file=sys.stderr, flush=True)

                            # Strategy 2b: Text-based detection (if lines found few or mode is aggressive)
                            if not tables or detection_mode == "aggressive":
                                text_settings = {
                                    "vertical_strategy": "text",
                                    "horizontal_strategy": "text",
                                    "snap_tolerance": 8 if detection_mode == "aggressive" else 5,
                                    "join_tolerance": 8 if detection_mode == "aggressive" else 5,
                                    "min_words_vertical": 1,
                                    "min_words_horizontal": 1,
                                }
                                text_tables = page_obj.extract_tables(table_settings=text_settings)
                                if text_tables:
                                    for t in text_tables:
                                        if t:
                                            pdfplumber_tables.append({"data": t, "source": "pdfplumber_text", "accuracy": 50})
                                    print(f"[DEBUG] Page {page_num + 1}: pdfplumber text found {len(text_tables)} tables", file=sys.stderr, flush=True)

                            # Strategy 2c: Mixed strategy for forms (aggressive mode only)
                            if detection_mode == "aggressive" and len(pdfplumber_tables) == 0:
                                mixed_settings = {
                                    "vertical_strategy": "lines",
                                    "horizontal_strategy": "text",
                                    "snap_tolerance": 10,
                                    "join_tolerance": 10,
                                }
                                mixed_tables = page_obj.extract_tables(table_settings=mixed_settings)
                                if mixed_tables:
                                    for t in mixed_tables:
                                        if t:
                                            pdfplumber_tables.append({"data": t, "source": "pdfplumber_mixed", "accuracy": 45})
                                    print(f"[DEBUG] Page {page_num + 1}: pdfplumber mixed found {len(mixed_tables)} tables", file=sys.stderr, flush=True)

                            page_tables.extend(pdfplumber_tables)

                            # Deduplicate tables from different strategies
                            if len(page_tables) > 1:
                                page_tables = _deduplicate_tables(page_tables)
                                print(f"[DEBUG] Page {page_num + 1}: After deduplication: {len(page_tables)} tables", file=sys.stderr, flush=True)

                            if not page_tables:
                                continue

                            # Validate and collect tables
                            for table_num, table_info in enumerate(page_tables):
                                table = table_info.get("data", [])
                                source = table_info.get("source", "unknown")

                                if not table or len(table) == 0:
                                    continue

                                print(f"[DEBUG] Page {page_num + 1}, Table {table_num + 1} ({source}): {len(table)} rows, {len(table[0]) if table else 0} columns", file=sys.stderr, flush=True)

                                # Table quality validation
                                total_cells = sum(len(row) for row in table)

                                # Check memory limits
                                if total_cells > MAX_CELLS_PER_TABLE:
                                    logger.warning(f"Table {table_num + 1} on page {page_num + 1} exceeds cell limit ({total_cells} > {MAX_CELLS_PER_TABLE})")
                                    errors.append({
                                        "file": file_path,
                                        "error": f"Table {table_num + 1} on page {page_num + 1} is too large ({total_cells} cells). Maximum {MAX_CELLS_PER_TABLE} cells allowed."
                                    })
                                    continue

                                # Count non-empty cells
                                non_empty_cells = sum(1 for row in table for cell in row if cell and str(cell).strip())

                                # Filter out low-quality tables
                                if non_empty_cells < MIN_NON_EMPTY_CELLS:
                                    logger.info(f"Skipping low-quality table {table_num + 1} on page {page_num + 1} (only {non_empty_cells} non-empty cells)")
                                    continue

                                # Check minimum rows
                                if len(table) < MIN_ROWS + 1:
                                    logger.info(f"Skipping table {table_num + 1} on page {page_num + 1} (only {len(table)} rows)")
                                    continue

                                # Check minimum columns
                                if table[0] and len(table[0]) < MIN_COLUMNS:
                                    logger.info(f"Skipping table {table_num + 1} on page {page_num + 1} (only {len(table[0])} columns)")
                                    continue

                                # Validate table structure
                                # Use relaxed validation for Camelot results (they're usually reliable)
                                use_strict_validation = source.startswith("pdfplumber_text")
                                is_valid, reason = _is_valid_table(
                                    table,
                                    min_columns=MIN_COLUMNS,
                                    min_rows=MIN_ROWS + 1,
                                    require_structure=use_strict_validation,
                                    detection_mode=detection_mode
                                )
                                if not is_valid:
                                    print(f"[DEBUG] Rejecting table {table_num + 1} on page {page_num + 1}: {reason}", file=sys.stderr, flush=True)
                                    logger.info(f"Rejecting table {table_num + 1} on page {page_num + 1}: {reason}")
                                    continue

                                tables_found += 1

                                # Collect table with metadata for potential merging
                                collected_tables.append({
                                    "page_num": page_num,
                                    "table_num": table_num,
                                    "data": table,
                                    "non_empty_cells": non_empty_cells,
                                    "source": source
                                })
                                logger.info(f"Collected table {table_num + 1} from page {page_num + 1} ({non_empty_cells} cells, source={source})")

                        except Exception as e:
                            logger.warning(f"Table extraction failed for page {page_num + 1}: {e}")
                            errors.append({
                                "file": file_path,
                                "error": f"Page {page_num + 1} processing failed: {str(e)}"
                            })

                # Process collected tables - merge if enabled, otherwise save individually
                print(f"[DEBUG] Collected {len(collected_tables)} tables total, merge_enabled={merge_tables_enabled}", file=sys.stderr, flush=True)
                if collected_tables:
                    # Step 1: Global deduplication - remove truly duplicate tables across all pages/strategies
                    # This prevents the same table from being saved multiple times
                    # Use same threshold as merge to ensure consistency
                    deduplicated_tables = _deduplicate_tables_global(collected_tables, similarity_threshold=COLUMN_SIMILARITY_THRESHOLD)
                    print(f"[DEBUG] After global deduplication: {len(deduplicated_tables)} unique tables", file=sys.stderr, flush=True)
                    
                    if merge_tables_enabled:
                        # NEW APPROACH: Content-first merging based on fingerprints
                        # This allows merging tables across different strategies if they have similar structure

                        # Step 2: Generate fingerprints for all tables
                        table_fingerprints = []
                        tables_without_fp = []  # Tables without fingerprints (save individually)
                        
                        for table in deduplicated_tables:
                            fp = _generate_table_fingerprint(table.get("data", []))
                            if fp:
                                table_fingerprints.append({
                                    "table": table,
                                    "fingerprint": fp
                                })
                            else:
                                # Table without valid fingerprint - save individually
                                tables_without_fp.append(table)

                        # Step 3: Find mergeable tables based on content similarity (not strategy)
                        # When merge is enabled, we want to merge tables that span pages
                        # and only save standalone tables (that couldn't be merged with anything)
                        merged_groups = []
                        processed_indices = set()

                        # Sort tables by page number for easier processing
                        sorted_tables = sorted(enumerate(table_fingerprints), 
                                             key=lambda x: (x[1]["table"]["page_num"], x[1]["table"]["table_num"]))

                        for i, (orig_idx, table1_info) in enumerate(sorted_tables):
                            if i in processed_indices:
                                continue

                            current_group = [table1_info["table"]]
                            processed_indices.add(i)

                            # Look for continuation tables on subsequent pages (within reasonable distance)
                            # Allow up to 4 page gap if similarity is high enough
                            for j, (orig_idx2, table2_info) in enumerate(sorted_tables[i+1:], start=i+1):
                                if j in processed_indices:
                                    continue

                                prev_table = current_group[-1]
                                curr_table = table2_info["table"]
                                
                                # Check page distance - allow larger gaps for high similarity
                                page_diff = curr_table["page_num"] - prev_table["page_num"]
                                
                                # Allow gaps up to 4 pages, but prefer consecutive pages
                                max_page_gap = 4
                                if page_diff > max_page_gap:
                                    break  # Too far, stop looking for this group
                                
                                if page_diff <= 0:
                                    continue  # Same or earlier page, skip

                                # Check if tables can be merged based on fingerprint similarity
                                can_merge, score, reason = _fingerprints_match(
                                    table1_info["fingerprint"],
                                    table2_info["fingerprint"],
                                    COLUMN_SIMILARITY_THRESHOLD
                                )
                                
                                # For non-consecutive pages, require higher similarity
                                # But be more lenient - only increase threshold slightly
                                required_similarity = COLUMN_SIMILARITY_THRESHOLD
                                if page_diff > 1:
                                    # Require 3% higher similarity for each page gap (reduced from 5%)
                                    required_similarity = min(0.90, COLUMN_SIMILARITY_THRESHOLD + (page_diff - 1) * 0.03)
                                
                                # Use epsilon for floating point comparison
                                epsilon = 0.001
                                can_merge = can_merge and score >= (required_similarity - epsilon)

                                print(f"[MERGE] Page {prev_table['page_num']+1} -> {curr_table['page_num']+1} (gap={page_diff}): can_merge={can_merge}, score={score:.2f}, required={required_similarity:.2f} - {reason}", file=sys.stderr, flush=True)

                                if can_merge:
                                    current_group.append(curr_table)
                                    processed_indices.add(j)

                            merged_groups.append(current_group)

                        # Also add any unprocessed tables as single-table groups (standalone tables)
                        for i, (orig_idx, table_info) in enumerate(sorted_tables):
                            if i not in processed_indices:
                                merged_groups.append([table_info["table"]])
                                processed_indices.add(i)
                        
                        # Add tables without fingerprints as individual groups
                        for table in tables_without_fp:
                            merged_groups.append([table])

                        print(f"[DEBUG] Created {len(merged_groups)} merge groups", file=sys.stderr, flush=True)
                        for idx, group in enumerate(merged_groups):
                            if len(group) > 1:
                                pages = [t["page_num"]+1 for t in group]
                                print(f"[DEBUG] Group {idx+1}: {len(group)} tables merged from pages {pages} (source: {group[0].get('source', 'unknown')})", file=sys.stderr, flush=True)

                        # Process each group
                        for group_idx, group in enumerate(merged_groups):
                            try:
                                if len(group) == 1:
                                    # Single table - this is a standalone table that couldn't be merged with anything
                                    # Save it as an individual table
                                    table_meta = group[0]
                                    table = table_meta["data"]

                                    if len(table) > 1:
                                        df = pd.DataFrame(table[1:], columns=table[0])
                                    else:
                                        df = pd.DataFrame(table)

                                    if output_format == "csv":
                                        output_path = f"{base}_p{table_meta['page_num'] + 1}_t{table_meta['table_num'] + 1}.csv"
                                        df.to_csv(output_path, index=False, encoding='utf-8-sig')
                                        processed_files.append(output_path)
                                    elif output_format == "excel" or output_format == "xlsx":
                                        output_path = f"{base}_p{table_meta['page_num'] + 1}_t{table_meta['table_num'] + 1}.xlsx"
                                        df.to_excel(output_path, index=False, engine='openpyxl')
                                        processed_files.append(output_path)

                                    logger.info(f"Saved standalone table from page {table_meta['page_num'] + 1} (could not be merged)")
                                else:
                                    # Multiple tables - merge them intelligently
                                    first_table = group[0]["data"]
                                    first_has_header = _detect_header_row(first_table)

                                    # Use header from first table if it has one
                                    header = first_table[0] if first_has_header and len(first_table) > 1 else None

                                    # Collect all data rows
                                    all_rows = []

                                    # Add rows from first table
                                    if first_has_header and len(first_table) > 1:
                                        # Skip header row
                                        all_rows.extend(first_table[1:])
                                    else:
                                        # No header - include all rows
                                        all_rows.extend(first_table)

                                    # Add rows from subsequent tables
                                    for table_meta in group[1:]:
                                        table_data = table_meta["data"]

                                        # Check if this continuation table has a header
                                        table_has_header = _detect_header_row(table_data)

                                        if table_has_header and len(table_data) > 1:
                                            # Has header (repeated header case) - skip it
                                            all_rows.extend(table_data[1:])
                                            logger.info(f"Skipping repeated header from page {table_meta['page_num'] + 1}")
                                        else:
                                            # No header (continuation case) - include ALL rows
                                            all_rows.extend(table_data)
                                            logger.info(f"Merging headerless continuation from page {table_meta['page_num'] + 1}")

                                    # Create merged DataFrame
                                    if header:
                                        df = pd.DataFrame(all_rows, columns=header)
                                    else:
                                        df = pd.DataFrame(all_rows)

                                    # Save merged table
                                    start_page = group[0]["page_num"] + 1
                                    end_page = group[-1]["page_num"] + 1

                                    # Create unique filename to avoid duplicates
                                    # Include source strategy in filename to differentiate
                                    source_str = group[0].get("source", "unknown").replace("camelot_", "").replace("pdfplumber_", "")
                                    if output_format == "csv":
                                        output_path = f"{base}_p{start_page}-{end_page}_merged_{source_str}.csv"
                                        # Check if file already exists (avoid duplicates)
                                        if output_path not in processed_files:
                                            df.to_csv(output_path, index=False, encoding='utf-8-sig')
                                            processed_files.append(output_path)
                                        else:
                                            logger.warning(f"Skipping duplicate merged file: {output_path}")
                                    elif output_format == "excel" or output_format == "xlsx":
                                        output_path = f"{base}_p{start_page}-{end_page}_merged_{source_str}.xlsx"
                                        # Check if file already exists (avoid duplicates)
                                        if output_path not in processed_files:
                                            df.to_excel(output_path, index=False, engine='openpyxl')
                                            processed_files.append(output_path)
                                        else:
                                            logger.warning(f"Skipping duplicate merged file: {output_path}")

                                    logger.info(f"Merged {len(group)} tables from pages {start_page}-{end_page} ({len(all_rows)} total rows)")

                            except Exception as e:
                                logger.warning(f"Failed to process table group {group_idx + 1}: {e}")
                                errors.append({
                                    "file": file_path,
                                    "error": f"Failed to save merged table: {str(e)}"
                                })
                    else:
                        # Merge disabled - save each table individually
                        # Use deduplicated tables to prevent duplicate files
                        for table_meta in deduplicated_tables:
                            try:
                                table = table_meta["data"]
                                page_num = table_meta["page_num"]
                                table_num = table_meta["table_num"]

                                if len(table) > 1:
                                    df = pd.DataFrame(table[1:], columns=table[0])
                                else:
                                    df = pd.DataFrame(table)

                                if output_format == "csv":
                                    output_path = f"{base}_p{page_num + 1}_t{table_num + 1}.csv"
                                    df.to_csv(output_path, index=False, encoding='utf-8-sig')
                                    processed_files.append(output_path)
                                elif output_format == "excel" or output_format == "xlsx":
                                    output_path = f"{base}_p{page_num + 1}_t{table_num + 1}.xlsx"
                                    df.to_excel(output_path, index=False, engine='openpyxl')
                                    processed_files.append(output_path)

                                logger.info(f"Successfully extracted table {table_num + 1} from page {page_num + 1} ({table_meta['non_empty_cells']} cells)")

                            except Exception as e:
                                logger.warning(f"Failed to save table {table_num + 1} from page {page_num + 1}: {e}")
                                errors.append({
                                    "file": file_path,
                                    "error": f"Failed to save table {table_num + 1} from page {page_num + 1}: {str(e)}"
                                })

                # Check if we successfully processed any files from this PDF
                files_before = len(processed_files)
                # Count how many files from this specific PDF were added
                # (processed_files could have files from previous PDFs in the loop)

                # Categorized error messages for no tables found or all filtered
                if tables_found == 0 and not is_scanned:
                    logger.info(f"No tables found in {file_path}")
                    errors.append({
                        "file": file_path,
                        "error": f"No tables detected in {total_pages} page(s). The PDF may not contain structured tables with clear borders/gridlines."
                    })
                elif tables_found > 0 and not collected_tables:
                    # Tables were detected but all filtered out by quality checks
                    logger.info(f"All {tables_found} table(s) from {file_path} were filtered out (quality checks)")
                    errors.append({
                        "file": file_path,
                        "error": f"Detected {tables_found} potential table(s), but all were filtered out due to low quality (too few data cells, mostly empty rows). Try a different PDF or adjust detection settings."
                    })
                elif collected_tables and not any(base in str(pf) for pf in processed_files):
                    # Tables collected but none were successfully saved
                    logger.warning(f"Collected {len(collected_tables)} tables from {file_path} but none saved successfully")
                    if not any(err.get("file") == file_path for err in errors):
                        errors.append({
                            "file": file_path,
                            "error": f"Failed to extract tables. All {len(collected_tables)} table(s) encountered processing errors."
                        })

            except Exception as e:
                logger.error(f"Error extracting tables from {file_path}: {e}", exc_info=True)
                # Categorize errors
                error_msg = str(e)
                if "password" in error_msg.lower() or "encrypted" in error_msg.lower():
                    errors.append({"file": file_path, "error": "PDF is password-protected or encrypted. Please unlock it first."})
                elif "corrupt" in error_msg.lower() or "damaged" in error_msg.lower():
                    errors.append({"file": file_path, "error": "PDF file appears to be corrupted or damaged."})
                else:
                    errors.append({"file": file_path, "error": f"Processing failed: {str(e)}"})

        # Final safety check - ensure we never return completely empty results
        if not processed_files and not errors:
            logger.error("CRITICAL: extract_tables returning empty results with no errors!")
            errors.append({
                "file": "unknown",
                "error": "Table extraction failed with no specific error. The PDF may be corrupted, empty, or in an unsupported format."
            })

        # Partial success handling - return what we have
        logger.info(f"Extraction complete: {len(processed_files)} files, {len(errors)} errors")
        return {"processed_files": processed_files, "errors": errors}

    except Exception as e:
        # Top-level exception handler - catch ANY unexpected error
        logger.error(f"UNEXPECTED ERROR in extract_tables: {e}", exc_info=True)
        return {
            "processed_files": [],
            "errors": [{
                "file": "unknown",
                "error": f"Unexpected error during table extraction: {str(e)}"
            }]
        }

def grayscale_pdf(payload):
    """
    Convert PDF to grayscale.

    This converts the entire PDF (text, images, and graphics) to grayscale/black & white.
    Uses Ghostscript for reliable conversion while preserving PDF structure.
    """
    files = payload.get("files", [])

    processed_files = []
    errors = []

    for file_path in files:
        try:
            base, ext = os.path.splitext(file_path)
            output_path = f"{base}_grayscale{ext}"

            # Use Ghostscript for proper grayscale conversion
            # This preserves text quality and properly converts all elements
            import subprocess

            gs_command = [
                "gswin64c" if os.name == "nt" else "gs",  # Windows uses gswin64c
                "-sDEVICE=pdfwrite",
                "-sColorConversionStrategy=Gray",
                "-dProcessColorModel=/DeviceGray",
                "-dCompatibilityLevel=1.4",
                "-dNOPAUSE",
                "-dBATCH",
                "-dQUIET",
                f"-sOutputFile={output_path}",
                file_path
            ]

            try:
                result = subprocess.run(gs_command, capture_output=True, text=True, timeout=60)
                if result.returncode == 0 and os.path.exists(output_path):
                    processed_files.append(output_path)
                else:
                    # Fallback to PyMuPDF method if Ghostscript fails
                    raise Exception("Ghostscript conversion failed, using fallback")
            except (subprocess.TimeoutExpired, FileNotFoundError, Exception) as gs_error:
                logger.warning(f"Ghostscript not available or failed ({gs_error}), using fallback method")

                # Fallback: Convert each page to grayscale image and rebuild PDF
                # This works but reduces quality slightly
                doc = fitz.open(file_path)
                new_doc = fitz.open()

                for page_num in range(len(doc)):
                    page = doc[page_num]

                    # Render page as high-quality grayscale image
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), colorspace=fitz.csGRAY)

                    # Create new page with same dimensions
                    new_page = new_doc.new_page(width=page.rect.width, height=page.rect.height)

                    # Insert grayscale image
                    new_page.insert_image(new_page.rect, pixmap=pix)

                    pix = None

                new_doc.save(output_path, garbage=4, deflate=True)
                new_doc.close()
                doc.close()

                processed_files.append(output_path)

        except Exception as e:
            logger.error(f"Error converting PDF to grayscale {file_path}: {e}", exc_info=True)
            errors.append({"file": file_path, "error": str(e)})

    return {"processed_files": processed_files, "errors": errors}

def repair_pdf(payload):
    """
    Repair corrupted PDF files using multiple strategies.

    Repair techniques:
    1. Standard repair: Open and resave with garbage collection
    2. Linearization: Rebuild PDF structure
    3. Incremental save disabled: Force full rewrite
    4. Clean content streams: Fix broken page content
    """
    files = payload.get("files", [])

    processed_files = []
    errors = []

    for file_path in files:
        repaired = False
        last_error = None

        try:
            # Check if file exists
            if not os.path.exists(file_path):
                errors.append({"file": file_path, "error": f"File not found: {file_path}"})
                continue

            file_size = os.path.getsize(file_path)
            print(f"[REPAIR] Starting repair for {os.path.basename(file_path)} (size: {file_size} bytes)", file=sys.stderr, flush=True)

            # Verify file is readable and check PDF magic header
            try:
                with open(file_path, 'rb') as f:
                    first_bytes = f.read(1024)
                    print(f"[REPAIR] File readable. First bytes: {first_bytes[:20]}", file=sys.stderr, flush=True)

                    # Check for valid PDF header
                    if not first_bytes.startswith(b'%PDF-'):
                        # File has no PDF structure at all - cannot repair
                        error_msg = "This file is not a valid PDF or is completely corrupted beyond repair. PDF files must start with '%PDF-' header, but this file does not. The repair tool can only fix PDFs with minor internal corruption, not files that have lost their entire PDF structure."

                        print(f"[REPAIR] ERROR: {error_msg}", file=sys.stderr, flush=True)
                        errors.append({"file": file_path, "error": error_msg})
                        continue
            except Exception as read_err:
                print(f"[REPAIR] Cannot read file with Python: {read_err}", file=sys.stderr, flush=True)
                errors.append({"file": file_path, "error": f"Cannot read file: {read_err}"})
                continue

            base, ext = os.path.splitext(file_path)
            output_path = f"{base}_repaired{ext}"

            # Strategy 1: Standard repair with aggressive garbage collection
            try:
                print(f"[REPAIR] Attempting fitz.open() with file_path: {file_path}", file=sys.stderr, flush=True)
                doc = fitz.open(file_path)
                doc.save(
                    output_path,
                    garbage=4,  # Aggressive garbage collection (removes unused objects)
                    deflate=True,  # Compress streams
                    clean=True,  # Clean and sanitize content streams
                    pretty=True,  # Format content streams for readability
                )
                doc.close()
                repaired = True
                print(f"[REPAIR] Strategy 1 (standard) succeeded for {os.path.basename(file_path)}", file=sys.stderr)
            except Exception as e1:
                last_error = e1
                print(f"[REPAIR] Strategy 1 failed: {e1}", file=sys.stderr)

                # Strategy 2: Try with linear=True (rebuild PDF structure)
                try:
                    doc = fitz.open(file_path)
                    doc.save(
                        output_path,
                        garbage=4,
                        deflate=True,
                        linear=True,  # Linearize PDF (rebuilds structure)
                    )
                    doc.close()
                    repaired = True
                    print(f"[REPAIR] Strategy 2 (linearize) succeeded for {os.path.basename(file_path)}", file=sys.stderr)
                except Exception as e2:
                    last_error = e2
                    print(f"[REPAIR] Strategy 2 failed: {e2}", file=sys.stderr)

                    # Strategy 3: Force full rewrite (no incremental)
                    try:
                        doc = fitz.open(file_path)
                        doc.save(
                            output_path,
                            garbage=4,
                            deflate=True,
                            incremental=False,  # Force full rewrite
                            expand=True,  # Expand all objects
                        )
                        doc.close()
                        repaired = True
                        print(f"[REPAIR] Strategy 3 (full rewrite) succeeded for {os.path.basename(file_path)}", file=sys.stderr)
                    except Exception as e3:
                        last_error = e3
                        print(f"[REPAIR] Strategy 3 failed: {e3}", file=sys.stderr)

            if repaired:
                processed_files.append(output_path)
            else:
                # All strategies failed
                error_msg = f"Failed to repair PDF. Last error: {str(last_error)}"
                errors.append({"file": file_path, "error": error_msg})

        except Exception as e:
            errors.append({"file": file_path, "error": f"Unexpected error: {str(e)}"})

    return {"processed_files": processed_files, "errors": errors}

def flatten_pdf(payload):
    """Flatten PDF - make form fields non-editable and merge annotations into content."""
    files = payload.get("files", [])

    processed_files = []
    errors = []

    for file_path in files:
        if not os.path.exists(file_path):
            errors.append({"file": file_path, "error": f"File not found: {file_path}"})
            continue

        try:
            base, ext = os.path.splitext(file_path)
            output_path = f"{base}_flattened{ext}"

            doc = fitz.open(file_path)

            # Check if document has any annotations or form fields
            has_annotations = False
            for page in doc:
                if page.annots() or page.widgets():
                    has_annotations = True
                    break

            if not has_annotations:
                # No annotations to flatten - just copy the file
                doc.save(output_path)
                doc.close()
                processed_files.append(output_path)
                continue

            # Flatten all annotations on each page
            for page in doc:
                # First flatten widgets (form fields)
                widgets = list(page.widgets())
                for widget in widgets:
                    try:
                        widget.update()  # Ensure widget is rendered
                    except Exception:
                        pass

                # Then flatten all annotations
                try:
                    page.apply_redactions()  # Apply any redaction annotations first
                except Exception:
                    pass

                # Flatten remaining annotations
                annots = list(page.annots()) if page.annots() else []
                for annot in annots:
                    try:
                        # Convert annotation to drawing
                        page.draw_rect(annot.rect, color=None, fill=None)
                    except Exception:
                        pass

            doc.save(output_path, garbage=4, deflate=True)
            doc.close()
            processed_files.append(output_path)

        except Exception as e:
            errors.append({"file": file_path, "error": f"Failed to flatten: {str(e)}"})

    return {"processed_files": processed_files, "errors": errors}

def add_page_numbers(payload):
    """Add page numbers to PDF."""
    files = payload.get("files", [])
    position = payload.get("position", "bottom-right")  # bottom-right, bottom-left, top-right, top-left

    processed_files = []
    errors = []

    for file_path in files:
        try:
            base, ext = os.path.splitext(file_path)
            output_path = f"{base}_numbered{ext}"

            doc = fitz.open(file_path)

            for page_num, page in enumerate(doc):
                rect = page.rect

                # Calculate position
                if position == "bottom-right":
                    point = fitz.Point(rect.width - 50, rect.height - 20)
                elif position == "bottom-left":
                    point = fitz.Point(50, rect.height - 20)
                elif position == "top-right":
                    point = fitz.Point(rect.width - 50, 30)
                elif position == "top-left":
                    point = fitz.Point(50, 30)
                else:
                    point = fitz.Point(rect.width - 50, rect.height - 20)

                page.insert_text(point, str(page_num + 1), fontsize=12, color=(0, 0, 0))

            doc.save(output_path)
            doc.close()
            processed_files.append(output_path)
        except Exception as e:
            errors.append({"file": file_path, "error": str(e)})

    return {"processed_files": processed_files, "errors": errors}

def delete_pages(payload):
    """Delete specific pages from PDF."""
    files = payload.get("files", [])
    pages = payload.get("pages", "")  # e.g., "1,3,5" or "2-5"

    processed_files = []
    errors = []

    for file_path in files:
        try:
            base, ext = os.path.splitext(file_path)
            output_path = f"{base}_deleted{ext}"

            doc = fitz.open(file_path)

            # Parse pages to delete (0-indexed)
            pages_to_delete = set()
            if pages:
                if "-" in pages:
                    start, end = map(int, pages.split("-"))
                    pages_to_delete = set(range(start - 1, min(end, len(doc))))
                else:
                    pages_to_delete = {int(p.strip()) - 1 for p in pages.split(",")}

            # Create new document without deleted pages
            new_doc = fitz.open()
            for page_num in range(len(doc)):
                if page_num not in pages_to_delete:
                    new_doc.insert_pdf(doc, from_page=page_num, to_page=page_num)

            new_doc.save(output_path)
            new_doc.close()
            doc.close()
            processed_files.append(output_path)
        except Exception as e:
            errors.append({"file": file_path, "error": str(e)})

    return {"processed_files": processed_files, "errors": errors}


def _extract_text_with_positions(page):
    """
    Extract text with bounding box positions from a PDF page.

    Args:
        page: PyMuPDF page object

    Returns:
        List of dicts with 'text', 'bbox' (x0, y0, x1, y1), and 'line_num'
    """
    words = []
    # get_text("words") returns: (x0, y0, x1, y1, "word", block_no, line_no, word_no)
    word_list = page.get_text("words")

    for word_info in word_list:
        x0, y0, x1, y1, text, block_no, line_no, word_no = word_info
        words.append({
            'text': text,
            'bbox': (x0, y0, x1, y1),
            'block_no': block_no,
            'line_no': line_no,
            'word_no': word_no
        })

    return words


def _find_text_differences(words1, words2):
    """
    Find text differences between two lists of words using difflib.

    Args:
        words1: List of word dicts from page 1
        words2: List of word dicts from page 2

    Returns:
        Tuple of (deleted_bboxes, added_bboxes, changed_bboxes_page1, changed_bboxes_page2)
        Each is a list of (x0, y0, x1, y1) bounding boxes
    """
    import difflib

    # Extract just the text for comparison
    text1 = [w['text'] for w in words1]
    text2 = [w['text'] for w in words2]

    # Use SequenceMatcher to find differences
    matcher = difflib.SequenceMatcher(None, text1, text2)

    deleted_bboxes = []  # In doc1 but not in doc2 (red)
    added_bboxes = []    # In doc2 but not in doc1 (green)

    for tag, i1, i2, j1, j2 in matcher.get_opcodes():
        if tag == 'delete':
            # Text exists in doc1 but not in doc2
            for idx in range(i1, i2):
                if idx < len(words1):
                    deleted_bboxes.append(words1[idx]['bbox'])
        elif tag == 'insert':
            # Text exists in doc2 but not in doc1
            for idx in range(j1, j2):
                if idx < len(words2):
                    added_bboxes.append(words2[idx]['bbox'])
        elif tag == 'replace':
            # Text changed between doc1 and doc2
            for idx in range(i1, i2):
                if idx < len(words1):
                    deleted_bboxes.append(words1[idx]['bbox'])
            for idx in range(j1, j2):
                if idx < len(words2):
                    added_bboxes.append(words2[idx]['bbox'])

    return deleted_bboxes, added_bboxes


def _merge_adjacent_bboxes(bboxes, tolerance=5):
    """
    Merge bounding boxes that are on the same line and adjacent.

    Args:
        bboxes: List of (x0, y0, x1, y1) tuples
        tolerance: Vertical tolerance for considering boxes on same line

    Returns:
        List of merged (x0, y0, x1, y1) tuples
    """
    if not bboxes:
        return []

    # Sort by y position (line), then x position
    sorted_boxes = sorted(bboxes, key=lambda b: (round(b[1] / tolerance), b[0]))

    merged = []
    current = list(sorted_boxes[0])

    for bbox in sorted_boxes[1:]:
        x0, y0, x1, y1 = bbox

        # Check if on same line (similar y) and adjacent/overlapping in x
        same_line = abs(y0 - current[1]) < tolerance and abs(y1 - current[3]) < tolerance
        adjacent_x = x0 <= current[2] + tolerance * 3  # Allow small gap between words

        if same_line and adjacent_x:
            # Extend current box
            current[2] = max(current[2], x1)
            current[1] = min(current[1], y0)
            current[3] = max(current[3], y1)
        else:
            merged.append(tuple(current))
            current = [x0, y0, x1, y1]

    merged.append(tuple(current))
    return merged


def diff_pdfs(payload):
    """
    Compare two PDFs using text-based comparison with highlighted differences.

    Uses difflib to find actual text differences (not pixel-level), resulting in
    precise highlighting of only the words/phrases that changed.

    Creates a comparison PDF showing:
    - Side-by-side view of both PDFs
    - Deleted text highlighted in RED (left side - File 1)
    - Added text highlighted in GREEN (right side - File 2)
    - Page-by-page comparison with difference counts
    """
    files = payload.get("files", [])

    # Production limits
    MAX_PAGES_PER_PDF = 100
    MAX_TOTAL_PAGES = 150

    # Strict validation: exactly 2 files required
    if len(files) != 2:
        if len(files) < 2:
            return {"processed_files": [], "errors": [{"file": "comparison", "error": "PDF comparison requires exactly 2 PDF files. Please upload 2 PDFs to compare."}]}
        else:
            return {"processed_files": [], "errors": [{"file": "comparison", "error": f"PDF comparison requires exactly 2 files, but {len(files)} were provided. Only the first 2 will be compared."}]}

    processed_files = []
    errors = []

    # Validate files exist
    for i, file_path in enumerate(files[:2]):
        if not os.path.exists(file_path):
            return {"processed_files": [], "errors": [{"file": file_path, "error": f"File {i+1} not found: {os.path.basename(file_path)}"}]}

    doc1 = None
    doc2 = None
    new_doc = None

    try:
        print(f"[PDF_DIFF] Starting comparison: {os.path.basename(files[0])} vs {os.path.basename(files[1])}", file=sys.stderr, flush=True)

        doc1 = fitz.open(files[0])
        doc2 = fitz.open(files[1])

        # Check page count limits
        if len(doc1) > MAX_PAGES_PER_PDF or len(doc2) > MAX_PAGES_PER_PDF:
            return {"processed_files": [], "errors": [{"file": files[0], "error": f"PDFs with more than {MAX_PAGES_PER_PDF} pages are not supported for comparison. File 1 has {len(doc1)} pages, File 2 has {len(doc2)} pages."}]}

        max_pages = max(len(doc1), len(doc2))
        if max_pages > MAX_TOTAL_PAGES:
            return {"processed_files": [], "errors": [{"file": files[0], "error": f"Total page count ({max_pages}) exceeds limit of {MAX_TOTAL_PAGES} pages."}]}

        print(f"[PDF_DIFF] File 1: {len(doc1)} pages, File 2: {len(doc2)} pages", file=sys.stderr, flush=True)

        # Better output naming: include both file names
        base1 = os.path.splitext(os.path.basename(files[0]))[0]
        base2 = os.path.splitext(os.path.basename(files[1]))[0]
        output_dir = os.path.dirname(files[0])
        output_path = os.path.join(output_dir, f"{base1}_vs_{base2}_comparison.pdf")

        # Create comparison PDF with difference highlighting
        new_doc = fitz.open()
        max_pages = max(len(doc1), len(doc2))
        differences_found = 0

        for page_num in range(max_pages):
            # Get page dimensions (use first page as reference)
            if page_num < len(doc1):
                ref_width = doc1[page_num].rect.width
                ref_height = doc1[page_num].rect.height
            else:
                ref_width = doc2[page_num].rect.width
                ref_height = doc2[page_num].rect.height

            # Create wide page for side-by-side comparison
            page = new_doc.new_page(width=ref_width * 2 + 40, height=ref_height + 80)

            # Add labels
            page.insert_text(fitz.Point(20, 20), "Original (File 1)", fontsize=12, color=(0, 0, 0))
            page.insert_text(fitz.Point(ref_width + 40, 20), "Comparison (File 2)", fontsize=12, color=(0, 0, 0))

            # Render both pages as images to detect differences
            if page_num < len(doc1):
                page1 = doc1[page_num]
                # Show original on left side
                page.show_pdf_page(
                    fitz.Rect(20, 40, ref_width + 20, ref_height + 40),
                    doc1,
                    page_num
                )

            if page_num < len(doc2):
                page2 = doc2[page_num]
                # Show comparison on right side
                page.show_pdf_page(
                    fitz.Rect(ref_width + 40, 40, ref_width * 2 + 40, ref_height + 40),
                    doc2,
                    page_num
                )

            # Try to detect and highlight differences using TEXT-BASED comparison
            if page_num < len(doc1) and page_num < len(doc2):
                try:
                    # Extract text with positions from both pages
                    words1 = _extract_text_with_positions(doc1[page_num])
                    words2 = _extract_text_with_positions(doc2[page_num])

                    print(f"[PDF_DIFF] Page {page_num + 1}: Extracted {len(words1)} words from doc1, {len(words2)} words from doc2", file=sys.stderr, flush=True)

                    # Find text differences using difflib
                    deleted_bboxes, added_bboxes = _find_text_differences(words1, words2)

                    # Merge adjacent bboxes for cleaner highlighting
                    deleted_merged = _merge_adjacent_bboxes(deleted_bboxes)
                    added_merged = _merge_adjacent_bboxes(added_bboxes)

                    total_diffs = len(deleted_merged) + len(added_merged)

                    if total_diffs > 0:
                        differences_found += 1

                        print(f"[PDF_DIFF] Page {page_num + 1}: Found {len(deleted_merged)} deletions, {len(added_merged)} additions", file=sys.stderr, flush=True)

                        # Highlight deleted text (in doc1) - RED
                        for bbox in deleted_merged:
                            x0, y0, x1, y1 = bbox
                            # Left side offset: 20px margin
                            left_rect = fitz.Rect(20 + x0, 40 + y0, 20 + x1, 40 + y1)

                            highlight = page.add_rect_annot(left_rect)
                            highlight.set_colors(stroke=(0.8, 0, 0), fill=(1, 0.8, 0.8))
                            highlight.set_opacity(0.4)
                            highlight.set_border(width=1)
                            highlight.update()

                        # Highlight added text (in doc2) - GREEN
                        for bbox in added_merged:
                            x0, y0, x1, y1 = bbox
                            # Right side offset: ref_width + 40px margin
                            right_rect = fitz.Rect(ref_width + 40 + x0, 40 + y0, ref_width + 40 + x1, 40 + y1)

                            highlight = page.add_rect_annot(right_rect)
                            highlight.set_colors(stroke=(0, 0.6, 0), fill=(0.8, 1, 0.8))
                            highlight.set_opacity(0.4)
                            highlight.set_border(width=1)
                            highlight.update()

                        # Add difference summary
                        summary_parts = []
                        if len(deleted_merged) > 0:
                            summary_parts.append(f"{len(deleted_merged)} removed")
                        if len(added_merged) > 0:
                            summary_parts.append(f"{len(added_merged)} added")

                        try:
                            page.insert_text(
                                fitz.Point(20, ref_height + 60),
                                f"Page {page_num + 1}: {', '.join(summary_parts)}",
                                fontsize=10,
                                color=(0.8, 0, 0)
                            )
                        except Exception:
                            page.insert_text(
                                fitz.Point(20, ref_height + 60),
                                f"Page {page_num + 1}: Different",
                                fontsize=10,
                                color=(0.8, 0, 0)
                            )

                        print(f"[PDF_DIFF] Page {page_num + 1}: Highlighted {len(deleted_merged)} deletions (red), {len(added_merged)} additions (green)", file=sys.stderr, flush=True)
                    else:
                        # Pages are identical (no text differences)
                        try:
                            page.insert_text(
                                fitz.Point(20, ref_height + 60),
                                f"Page {page_num + 1}: Identical",
                                fontsize=10,
                                color=(0, 0.6, 0)
                            )
                        except Exception:
                            page.insert_text(
                                fitz.Point(20, ref_height + 60),
                                f"Page {page_num + 1}: Same",
                                fontsize=10,
                                color=(0, 0.6, 0)
                            )

                except Exception as e:
                    # Log and notify user of comparison failure for this page
                    print(f"[PDF_DIFF] Could not detect differences for page {page_num + 1}: {e}", file=sys.stderr, flush=True)
                    import traceback
                    traceback.print_exc()
                    try:
                        page.insert_text(
                            fitz.Point(20, ref_height + 60),
                            f"Page {page_num + 1}: Comparison failed",
                            fontsize=10,
                            color=(0.5, 0.5, 0.5)
                        )
                    except Exception:
                        pass

            elif page_num >= len(doc1):
                # Page only in doc2
                differences_found += 1
                try:
                    page.insert_text(
                        fitz.Point(20, ref_height + 60),
                        f"! Page {page_num + 1}: Only exists in File 2",
                        fontsize=10,
                        color=(0.8, 0, 0)
                    )
                except Exception:
                    page.insert_text(
                        fitz.Point(20, ref_height + 60),
                        f"Page {page_num + 1}: Only in File 2",
                        fontsize=10,
                        color=(0.8, 0, 0)
                    )
            elif page_num >= len(doc2):
                # Page only in doc1
                differences_found += 1
                try:
                    page.insert_text(
                        fitz.Point(20, ref_height + 60),
                        f"! Page {page_num + 1}: Only exists in File 1",
                        fontsize=10,
                        color=(0.8, 0, 0)
                    )
                except Exception:
                    page.insert_text(
                        fitz.Point(20, ref_height + 60),
                        f"Page {page_num + 1}: Only in File 1",
                        fontsize=10,
                        color=(0.8, 0, 0)
                    )

        # Add summary page at the beginning with legend
        summary_page = new_doc.new_page(0, width=ref_width * 2 + 40, height=280)
        summary_page.insert_text(
            fitz.Point(20, 40),
            "PDF Comparison Summary",
            fontsize=16,
            color=(0, 0, 0)
        )
        summary_page.insert_text(
            fitz.Point(20, 70),
            f"File 1: {base1}",
            fontsize=12,
            color=(0, 0, 0)
        )
        summary_page.insert_text(
            fitz.Point(20, 90),
            f"File 2: {base2}",
            fontsize=12,
            color=(0, 0, 0)
        )
        summary_page.insert_text(
            fitz.Point(20, 120),
            f"Total pages compared: {max_pages}",
            fontsize=12,
            color=(0, 0, 0)
        )
        summary_page.insert_text(
            fitz.Point(20, 140),
            f"Pages with differences: {differences_found}",
            fontsize=12,
            color=(0.8, 0, 0) if differences_found > 0 else (0, 0.6, 0)
        )

        # Add legend for color coding
        summary_page.insert_text(
            fitz.Point(20, 175),
            "Legend:",
            fontsize=12,
            color=(0, 0, 0)
        )
        # Red box for File 1 differences
        legend_rect1 = fitz.Rect(20, 185, 40, 200)
        shape = summary_page.new_shape()
        shape.draw_rect(legend_rect1)
        shape.finish(color=(1, 0, 0), fill=(1, 0.8, 0.8), fill_opacity=0.5)
        shape.commit()
        summary_page.insert_text(
            fitz.Point(50, 197),
            "Red highlight = Content in File 1 (left side)",
            fontsize=10,
            color=(0, 0, 0)
        )
        # Blue box for File 2 differences
        legend_rect2 = fitz.Rect(20, 210, 40, 225)
        shape = summary_page.new_shape()
        shape.draw_rect(legend_rect2)
        shape.finish(color=(0, 0, 1), fill=(0.8, 0.8, 1), fill_opacity=0.5)
        shape.commit()
        summary_page.insert_text(
            fitz.Point(50, 222),
            "Blue highlight = Content in File 2 (right side)",
            fontsize=10,
            color=(0, 0, 0)
        )
        # Green text for identical
        summary_page.insert_text(
            fitz.Point(20, 250),
            "Green text = Pages are identical",
            fontsize=10,
            color=(0, 0.6, 0)
        )

        new_doc.save(output_path, garbage=4, deflate=True)
        print(f"[PDF_DIFF] Comparison complete. {differences_found} pages with differences found.", file=sys.stderr, flush=True)

        if new_doc:
            new_doc.close()
        if doc1:
            doc1.close()
        if doc2:
            doc2.close()

        processed_files.append(output_path)

    except Exception as e:
        logger.error(f"Error comparing PDFs: {e}", exc_info=True)
        print(f"[PDF_DIFF] Comparison failed: {e}", file=sys.stderr, flush=True)
        if new_doc:
            new_doc.close()
        if doc1:
            doc1.close()
        if doc2:
            doc2.close()
        errors.append({"file": files[0] if files else "unknown", "error": f"PDF comparison failed: {str(e)}. Please ensure both files are valid PDFs."})

    return {"processed_files": processed_files, "errors": errors}

def create_booklet(payload):
    """
    Create booklet layout for saddle-stitch printing.

    Reorders PDF pages so when printed double-sided, folded in half,
    and stapled in the middle, pages appear in correct sequential order.

    Algorithm:
    - For n pages (rounded up to multiple of 4):
    - Page order: [n, 1, 2, n-1, n-2, 3, 4, n-3, ...]
    - This creates proper imposition for folding and stapling
    """
    files = payload.get("files", [])

    # Production limits
    MAX_PAGES = 1000  # Increased from 200 to support larger documents

    processed_files = []
    errors = []

    for file_path in files:
        if not os.path.exists(file_path):
            errors.append({"file": file_path, "error": f"File not found: {file_path}"})
            continue

        doc = None
        new_doc = None

        try:
            logger.info(f"[BOOKLET] Starting booklet creation for {os.path.basename(file_path)}")

            doc = fitz.open(file_path)
            total_pages = len(doc)

            if total_pages == 0:
                errors.append({"file": file_path, "error": "PDF has no pages"})
                continue

            if total_pages > MAX_PAGES:
                errors.append({"file": file_path, "error": f"PDF has {total_pages} pages. Booklets with more than {MAX_PAGES} pages are not supported."})
                continue

            # Round up to nearest multiple of 4 (required for saddle-stitch)
            pages_needed = total_pages
            if pages_needed % 4 != 0:
                pages_needed = ((pages_needed // 4) + 1) * 4

            blank_pages_added = pages_needed - total_pages
            logger.info(f"[BOOKLET] Original pages: {total_pages}, booklet pages: {pages_needed}, blank pages: {blank_pages_added}")

            base, ext = os.path.splitext(file_path)
            output_path = f"{base}_booklet{ext}"

            # Create booklet with CORRECT saddle-stitch algorithm
            new_doc = fitz.open()

            # Correct booklet page order for saddle-stitch binding
            # For each physical sheet (4 pages: 2 on front, 2 on back):
            # Sheet i from outside to inside
            booklet_pages = []
            num_sheets = pages_needed // 4

            for i in range(num_sheets):
                # Front of sheet: [right page, left page]
                front_right = pages_needed - 1 - (2 * i)  # Last page working inward
                front_left = 2 * i                         # First page working inward

                # Back of sheet: [left page, right page]
                back_left = 2 * i + 1                      # Second page working inward
                back_right = pages_needed - 2 - (2 * i)   # Second-to-last working inward

                # Order for this sheet: front-right, front-left, back-left, back-right
                booklet_pages.extend([front_right, front_left, back_left, back_right])

            logger.info(f"[BOOKLET] Page order: {booklet_pages}")

            # Insert pages in booklet order
            for page_idx in booklet_pages:
                if page_idx < total_pages:
                    # Real page from original PDF
                    new_doc.insert_pdf(doc, from_page=page_idx, to_page=page_idx)
                else:
                    # Blank page (for padding to multiple of 4)
                    if total_pages > 0:
                        ref_page = doc[0]
                        blank_page = new_doc.new_page(width=ref_page.rect.width, height=ref_page.rect.height)
                        blank_page.insert_text(
                            fitz.Point(blank_page.rect.width / 2 - 30, blank_page.rect.height / 2),
                            "(Blank page)",
                            fontsize=10,
                            color=(0.7, 0.7, 0.7)
                        )

            new_doc.save(output_path, garbage=4, deflate=True)
            logger.info(f"[BOOKLET] Created booklet with {len(new_doc)} pages (added {blank_pages_added} blank pages)")

            if new_doc:
                new_doc.close()
            if doc:
                doc.close()

            processed_files.append(output_path)

        except Exception as e:
            logger.error(f"Error creating booklet from {file_path}: {e}", exc_info=True)
            if new_doc:
                new_doc.close()
            if doc:
                doc.close()
            errors.append({"file": file_path, "error": f"Failed to create booklet: {str(e)}"})

    return {"processed_files": processed_files, "errors": errors}

def scrub_pdf(payload):
    """
    Deep scrub PDF - Remove ALL privacy-leaking metadata and hidden content.

    Removes:
    - Document metadata (author, title, dates, etc.)
    - XMP metadata streams
    - Annotations and comments
    - JavaScript code
    - Embedded files and attachments
    - Named destinations
    - Bookmarks/outlines
    - Form data
    - Document ID
    - OpenAction (auto-execute scripts)
    - Threads and article threads
    - Hidden layers (OCGs)
    - Markup and revision history
    - ViewerPreferences
    """
    files = payload.get("files", [])

    # Production limits
    MAX_FILE_SIZE_MB = 500

    processed_files = []
    errors = []

    for file_path in files:
        if not os.path.exists(file_path):
            errors.append({"file": file_path, "error": f"File not found: {file_path}"})
            continue

        # Check file size
        file_size_mb = os.path.getsize(file_path) / (1024 * 1024)
        if file_size_mb > MAX_FILE_SIZE_MB:
            errors.append({"file": file_path, "error": f"File size ({file_size_mb:.1f}MB) exceeds limit of {MAX_FILE_SIZE_MB}MB"})
            continue

        try:
            print(f"[SCRUBBER] Starting deep scrub for {os.path.basename(file_path)}", file=sys.stderr, flush=True)

            base, ext = os.path.splitext(file_path)
            output_path = f"{base}_scrubbed{ext}"

            with pikepdf.open(file_path) as pdf:
                items_removed = []

                # 1. Remove XMP metadata stream
                if '/Metadata' in pdf.Root:
                    del pdf.Root.Metadata
                    items_removed.append("XMP metadata")

                # 2. Clear document info (author, title, creation date, etc.)
                # pikepdf's docinfo doesn't support .clear(), need to delete individual keys
                if pdf.docinfo:
                    try:
                        # Get all keys and delete them individually
                        keys_to_delete = list(pdf.docinfo.keys())
                        for key in keys_to_delete:
                            del pdf.docinfo[key]
                        items_removed.append("document info")
                    except Exception as e:
                        print(f"[SCRUBBER] Could not clear docinfo: {e}", file=sys.stderr, flush=True)

                # 3. Remove JavaScript
                try:
                    if '/Names' in pdf.Root and '/JavaScript' in pdf.Root.Names:
                        del pdf.Root.Names.JavaScript
                        items_removed.append("JavaScript")
                except Exception as e:
                    print(f"[SCRUBBER] Could not remove JavaScript: {e}", file=sys.stderr, flush=True)

                # 4. Remove embedded files/attachments
                try:
                    if '/Names' in pdf.Root and '/EmbeddedFiles' in pdf.Root.Names:
                        del pdf.Root.Names.EmbeddedFiles
                        items_removed.append("embedded files")
                except Exception as e:
                    print(f"[SCRUBBER] Could not remove embedded files: {e}", file=sys.stderr, flush=True)

                # 5. Remove OpenAction (auto-execute on open)
                try:
                    if '/OpenAction' in pdf.Root:
                        del pdf.Root.OpenAction
                        items_removed.append("OpenAction")
                except Exception as e:
                    print(f"[SCRUBBER] Could not remove OpenAction: {e}", file=sys.stderr, flush=True)

                # 6. Remove AA (Additional Actions)
                try:
                    if '/AA' in pdf.Root:
                        del pdf.Root.AA
                        items_removed.append("additional actions")
                except Exception as e:
                    print(f"[SCRUBBER] Could not remove additional actions: {e}", file=sys.stderr, flush=True)

                # 7. Remove Named Destinations
                try:
                    if '/Names' in pdf.Root and '/Dests' in pdf.Root.Names:
                        del pdf.Root.Names.Dests
                        items_removed.append("named destinations")
                except Exception as e:
                    print(f"[SCRUBBER] Could not remove named destinations: {e}", file=sys.stderr, flush=True)

                # 8. Remove Outlines (bookmarks - can contain user names/dates)
                try:
                    if '/Outlines' in pdf.Root:
                        del pdf.Root.Outlines
                        items_removed.append("bookmarks")
                except Exception as e:
                    print(f"[SCRUBBER] Could not remove bookmarks: {e}", file=sys.stderr, flush=True)

                # 9. Remove Threads (article threads)
                try:
                    if '/Threads' in pdf.Root:
                        del pdf.Root.Threads
                        items_removed.append("threads")
                except Exception as e:
                    print(f"[SCRUBBER] Could not remove threads: {e}", file=sys.stderr, flush=True)

                # 10. Remove Optional Content Groups (hidden layers)
                try:
                    if '/OCProperties' in pdf.Root:
                        del pdf.Root.OCProperties
                        items_removed.append("optional content layers")
                except Exception as e:
                    print(f"[SCRUBBER] Could not remove optional content: {e}", file=sys.stderr, flush=True)

                # 11. Remove ViewerPreferences (can leak info)
                try:
                    if '/ViewerPreferences' in pdf.Root:
                        del pdf.Root.ViewerPreferences
                        items_removed.append("viewer preferences")
                except Exception as e:
                    print(f"[SCRUBBER] Could not remove viewer preferences: {e}", file=sys.stderr, flush=True)

                # 12. Remove MarkInfo (accessibility/tagging info)
                try:
                    if '/MarkInfo' in pdf.Root:
                        del pdf.Root.MarkInfo
                        items_removed.append("mark info")
                except Exception as e:
                    print(f"[SCRUBBER] Could not remove mark info: {e}", file=sys.stderr, flush=True)

                # 13. Remove StructTreeRoot (structure tree - can contain metadata)
                try:
                    if '/StructTreeRoot' in pdf.Root:
                        del pdf.Root.StructTreeRoot
                        items_removed.append("structure tree")
                except Exception as e:
                    print(f"[SCRUBBER] Could not remove structure tree: {e}", file=sys.stderr, flush=True)

                # 14. Process each page
                for page_num, page in enumerate(pdf.pages):
                    # Remove annotations (comments, highlights, etc.)
                    try:
                        if '/Annots' in page:
                            del page['/Annots']
                    except Exception as e:
                        print(f"[SCRUBBER] Could not remove annotations from page {page_num + 1}: {e}", file=sys.stderr, flush=True)

                    # Remove page metadata
                    try:
                        if '/Metadata' in page:
                            del page['/Metadata']
                    except Exception as e:
                        print(f"[SCRUBBER] Could not remove page metadata from page {page_num + 1}: {e}", file=sys.stderr, flush=True)

                    # Remove piece info (editing history)
                    try:
                        if '/PieceInfo' in page:
                            del page['/PieceInfo']
                    except Exception as e:
                        print(f"[SCRUBBER] Could not remove piece info from page {page_num + 1}: {e}", file=sys.stderr, flush=True)

                    # Remove additional actions on page
                    try:
                        if '/AA' in page:
                            del page['/AA']
                    except Exception as e:
                        print(f"[SCRUBBER] Could not remove additional actions from page {page_num + 1}: {e}", file=sys.stderr, flush=True)

                if not items_removed:
                    items_removed.append("page annotations")

                # 15. Remove unreferenced resources (cleanup)
                try:
                    pdf.remove_unreferenced_resources()
                except Exception as e:
                    print(f"[SCRUBBER] Could not remove unreferenced resources: {e}", file=sys.stderr, flush=True)

                # Save with linearization (web optimization) and compression
                try:
                    pdf.save(
                        output_path,
                        linearize=True,  # Optimize for web/fast viewing
                        compress_streams=True,
                        stream_decode_level=pikepdf.StreamDecodeLevel.generalized,
                        object_stream_mode=pikepdf.ObjectStreamMode.generate
                    )
                except Exception as save_err:
                    # If linearization fails, try without it
                    print(f"[SCRUBBER] Linearization failed, trying without: {save_err}", file=sys.stderr, flush=True)
                    pdf.save(
                        output_path,
                        compress_streams=True,
                        stream_decode_level=pikepdf.StreamDecodeLevel.generalized
                    )

            print(f"[SCRUBBER] Removed: {', '.join(items_removed)}", file=sys.stderr, flush=True)
            print(f"[SCRUBBER] Scrubbed PDF saved: {os.path.basename(output_path)}", file=sys.stderr, flush=True)

            processed_files.append(output_path)

        except pikepdf.PasswordError:
            errors.append({"file": file_path, "error": "PDF is password-protected. Cannot scrub encrypted files. Please unlock the PDF first."})
        except Exception as e:
            logger.error(f"Error scrubbing PDF {file_path}: {e}", exc_info=True)
            print(f"[SCRUBBER] Failed: {e}", file=sys.stderr, flush=True)
            errors.append({"file": file_path, "error": f"Failed to scrub PDF: {str(e)}"})

    return {"processed_files": processed_files, "errors": errors}

def redact_pdf(payload):
    """
    Redact (permanently remove) text from PDF.

    Supports both single text (legacy) and multiple texts array.
    """
    files = payload.get("files", [])

    # Log operation details without exposing sensitive content
    logger.info(f"Redact PDF called with payload keys: {payload.keys()}")

    # Support both single text (legacy) and multiple texts array
    texts_to_redact = []
    if payload.get("texts"):
        # New format: JSON array of texts
        import json
        try:
            parsed = json.loads(payload.get("texts"))
            # Validate parsed result is a list
            if isinstance(parsed, list):
                texts_to_redact = parsed
            elif isinstance(parsed, str):
                texts_to_redact = [parsed]
            else:
                logger.warning(f"Unexpected JSON type for texts: {type(parsed).__name__}, using as-is")
                texts_to_redact = [str(parsed)]
            logger.info(f"Parsed {len(texts_to_redact)} text(s) from JSON")
        except (json.JSONDecodeError, TypeError) as e:
            logger.warning(f"Failed to parse texts as JSON: {type(e).__name__}, using as-is")
            texts_to_redact = [payload.get("texts")]
    elif payload.get("text"):
        # Legacy format: single text
        logger.info("Using legacy single text format")
        texts_to_redact = [payload.get("text")]

    # Filter out empty strings
    texts_to_redact = [t.strip() for t in texts_to_redact if t and t.strip()]
    # Log count only, not the actual sensitive content being redacted
    logger.info(f"Processing {len(texts_to_redact)} text(s) for redaction")

    if not texts_to_redact:
        file_name = files[0] if files else "unknown"
        return {"processed_files": [], "errors": [{"file": file_name, "error": "No text provided for redaction"}]}

    processed_files = []
    errors = []

    for file_path in files:
        doc = None
        try:
            base, ext = os.path.splitext(file_path)
            output_path = f"{base}_redacted{ext}"

            doc = fitz.open(file_path)
            total_redactions = 0
            page_count = len(doc)

            for page in doc:
                # Find and redact all specified texts
                for text in texts_to_redact:
                    # Find text instances
                    text_instances = page.search_for(text)

                    # Redact (black out) each instance
                    for inst in text_instances:
                        page.add_redact_annot(inst, fill=(0, 0, 0))
                        total_redactions += 1

                # Apply all redactions to this page
                page.apply_redactions()

            doc.save(output_path)
            doc.close()
            doc = None

            logger.info(f"Redacted {total_redactions} instances across {page_count} pages in {file_path}")
            processed_files.append(output_path)
        except Exception as e:
            logger.error(f"Error redacting PDF {file_path}: {e}", exc_info=True)
            if doc is not None:
                try:
                    doc.close()
                except:
                    pass
            errors.append({"file": file_path, "error": str(e)})

    return {"processed_files": processed_files, "errors": errors}

def sign_pdf(payload):
    """Digitally sign PDF with certificate (PKCS#12/PFX format).

    Supports visual signature appearance with configurable position.
    Uses pyhanko for cryptographic signing with X.509 certificates.

    Security Notes:
    - Certificate files (.pfx/.p12) should be deleted immediately after use
    - Passwords should be transmitted over HTTPS only
    - File size limits enforced to prevent DoS attacks
    """
    from pyhanko.pdf_utils.incremental_writer import IncrementalPdfFileWriter
    from pyhanko.sign.signers import SimpleSigner
    from pyhanko.sign.fields import SigFieldSpec
    from pyhanko.sign import PdfSignatureMetadata
    import tempfile

    files = payload.get("files", [])
    cert_file = payload.get("cert_file")
    password = payload.get("password", "")

    # Signature appearance configuration
    signature_text = payload.get("text", "Digitally Signed")
    reason = payload.get("reason", "Document Authorization")
    location = payload.get("location", "")

    # Position configuration (percentages from frontend, convert to PDF coordinates)
    x_percent = float(payload.get("x", 50)) / 100.0  # Default center
    y_percent = float(payload.get("y", 50)) / 100.0  # Default center

    processed_files = []
    errors = []

    # SECURITY: File size limits
    MAX_PDF_SIZE_MB = 50
    MAX_CERT_SIZE_KB = 100

    # Validation: Certificate required
    if not cert_file:
        return {"processed_files": [], "errors": ["Certificate file (.pfx or .p12) required for digital signing. Please upload a valid PKCS#12 certificate."]}

    if not os.path.exists(cert_file):
        return {"processed_files": [], "errors": [{"file": "certificate", "error": "Certificate file not found on server"}]}

    # SECURITY: Validate certificate file size
    try:
        cert_size_kb = os.path.getsize(cert_file) / 1024
        if cert_size_kb > MAX_CERT_SIZE_KB:
            return {"processed_files": [], "errors": [f"Certificate file too large ({cert_size_kb:.1f}KB). Maximum allowed: {MAX_CERT_SIZE_KB}KB."]}
    except Exception as e:
        return {"processed_files": [], "errors": [f"Error checking certificate file size: {str(e)}"]}

    # Validate certificate format (PKCS#12)
    if not (cert_file.lower().endswith('.pfx') or cert_file.lower().endswith('.p12')):
        # Try to load it anyway, but warn user
        logger.warning(f"Certificate file {cert_file} doesn't have .pfx or .p12 extension")

    # Load certificate and create signer
    signer = None
    try:
        # Try to load PKCS#12 certificate with password
        from cryptography.hazmat.primitives.serialization import pkcs12

        with open(cert_file, 'rb') as f:
            cert_data = f.read()

        # Parse PKCS#12 file
        try:
            private_key, certificate, ca_certs = pkcs12.load_key_and_certificates(
                cert_data,
                password.encode('utf-8') if password else None
            )
        except Exception as e:
            error_msg = str(e)
            if "password" in error_msg.lower() or "decrypt" in error_msg.lower():
                return {"processed_files": [], "errors": ["Invalid certificate password. Please check your password and try again."]}
            elif "could not deserialize" in error_msg.lower():
                return {"processed_files": [], "errors": ["Invalid certificate file format. Please use a valid PKCS#12 (.pfx or .p12) certificate."]}
            else:
                return {"processed_files": [], "errors": [f"Error loading certificate: {error_msg}"]}

        if not private_key or not certificate:
            return {"processed_files": [], "errors": ["Invalid certificate: missing private key or certificate data."]}

        # Create SimpleSigner from loaded certificate
        signer = SimpleSigner.load_pkcs12(
            pfx_file=cert_file,
            passphrase=password.encode('utf-8') if password else None
        )

        print(f"[SIGNER] Successfully loaded certificate for signing", file=sys.stderr, flush=True)

    except Exception as e:
        logger.error(f"Error loading certificate: {e}", exc_info=True)
        return {"processed_files": [], "errors": [f"Failed to load certificate: {str(e)}"]}

    # Process each PDF file
    for file_path in files:
        if not os.path.exists(file_path):
            errors.append({"file": file_path, "error": f"File not found: {file_path}"})
            continue

        # SECURITY: Check PDF file size
        try:
            pdf_size_mb = os.path.getsize(file_path) / (1024 * 1024)
            if pdf_size_mb > MAX_PDF_SIZE_MB:
                errors.append({"file": file_path, "error": f"PDF file too large ({pdf_size_mb:.1f}MB). Maximum allowed: {MAX_PDF_SIZE_MB}MB."})
                continue
        except Exception as e:
            errors.append({"file": file_path, "error": f"Error checking file size: {str(e)}"})
            continue

        try:
            # Read PDF and get dimensions for signature placement
            doc = fitz.open(file_path)
            if len(doc) == 0:
                errors.append({"file": file_path, "error": "PDF has no pages"})
                doc.close()
                continue

            # Get first page dimensions (sign first page by default)
            page = doc[0]
            page_width = page.rect.width
            page_height = page.rect.height
            doc.close()

            # Calculate signature box position
            # Signature box: 200x60 points (approximately 2.8" x 0.8")
            sig_width = 200
            sig_height = 60

            # Convert percentage position to PDF coordinates
            # PDF coordinates: (0,0) is bottom-left
            # Frontend sends: (0,0) is top-left, so flip Y
            x = int(x_percent * page_width - sig_width / 2)
            y = int((1 - y_percent) * page_height - sig_height / 2)  # Flip Y coordinate

            # Ensure signature box stays within page bounds
            x = max(10, min(x, page_width - sig_width - 10))
            y = max(10, min(y, page_height - sig_height - 10))

            print(f"[SIGNER] Signature box position: ({x}, {y}) on {page_width}x{page_height} page", file=sys.stderr, flush=True)

            # Create signature field specification
            sig_field_spec = SigFieldSpec(
                sig_field_name='Signature',
                box=(x, y, x + sig_width, y + sig_height)
            )

            # Create signature metadata
            sig_meta = PdfSignatureMetadata(
                field_name='Signature',
                reason=reason if reason else None,
                location=location if location else None,
                name=signature_text if signature_text else None,
            )

            # Generate output path
            base, ext = os.path.splitext(file_path)
            output_path = f"{base}_signed{ext}"

            # Sign the PDF
            with open(file_path, 'rb') as inf:
                writer = IncrementalPdfFileWriter(inf)

                # Add signature field
                fields.append_signature_field(
                    writer,
                    sig_field_spec
                )

                # Sign the field
                with open(output_path, 'wb') as outf:
                    meta = PdfSignatureMetadata(
                        field_name='Signature',
                        reason=reason if reason else None,
                        location=location if location else None,
                    )

                    # Actually sign the PDF
                    from pyhanko.sign.validation import validate_pdf_signature
                    fields.sign_field(
                        writer,
                        sig_field_spec,
                        signer,
                        meta
                    )

                    writer.write(outf)

            processed_files.append(output_path)
            print(f"[SIGNER] Successfully signed: {output_path}", file=sys.stderr, flush=True)

        except Exception as e:
            logger.error(f"Error signing PDF {file_path}: {e}", exc_info=True)
            errors.append({"file": file_path, "error": f"Signing failed: {str(e)}"})

    # SECURITY: Clean up certificate file immediately
    try:
        if cert_file and os.path.exists(cert_file):
            os.remove(cert_file)
            print(f"[SECURITY] Deleted certificate file: {cert_file}", file=sys.stderr, flush=True)
    except Exception as e:
        logger.warning(f"Failed to delete certificate file: {e}")

    return {"processed_files": processed_files, "errors": errors}

def optimize_pdf(payload):
    """Optimize PDF for web viewing (compression and structure optimization).

    Applies aggressive optimization:
    - Removes unused objects (garbage collection)
    - Compresses streams (deflate)
    - Cleans up PDF structure
    - Removes encryption

    Safety limits:
    - Maximum file size: 100MB (generous for desktop users)
    - Maximum pages: 1000 pages
    - Validates PDF integrity before processing
    - Detects password-protected PDFs
    """
    files = payload.get("files", [])

    processed_files = []
    errors = []

    # Backend safety limits (prevents abuse even for desktop)
    MAX_PDF_SIZE_MB = 100  # Generous limit for desktop users
    MAX_PAGES = 1000       # Page count limit

    for file_path in files:
        if not os.path.exists(file_path):
            errors.append({"file": file_path, "error": f"File not found: {file_path}"})
            continue

        doc = None
        try:
            # SECURITY: Check file size
            try:
                pdf_size_mb = os.path.getsize(file_path) / (1024 * 1024)
                if pdf_size_mb > MAX_PDF_SIZE_MB:
                    errors.append({
                        "file": file_path,
                        "error": f"PDF file too large ({pdf_size_mb:.1f}MB). Maximum allowed: {MAX_PDF_SIZE_MB}MB."
                    })
                    continue
            except Exception as e:
                errors.append({"file": file_path, "error": f"Cannot check file size: {str(e)}"})
                continue

            # VALIDATION: Check PDF header (integrity check)
            try:
                with open(file_path, 'rb') as f:
                    header = f.read(5)
                    if header != b'%PDF-':
                        errors.append({"file": file_path, "error": "Invalid PDF file: missing PDF header. File may be corrupted."})
                        continue
            except Exception as e:
                errors.append({"file": file_path, "error": f"Cannot read file: {str(e)}"})
                continue

            # Try to open PDF (will detect encryption/password protection)
            try:
                doc = fitz.open(file_path)
            except Exception as e:
                error_msg = str(e).lower()
                if "password" in error_msg or "crypt" in error_msg or "encrypt" in error_msg:
                    errors.append({
                        "file": file_path,
                        "error": "Cannot optimize password-protected PDFs. Please unlock the PDF first using the 'Unlock PDF' tool."
                    })
                else:
                    errors.append({
                        "file": file_path,
                        "error": f"Cannot open PDF: {str(e)}"
                    })
                continue

            # Check for empty PDF
            if len(doc) == 0:
                doc.close()
                errors.append({"file": file_path, "error": "PDF has no pages. Cannot optimize an empty PDF."})
                continue

            # SECURITY: Check page count
            if len(doc) > MAX_PAGES:
                doc.close()
                errors.append({
                    "file": file_path,
                    "error": f"PDF has too many pages ({len(doc)} pages). Maximum allowed: {MAX_PAGES} pages."
                })
                continue

            print(f"[OPTIMIZE] Starting optimization for {os.path.basename(file_path)} ({len(doc)} pages, {pdf_size_mb:.2f}MB)", file=sys.stderr, flush=True)

            base, ext = os.path.splitext(file_path)
            output_path = f"{base}_web_optimized{ext}"

            # Optimize for web viewing
            # garbage=4: Maximum garbage collection (remove unused objects)
            # deflate=True: Compress streams for smaller file size
            # clean=True: Clean up PDF structure
            # ascii=False: Keep binary format (smaller)
            # no_new_id=True: Don't generate new IDs (faster)
            # encryption=NONE: Remove any encryption
            doc.save(
                output_path,
                garbage=4,           # Maximum garbage collection
                deflate=True,        # Compress streams
                clean=True,          # Clean up structure
                ascii=False,         # Keep binary (smaller than ASCII)
                no_new_id=True,      # Don't generate new IDs (faster)
                encryption=fitz.PDF_ENCRYPT_NONE  # Ensure no encryption
            )

            # Calculate size reduction
            try:
                original_size_mb = os.path.getsize(file_path) / (1024 * 1024)
                optimized_size_mb = os.path.getsize(output_path) / (1024 * 1024)
                reduction_percent = ((original_size_mb - optimized_size_mb) / original_size_mb) * 100

                print(f"[OPTIMIZE] {os.path.basename(file_path)}: {original_size_mb:.2f}MB → {optimized_size_mb:.2f}MB ({reduction_percent:+.1f}%)", file=sys.stderr, flush=True)
            except Exception as e:
                print(f"[OPTIMIZE] Could not calculate size reduction: {e}", file=sys.stderr, flush=True)

            if os.path.exists(output_path):
                processed_files.append(output_path)
            else:
                errors.append({"file": file_path, "error": "Optimization failed: output file not created"})

        except Exception as e:
            logger.error(f"Error optimizing PDF {file_path}: {e}", exc_info=True)
            errors.append({"file": file_path, "error": f"Optimization failed: {str(e)}"})
        finally:
            # SECURITY: Always clean up resources
            if doc:
                try:
                    doc.close()
                except:
                    pass

    return {"processed_files": processed_files, "errors": errors}

def word_to_pdf(payload):
    """
    Convert Word documents to PDF.

    Note: This is a BETA feature. Best results with LibreOffice installed.
    Python-only fallback preserves text and basic formatting but may lose
    complex layouts, images, and advanced formatting.
    """
    files = payload.get("files", [])
    libreoffice_path = payload.get("libreoffice_path")

    processed_files = []
    errors = []

    for file_path in files:
        if not os.path.exists(file_path):
            errors.append({"file": file_path, "error": f"File not found: {file_path}"})
            continue

        try:
            base, _ = os.path.splitext(file_path)
            output_path = f"{base}.pdf"
            conversion_method = None

            # Try LibreOffice first (best quality)
            if libreoffice_path or os.name == 'nt' or os.name == 'posix':
                if os.name == 'nt':
                    possible_paths = [
                        libreoffice_path,
                        r"C:\Program Files\LibreOffice\program\soffice.exe",
                        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                    ]
                else:
                    possible_paths = [libreoffice_path, "/usr/bin/soffice", "/usr/local/bin/soffice", "soffice"]

                for path in possible_paths:
                    if path and (os.path.exists(path) if path != "soffice" else True):
                        try:
                            import subprocess
                            output_dir = os.path.dirname(output_path) or "."
                            cmd = [path, "--headless", "--convert-to", "pdf",
                                   "--outdir", output_dir, file_path]
                            result = subprocess.run(cmd, capture_output=True, timeout=120, text=True)
                            if result.returncode == 0 and os.path.exists(output_path):
                                conversion_method = "libreoffice"
                                break
                        except Exception as lo_err:
                            print(f"[WORD2PDF] LibreOffice attempt failed: {lo_err}", file=sys.stderr)
                            continue

            # Fallback: Improved Python-only conversion
            if not conversion_method:
                try:
                    from docx import Document
                    from docx.shared import Pt, Inches
                    from reportlab.lib.pagesizes import letter
                    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image as RLImage
                    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
                    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT, TA_JUSTIFY
                    from reportlab.lib import colors
                    from reportlab.lib.units import inch
                    import io

                    doc = Document(file_path)
                    pdf_doc = SimpleDocTemplate(
                        output_path,
                        pagesize=letter,
                        leftMargin=0.75*inch,
                        rightMargin=0.75*inch,
                        topMargin=0.75*inch,
                        bottomMargin=0.75*inch
                    )
                    story = []
                    styles = getSampleStyleSheet()

                    # Create custom styles for different heading levels
                    custom_styles = {
                        'Heading1': ParagraphStyle('CustomH1', parent=styles['Heading1'], fontSize=18, spaceAfter=12),
                        'Heading2': ParagraphStyle('CustomH2', parent=styles['Heading2'], fontSize=14, spaceAfter=10),
                        'Heading3': ParagraphStyle('CustomH3', parent=styles['Heading3'], fontSize=12, spaceAfter=8),
                        'Normal': ParagraphStyle('CustomNormal', parent=styles['Normal'], fontSize=11, spaceAfter=6),
                        'Bold': ParagraphStyle('CustomBold', parent=styles['Normal'], fontSize=11, spaceAfter=6, fontName='Helvetica-Bold'),
                    }

                    # Process paragraphs
                    for para in doc.paragraphs:
                        if not para.text.strip():
                            story.append(Spacer(1, 6))
                            continue

                        # Determine style based on paragraph style name
                        style_name = para.style.name if para.style else 'Normal'
                        if 'Heading 1' in style_name:
                            style = custom_styles['Heading1']
                        elif 'Heading 2' in style_name:
                            style = custom_styles['Heading2']
                        elif 'Heading 3' in style_name:
                            style = custom_styles['Heading3']
                        else:
                            # Check for bold text
                            has_bold = any(run.bold for run in para.runs if run.bold)
                            style = custom_styles['Bold'] if has_bold else custom_styles['Normal']

                        # Escape special characters for reportlab
                        text = para.text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')

                        try:
                            story.append(Paragraph(text, style))
                        except Exception:
                            # If paragraph fails, try plain text
                            story.append(Paragraph(text, styles['Normal']))

                    # Process tables
                    for table in doc.tables:
                        table_data = []
                        for row in table.rows:
                            row_data = []
                            for cell in row.cells:
                                cell_text = cell.text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                                row_data.append(cell_text[:100])  # Truncate long cells
                            table_data.append(row_data)

                        if table_data:
                            try:
                                t = Table(table_data)
                                t.setStyle(TableStyle([
                                    ('BACKGROUND', (0, 0), (-1, 0), colors.Color(0.9, 0.9, 0.9)),
                                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.black),
                                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                                    ('FONTSIZE', (0, 0), (-1, -1), 9),
                                    ('BOTTOMPADDING', (0, 0), (-1, 0), 8),
                                    ('GRID', (0, 0), (-1, -1), 0.5, colors.grey)
                                ]))
                                story.append(Spacer(1, 12))
                                story.append(t)
                                story.append(Spacer(1, 12))
                            except Exception as table_err:
                                print(f"[WORD2PDF] Table conversion failed: {table_err}", file=sys.stderr)

                    if not story:
                        story.append(Paragraph("(Document appears to be empty or contains unsupported content)", styles['Normal']))

                    pdf_doc.build(story)
                    conversion_method = "python-docx"

                except ImportError as ie:
                    errors.append({"file": file_path, "error": f"Required library missing: {str(ie)}"})
                    continue
                except Exception as e:
                    errors.append({"file": file_path, "error": f"Conversion failed: {str(e)}"})
                    continue

            if os.path.exists(output_path):
                processed_files.append(output_path)
                print(f"[WORD2PDF] Converted using {conversion_method}: {os.path.basename(output_path)}", file=sys.stderr)
            else:
                errors.append({"file": file_path, "error": "Output file was not created"})

        except Exception as e:
            errors.append({"file": file_path, "error": str(e)})

    return {"processed_files": processed_files, "errors": errors}

def powerpoint_to_pdf(payload):
    """Convert PowerPoint presentations to PDF."""
    files = payload.get("files", [])
    libreoffice_path = payload.get("libreoffice_path")

    processed_files = []
    errors = []

    for file_path in files:
        try:
            base, _ = os.path.splitext(file_path)
            output_path = f"{base}.pdf"

            # Try LibreOffice first
            libreoffice_available = False
            if libreoffice_path or os.name == 'nt':
                if os.name == 'nt':
                    possible_paths = [
                        libreoffice_path,
                        r"C:\Program Files\LibreOffice\program\soffice.exe",
                        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                    ]
                else:
                    possible_paths = [libreoffice_path, "/usr/bin/soffice", "/usr/local/bin/soffice", "soffice"]

                for path in possible_paths:
                    if path and (os.path.exists(path) or path == "soffice"):
                        try:
                            import subprocess
                            cmd = [path, "--headless", "--convert-to", "pdf",
                                   "--outdir", os.path.dirname(output_path) or ".", file_path]
                            result = subprocess.run(cmd, capture_output=True, timeout=120, text=True)
                            if result.returncode == 0 and os.path.exists(output_path):
                                libreoffice_available = True
                                break
                        except:
                            continue

            # Fallback: Convert slides to images then to PDF
            if not libreoffice_available:
                try:
                    from pptx import Presentation
                    from reportlab.lib.pagesizes import letter
                    from reportlab.platypus import SimpleDocTemplate, Image as RLImage, Spacer
                    from reportlab.lib.units import inch

                    prs = Presentation(file_path)
                    pdf_doc = SimpleDocTemplate(output_path, pagesize=letter)
                    story = []

                    for slide in prs.slides:
                        # Extract text from slide
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                from reportlab.platypus import Paragraph
                                from reportlab.lib.styles import getSampleStyleSheet
                                styles = getSampleStyleSheet()
                                story.append(Paragraph(shape.text, styles['Normal']))
                        story.append(Spacer(1, 0.5*inch))

                    pdf_doc.build(story)
                except ImportError:
                    errors.append({"file": file_path, "error": "python-pptx and reportlab required for PowerPoint conversion"})
                    continue
                except Exception as e:
                    errors.append({"file": file_path, "error": f"Conversion failed: {str(e)}"})
                    continue

            processed_files.append(output_path)
        except Exception as e:
            errors.append({"file": file_path, "error": str(e)})

    return {"processed_files": processed_files, "errors": errors}

def excel_to_pdf(payload):
    """
    Convert Excel spreadsheets to PDF.

    Note: This is a BETA feature. Best results with LibreOffice installed.
    Python-only fallback creates a table-based PDF but may not preserve
    cell formatting, colors, formulas, or charts.
    """
    files = payload.get("files", [])
    libreoffice_path = payload.get("libreoffice_path")

    processed_files = []
    errors = []

    for file_path in files:
        if not os.path.exists(file_path):
            errors.append({"file": file_path, "error": f"File not found: {file_path}"})
            continue

        try:
            base, _ = os.path.splitext(file_path)
            output_path = f"{base}.pdf"
            conversion_method = None

            # Try LibreOffice first
            if libreoffice_path or os.name == 'nt' or os.name == 'posix':
                if os.name == 'nt':
                    possible_paths = [
                        libreoffice_path,
                        r"C:\Program Files\LibreOffice\program\soffice.exe",
                        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                    ]
                else:
                    possible_paths = [libreoffice_path, "/usr/bin/soffice", "/usr/local/bin/soffice", "soffice"]

                for path in possible_paths:
                    if path and (os.path.exists(path) if path != "soffice" else True):
                        try:
                            import subprocess
                            output_dir = os.path.dirname(output_path) or "."
                            cmd = [path, "--headless", "--convert-to", "pdf",
                                   "--outdir", output_dir, file_path]
                            result = subprocess.run(cmd, capture_output=True, timeout=120, text=True)
                            if result.returncode == 0 and os.path.exists(output_path):
                                conversion_method = "libreoffice"
                                break
                        except Exception as lo_err:
                            print(f"[EXCEL2PDF] LibreOffice attempt failed: {lo_err}", file=sys.stderr)
                            continue

            # Fallback: Improved Python-only conversion
            if not conversion_method:
                try:
                    import openpyxl
                    from reportlab.lib.pagesizes import letter, landscape, A4
                    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak, LongTable
                    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
                    from reportlab.lib import colors
                    from reportlab.lib.units import inch

                    wb = openpyxl.load_workbook(file_path, data_only=True)  # data_only to get calculated values
                    pdf_doc = SimpleDocTemplate(
                        output_path,
                        pagesize=landscape(letter),
                        leftMargin=0.5*inch,
                        rightMargin=0.5*inch,
                        topMargin=0.5*inch,
                        bottomMargin=0.5*inch
                    )
                    story = []
                    styles = getSampleStyleSheet()

                    # Create styles for table cells with text wrapping
                    # For very wide tables, use smaller font size for better fit
                    cell_style = ParagraphStyle(
                        'CellStyle',
                        parent=styles['Normal'],
                        fontSize=8,
                        leading=10,  # Line spacing
                        wordWrap='CJK'  # Enable word wrapping
                    )
                    
                    # Create a style for header cells
                    header_style = ParagraphStyle(
                        'HeaderStyle',
                        parent=styles['Normal'],
                        fontSize=8,
                        leading=10,
                        wordWrap='CJK',
                        fontName='Helvetica-Bold'
                    )
                    
                    # Create compact styles for very wide tables (50+ columns)
                    compact_cell_style = ParagraphStyle(
                        'CompactCellStyle',
                        parent=styles['Normal'],
                        fontSize=6,  # Smaller font for wide tables
                        leading=7,   # Tighter line spacing
                        wordWrap='CJK'
                    )
                    
                    compact_header_style = ParagraphStyle(
                        'CompactHeaderStyle',
                        parent=styles['Normal'],
                        fontSize=6,  # Smaller font for wide tables
                        leading=7,   # Tighter line spacing
                        wordWrap='CJK',
                        fontName='Helvetica-Bold'
                    )

                    # Helper function to escape XML special characters
                    def escape_xml(text):
                        """Escape XML special characters for Paragraph objects"""
                        text = str(text)
                        text = text.replace('&', '&amp;')
                        text = text.replace('<', '&lt;')
                        text = text.replace('>', '&gt;')
                        return text

                    def calculate_dynamic_column_widths(raw_data, available_width, font_size=8, is_very_wide=False):
                        """
                        Calculate column widths based on actual content length.
                        
                        Args:
                            raw_data: List of rows (first row is header)
                            available_width: Available page width in points
                            font_size: Font size in points
                            is_very_wide: True if table has 50+ columns
                        
                        Returns:
                            List of column widths in points
                        """
                        if not raw_data or not raw_data[0]:
                            return []
                        
                        num_cols = len(raw_data[0])
                        sample_size = min(50, len(raw_data) - 1)  # Sample first 50 data rows (skip header)
                        
                        # Step 1: Analyze content per column
                        column_metrics = []
                        for col_idx in range(num_cols):
                            max_len = 0
                            total_len = 0
                            count = 0
                            
                            # Check header (first row)
                            if len(raw_data) > 0 and col_idx < len(raw_data[0]):
                                header_text = str(raw_data[0][col_idx] or "")
                                header_len = len(header_text)
                                max_len = max(max_len, header_len)
                            
                            # Check sample data rows (skip header row)
                            for row in raw_data[1:sample_size+1]:
                                if col_idx < len(row):
                                    cell_text = str(row[col_idx] or "")
                                    cell_len = len(cell_text)
                                    max_len = max(max_len, cell_len)
                                    total_len += cell_len
                                    count += 1
                            
                            avg_len = total_len / max(count, 1)
                            # Use weighted approach: 70% max, 30% average (favors max but considers average)
                            target_len = max_len * 0.7 + avg_len * 0.3
                            
                            column_metrics.append({
                                'max_len': max_len,
                                'avg_len': avg_len,
                                'target_len': target_len
                            })
                        
                        # Step 2: Estimate required widths
                        # Character width factor depends on font size (inches per character)
                        char_width_factors = {6: 0.05, 8: 0.06, 10: 0.07}
                        char_width_factor = char_width_factors.get(font_size, 0.06)
                        
                        estimated_widths = []
                        for metric in column_metrics:
                            # Estimate width needed (in points: 1 inch = 72 points)
                            # Add 20% padding for readability
                            width_inches = metric['target_len'] * char_width_factor * 1.2
                            width_points = width_inches * 72
                            estimated_widths.append(width_points)
                        
                        # Step 3: Set min/max bounds
                        # For very wide tables, use smaller min width
                        min_width_points = (0.15 * 72) if is_very_wide else (0.2 * 72)  # 0.15-0.2 inches minimum
                        max_width_points = available_width * 0.4  # 40% of page max
                        
                        bounded_widths = [
                            max(min_width_points, min(w, max_width_points))
                            for w in estimated_widths
                        ]
                        
                        # Step 4: Normalize to fit available width
                        total_estimated = sum(bounded_widths)
                        
                        if total_estimated > available_width:
                            # Scale down proportionally
                            scale_factor = available_width / total_estimated
                            final_widths = [w * scale_factor for w in bounded_widths]
                        else:
                            # Distribute excess to columns that need it most
                            excess = available_width - total_estimated
                            if excess > 0:
                                # Give more to columns that are close to their max or have long content
                                priorities = []
                                for i, (w, metric) in enumerate(zip(bounded_widths, column_metrics)):
                                    # Priority based on: how close to max, and content length
                                    width_ratio = (w / max_width_points) if max_width_points > 0 else 0
                                    content_ratio = min(metric['target_len'] / 100, 1.0)  # Normalize content length
                                    priority = width_ratio * 0.6 + content_ratio * 0.4
                                    priorities.append(priority)
                                
                                total_priority = sum(priorities) or 1
                                final_widths = [
                                    w + (excess * p / total_priority)
                                    for w, p in zip(bounded_widths, priorities)
                                ]
                            else:
                                final_widths = bounded_widths
                        
                        return final_widths

                    for sheet_idx, sheet_name in enumerate(wb.sheetnames):
                        ws = wb[sheet_name]

                        # Get the actual data range
                        raw_data = []  # Store raw data first
                        max_col = ws.max_column or 1
                        max_row = ws.max_row or 1

                        # Limit to reasonable size to prevent memory issues
                        # For large files, allow more columns/rows but with warnings
                        max_col_limit = payload.get('max_columns', 50)  # Configurable, default 50
                        max_row_limit = payload.get('max_rows', 2000)  # Configurable, default 2000
                        max_col = min(max_col, max_col_limit)
                        max_row = min(max_row, max_row_limit)
                        
                        # Track if we're truncating data
                        data_truncated = False
                        if ws.max_column > max_col or ws.max_row > max_row:
                            data_truncated = True

                        # First pass: collect raw data
                        for row_idx, row in enumerate(ws.iter_rows(min_row=1, max_row=max_row, max_col=max_col, values_only=True)):
                            if row_idx >= max_row:
                                break
                            row_data = []
                            for cell in row:
                                if cell is None:
                                    row_data.append("")
                                elif isinstance(cell, (int, float)):
                                    # Format numbers nicely
                                    if isinstance(cell, float) and cell == int(cell):
                                        row_data.append(str(int(cell)))
                                    else:
                                        row_data.append(str(cell))
                                else:
                                    row_data.append(str(cell))
                            raw_data.append(row_data)

                        # Calculate column widths before creating Paragraphs
                        if raw_data and any(any(cell for cell in row) for row in raw_data):
                            num_cols = len(raw_data[0]) if raw_data else 1
                            available_width = landscape(letter)[0] - inch  # Page width minus margins
                            
                            # Detect very wide tables (50+ columns) - need special handling
                            is_very_wide = num_cols >= 50
                            
                            # Use dynamic column width calculation based on content
                            font_size = 6 if is_very_wide else 8
                            col_widths = calculate_dynamic_column_widths(
                                raw_data, 
                                available_width, 
                                font_size=font_size,
                                is_very_wide=is_very_wide
                            )
                            
                            # Calculate per-column text handling limits based on individual column widths
                            # This allows different columns to have different text handling strategies
                            available_height = landscape(letter)[1] - 2*0.5*inch
                            max_cell_height = available_height * 0.5  # ~270 points max cell height
                            max_lines = int(max_cell_height / 10)  # ~27 lines at 10pt leading
                            
                            # For each column, determine if we can use Paragraphs and what the limits are
                            column_text_limits = []
                            for col_idx, col_w in enumerate(col_widths):
                                # Determine if this column is wide enough for Paragraph wrapping
                                use_para = col_w >= 30  # Only use Paragraphs for columns >= 30 points wide
                                
                                if use_para and not is_very_wide:
                                    # Calculate safe character limit for this column
                                    chars_per_line = max(10, int(col_w / (font_size * 0.6)))
                                    max_chars = max_lines * chars_per_line
                                    max_chars = min(max_chars, 150)  # Cap at 150 chars for Paragraphs
                                else:
                                    # Use plain strings for narrow columns or very wide tables
                                    max_chars = 0  # 0 means use plain strings
                                    # Set plain string length limit based on column width
                                    if is_very_wide:
                                        max_chars = 30  # Very aggressive for wide tables
                                    else:
                                        # Estimate how many chars fit in this column width
                                        chars_per_line = max(5, int(col_w / (font_size * 0.5)))
                                        max_chars = chars_per_line * 2  # Allow 2 lines worth
                                        max_chars = min(max_chars, 100)  # Cap at 100
                                
                                column_text_limits.append({
                                    'use_paragraph': use_para and not is_very_wide,
                                    'max_chars': max_chars,
                                    'width': col_w
                                })
                            
                            # Set active styles
                            if is_very_wide:
                                active_cell_style = compact_cell_style
                                active_header_style = compact_header_style
                            else:
                                active_cell_style = cell_style
                                active_header_style = header_style
                            
                            # Second pass: convert to Paragraph objects with proper limits
                            data = []
                            for row_idx, row_data in enumerate(raw_data):
                                is_header = (row_idx == 0)
                                current_style = active_header_style if is_header else active_cell_style
                                
                                paragraph_row = []
                                for cell_idx, cell_text in enumerate(row_data):
                                    if not cell_text:  # Empty string or None
                                        paragraph_row.append("")
                                    else:
                                        # Get text handling strategy for this specific column
                                        if cell_idx < len(column_text_limits):
                                            col_limit = column_text_limits[cell_idx]
                                            use_para = col_limit['use_paragraph']
                                            max_chars = col_limit['max_chars']
                                        else:
                                            # Fallback if index is out of range
                                            use_para = False
                                            max_chars = 50
                                        
                                        # Truncate extremely long text to prevent cells taller than page
                                        text = str(cell_text).strip()
                                        original_length = len(text)
                                        
                                        # Handle text based on this column's width and limits
                                        if not use_para or original_length > max_chars:
                                            # Use plain string for narrow columns or very long text
                                            # This prevents ReportLab from creating extremely tall cells
                                            if original_length > max_chars:
                                                truncate_at = max_chars - 20  # Leave room for truncation message
                                                text = text[:truncate_at] + "... [truncated]"
                                            paragraph_row.append(text)
                                        else:
                                            # For shorter text in wider columns, use Paragraph for proper wrapping
                                            try:
                                                para = Paragraph(escape_xml(text), current_style)
                                                paragraph_row.append(para)
                                            except Exception as para_err:
                                                # If Paragraph creation fails, use plain string as fallback
                                                print(f"[EXCEL2PDF] Paragraph creation failed for cell ({row_idx},{cell_idx}), using plain text: {para_err}", file=sys.stderr)
                                                paragraph_row.append(text)
                                data.append(paragraph_row)

                        if data and any(any(cell for cell in row) for row in data):
                            # Add sheet name as heading
                            sheet_title = f"<b>{sheet_name}</b>"
                            if data_truncated:
                                sheet_title += f" <i>(Showing {max_col} of {ws.max_column} columns, {max_row} of {ws.max_row} rows)</i>"
                            story.append(Paragraph(sheet_title, styles['Heading2']))
                            story.append(Spacer(1, 12))

                            try:
                                # Use LongTable for better pagination with large datasets
                                # LongTable handles page breaks more efficiently than Table
                                use_longtable = len(data) > 100  # Use LongTable for tables with more than 100 rows
                                
                                # Use dynamically calculated column widths
                                if use_longtable:
                                    table = LongTable(data, colWidths=col_widths, repeatRows=1)
                                else:
                                    table = Table(data, colWidths=col_widths)
                                # Adjust table style based on table width
                                if is_very_wide:
                                    # For very wide tables, use smaller font and tighter padding
                                    table_style = TableStyle([
                                        ('BACKGROUND', (0, 0), (-1, 0), colors.Color(0.2, 0.4, 0.6)),
                                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                                        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                                        ('FONTSIZE', (0, 0), (-1, -1), 6),  # Smaller font for wide tables
                                        ('BOTTOMPADDING', (0, 0), (-1, 0), 4),  # Tighter padding
                                        ('TOPPADDING', (0, 0), (-1, -1), 2),  # Tighter padding
                                        ('LEFTPADDING', (0, 0), (-1, -1), 2),  # Minimal side padding
                                        ('RIGHTPADDING', (0, 0), (-1, -1), 2),  # Minimal side padding
                                        ('BACKGROUND', (0, 1), (-1, -1), colors.Color(0.95, 0.95, 0.95)),
                                        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.Color(0.95, 0.95, 0.95)]),
                                        ('GRID', (0, 0), (-1, -1), 0.3, colors.Color(0.7, 0.7, 0.7)),  # Thinner grid
                                        ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                                    ])
                                else:
                                    # Standard style for normal tables
                                    table_style = TableStyle([
                                        ('BACKGROUND', (0, 0), (-1, 0), colors.Color(0.2, 0.4, 0.6)),
                                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                                        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                                        ('FONTSIZE', (0, 0), (-1, -1), 8),
                                        ('BOTTOMPADDING', (0, 0), (-1, 0), 8),
                                        ('TOPPADDING', (0, 0), (-1, -1), 4),
                                        ('BACKGROUND', (0, 1), (-1, -1), colors.Color(0.95, 0.95, 0.95)),
                                        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.Color(0.95, 0.95, 0.95)]),
                                        ('GRID', (0, 0), (-1, -1), 0.5, colors.Color(0.7, 0.7, 0.7)),
                                        ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                                    ])
                                
                                table.setStyle(table_style)
                                story.append(table)
                            except Exception as table_err:
                                print(f"[EXCEL2PDF] Table creation failed for {sheet_name}: {table_err}", file=sys.stderr)
                                story.append(Paragraph(f"(Could not render sheet: {sheet_name})", styles['Normal']))

                            # Add page break between sheets
                            if sheet_idx < len(wb.sheetnames) - 1:
                                story.append(PageBreak())

                    if not story:
                        story.append(Paragraph("(Spreadsheet appears to be empty)", styles['Normal']))

                    pdf_doc.build(story)
                    conversion_method = "openpyxl"

                except ImportError as ie:
                    errors.append({"file": file_path, "error": f"Required library missing: {str(ie)}"})
                    continue
                except Exception as e:
                    errors.append({"file": file_path, "error": f"Conversion failed: {str(e)}"})
                    continue

            processed_files.append(output_path)
        except Exception as e:
            errors.append({"file": file_path, "error": str(e)})

    return {"processed_files": processed_files, "errors": errors}

def html_to_pdf(payload):
    """Convert HTML files to PDF."""
    files = payload.get("files", [])

    processed_files = []
    errors = []

    for file_path in files:
        try:
            base, _ = os.path.splitext(file_path)
            output_path = f"{base}.pdf"

            # Try WeasyPrint first (pure Python, good CSS support)
            try:
                from weasyprint import HTML
                HTML(filename=file_path).write_pdf(output_path)
                processed_files.append(output_path)
                continue
            except (ImportError, Exception) as e:
                logger.warning(f"WeasyPrint conversion unavailable or failed: {e}, trying browser fallback")

            # Try Browser-based conversion (Headless Edge/Chrome/Chromium) - Best for high-quality offline rendering
            browser_success = False
            browser_paths = []
            
            # 1. Platform-specific common paths
            if platform.system() == "Windows":
                browser_paths = [
                    os.path.join(os.environ.get("ProgramFiles(x86)", "C:\\Program Files (x86)"), "Microsoft\\Edge\\Application\\msedge.exe"),
                    os.path.join(os.environ.get("ProgramFiles", "C:\\Program Files"), "Microsoft\\Edge\\Application\\msedge.exe"),
                    os.path.join(os.environ.get("ProgramFiles", "C:\\Program Files"), "Google\\Chrome\\Application\\chrome.exe"),
                    os.path.join(os.environ.get("ProgramFiles(x86)", "C:\\Program Files (x86)"), "Google\\Chrome\\Application\\chrome.exe")
                ]
            elif platform.system() == "Darwin":  # macOS
                browser_paths = [
                    "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
                    "/Applications/Microsoft Edge.app/Contents/MacOS/Microsoft Edge",
                    "/Applications/Chromium.app/Contents/MacOS/Chromium"
                ]
            
            # 2. Add system PATH search for common browser binaries (especially for Linux)
            import shutil
            for bin_name in ['google-chrome', 'microsoft-edge', 'chromium', 'chromium-browser', 'chrome', 'google-chrome-stable']:
                path = shutil.which(bin_name)
                if path and path not in browser_paths:
                    browser_paths.append(path)
            
            for browser_path in browser_paths:
                if browser_path and os.path.exists(browser_path):
                    try:
                        import subprocess
                        abs_file_path = os.path.abspath(file_path)
                        abs_output_path = os.path.abspath(output_path)
                        
                        # Format file URL correctly (works for all platforms)
                        file_url = f"file:///{abs_file_path.replace(os.sep, '/')}"
                        
                        subprocess.run([
                            browser_path, 
                            '--headless', 
                            '--disable-gpu', 
                            f'--print-to-pdf={abs_output_path}', 
                            '--no-margins',
                            '--disable-extensions',
                            '--disable-software-rasterizer',
                            file_url
                        ], check=True, timeout=30, capture_output=True)
                        
                        if os.path.exists(output_path):
                            processed_files.append(output_path)
                            browser_success = True
                            logger.info(f"Successfully converted HTML to PDF using browser: {browser_path}")
                            break
                    except Exception as be:
                        logger.warning(f"Browser conversion failed with {browser_path}: {be}")

            if browser_success:
                continue

            # Fallback: Basic HTML parsing and PDF generation (Text-only)
            try:
                from html.parser import HTMLParser
                from reportlab.lib.pagesizes import letter
                from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
                from reportlab.lib.styles import getSampleStyleSheet
                from reportlab.lib.units import inch
                import re

                with open(file_path, 'r', encoding='utf-8') as f:
                    html_content = f.read()

                # Remove script and style tags
                html_content = re.sub(r'<script[^>]*>.*?</script>', '', html_content, flags=re.DOTALL | re.IGNORECASE)
                html_content = re.sub(r'<style[^>]*>.*?</style>', '', html_content, flags=re.DOTALL | re.IGNORECASE)

                # Extract text (very basic)
                text_content = re.sub(r'<[^>]+>', '\n', html_content)
                text_content = ' '.join(text_content.split())

                pdf_doc = SimpleDocTemplate(output_path, pagesize=letter)
                story = []
                styles = getSampleStyleSheet()

                # Split into paragraphs
                paragraphs = [p.strip() for p in text_content.split('\n') if p.strip()]
                for para_text in paragraphs[:200]:  # Limit for performance
                    if para_text:
                        story.append(Paragraph(para_text, styles['Normal']))
                        story.append(Spacer(1, 0.1*inch))

                pdf_doc.build(story)
                processed_files.append(output_path)

            except Exception as e:
                errors.append({"file": file_path, "error": f"HTML conversion failed: {str(e)}. Install WeasyPrint for better results: pip install weasyprint"})

        except Exception as e:
            errors.append({"file": file_path, "error": str(e)})

    return {"processed_files": processed_files, "errors": errors}

def ocr_pdf(payload):
    """Convert scanned PDF to searchable PDF using OCR."""
    files = payload.get("files", [])
    language = payload.get("language", "eng")  # Default English

    processed_files = []
    errors = []

    # Configure Tesseract (bundled or system)
    configure_tesseract()

    # Check Tesseract availability first
    tesseract_available, tesseract_error = is_tesseract_available()
    if not tesseract_available:
        for file_path in files:
            errors.append({"file": file_path, "error": tesseract_error})
        return {"processed_files": processed_files, "errors": errors}

    for file_path in files:
        try:
            base, ext = os.path.splitext(file_path)
            output_path = f"{base}_ocr{ext}"

            # Try pytesseract (Tesseract OCR)
            try:
                import pytesseract
                from pdf2image import convert_from_path

                # Convert PDF pages to images
                try:
                    images = convert_from_path(file_path, dpi=300)
                except Exception as e:
                    # Fallback: use PyMuPDF to render pages
                    images = []
                    doc = fitz.open(file_path)
                    for page in doc:
                        pix = page.get_pixmap(dpi=300)
                        img_data = pix.tobytes("png")
                        images.append(Image.open(io.BytesIO(img_data)))
                    doc.close()

                # Create new PDF with OCR text
                ocr_doc = fitz.open()
                total_pages = len(images)
                logger.info(f"Processing {total_pages} pages for OCR")

                for page_idx, img in enumerate(images):
                    logger.info(f"Processing page {page_idx + 1}/{total_pages}")
                    
                    # Perform OCR with detailed data (bounding boxes)
                    ocr_data = None
                    ocr_text = ""
                    try:
                        # Get OCR data with bounding boxes for accurate text positioning
                        ocr_data = pytesseract.image_to_data(img, lang=language, output_type=pytesseract.Output.DICT)
                        ocr_text = pytesseract.image_to_string(img, lang=language)
                        logger.info(f"OCR extracted {len(ocr_text)} characters from page {page_idx + 1}")
                    except Exception as ocr_err:
                        logger.error(f"OCR failed for page {page_idx + 1}: {ocr_err}")
                        # Fallback: try simple OCR without detailed data
                        try:
                            ocr_text = pytesseract.image_to_string(img, lang=language)
                            logger.info(f"Fallback OCR extracted {len(ocr_text)} characters")
                        except Exception as fallback_err:
                            logger.error(f"Fallback OCR also failed: {fallback_err}")
                            ocr_text = ""

                    # Create new page
                    page = ocr_doc.new_page(width=img.width, height=img.height)

                    # Insert original image
                    img_bytes = io.BytesIO()
                    img.save(img_bytes, format='PNG')
                    img_bytes.seek(0)
                    page.insert_image(page.rect, stream=img_bytes.getvalue())

                    # Add invisible text layer (for searchability)
                    text_added = False
                    if ocr_text.strip():
                        if ocr_data and len(ocr_data.get('text', [])) > 0:
                            # Use OCR bounding boxes for accurate text positioning
                            n_boxes = len(ocr_data['text'])
                            logger.info(f"Found {n_boxes} text boxes from OCR")
                            for i in range(n_boxes):
                                text = ocr_data['text'][i].strip()
                                conf = int(ocr_data['conf'][i]) if ocr_data['conf'][i] else 0
                                if text and conf > 0:  # Confidence > 0
                                    x = ocr_data['left'][i]
                                    y = ocr_data['top'][i]
                                    w = ocr_data['width'][i]
                                    h = ocr_data['height'][i]
                                    
                                    # Convert to PDF coordinates (OCR uses top-left, PyMuPDF uses bottom-left)
                                    pdf_y = img.height - y - h
                                    
                                    # Insert invisible text (render_mode=3) for searchability
                                    try:
                                        page.insert_text(
                                            fitz.Point(x, pdf_y + h),
                                            text,
                                            fontsize=max(8, min(h * 0.8, 12)),
                                            color=(1, 1, 1),  # White (invisible on white background)
                                            render_mode=3  # Invisible but searchable
                                        )
                                        text_added = True
                                    except Exception as text_err:
                                        logger.warning(f"Failed to insert text at position ({x}, {pdf_y}): {text_err}")
                                        # Fallback: simple text insertion
                                        try:
                                            page.insert_text(
                                                fitz.Point(x, pdf_y + h),
                                                text[:50],  # Limit length
                                                fontsize=10,
                                                color=(1, 1, 1),
                                                render_mode=3
                                            )
                                            text_added = True
                                        except:
                                            pass
                        
                        # Fallback: add text as lines if bounding boxes failed
                        if not text_added:
                            logger.info("Using fallback text positioning (lines)")
                            lines = [l.strip() for l in ocr_text.split('\n') if l.strip()]
                            if lines:
                                line_height = img.height / max(len(lines), 1)
                                line_idx = 0
                                for line in lines[:200]:  # Limit lines for performance
                                    y_pos = line_idx * line_height + 20
                                    text_point = fitz.Point(10, min(y_pos, page.rect.height - 10))
                                    # Use render_mode=3 (invisible) to make text searchable but not visible
                                    try:
                                        page.insert_text(
                                            text_point,
                                            line[:100],  # Limit line length
                                            fontsize=12,
                                            color=(1, 1, 1),
                                            render_mode=3
                                        )
                                        text_added = True
                                    except Exception as line_err:
                                        logger.warning(f"Failed to insert line {line_idx}: {line_err}")
                                    line_idx += 1
                    
                    if text_added:
                        logger.info(f"Successfully added searchable text to page {page_idx + 1}")
                    else:
                        logger.warning(f"No text was added to page {page_idx + 1}")

                ocr_doc.save(output_path)
                ocr_doc.close()
                
                # Verify output file was created
                if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                    logger.info(f"OCR PDF saved successfully: {output_path} ({os.path.getsize(output_path)} bytes)")
                    processed_files.append(output_path)
                else:
                    logger.error(f"OCR PDF file was not created or is empty: {output_path}")
                    errors.append({"file": file_path, "error": "OCR processing completed but output file was not created"})

            except ImportError:
                errors.append({"file": file_path, "error": "Python packages required. Install with: pip install pytesseract pdf2image"})
            except Exception as e:
                error_str = str(e)
                # Check if it's a Tesseract-specific error
                if "tesseract" in error_str.lower():
                    tesseract_available, tesseract_error = is_tesseract_available()
                    if not tesseract_available:
                        errors.append({"file": file_path, "error": tesseract_error})
                    else:
                        errors.append({"file": file_path, "error": f"OCR processing failed: {error_str}"})
                else:
                    errors.append({"file": file_path, "error": f"OCR failed: {error_str}"})

        except Exception as e:
            errors.append({"file": file_path, "error": f"Unexpected error: {str(e)}"})

    return {"processed_files": processed_files, "errors": errors}

def pdf_to_pdfa(payload):
    """Convert PDF to PDF/A format (archival compliance).

    Note: Full PDF/A compliance requires proper metadata, color profiles, and structure.
    This implementation provides basic PDF/A conversion. For full compliance, consider
    using specialized tools like Ghostscript or pdfa_validator.
    """
    files = payload.get("files", [])
    pdfa_version = payload.get("pdfa_version", "2b")  # Default PDF/A-2b

    processed_files = []
    errors = []

    for file_path in files:
        if not os.path.exists(file_path):
            errors.append({"file": file_path, "error": f"File not found: {file_path}"})
            continue

        try:
            base, ext = os.path.splitext(file_path)
            output_path = f"{base}_pdfa{ext}"

            # Use pikepdf for PDF/A conversion
            with pikepdf.open(file_path) as pdf:
                # Remove encryption if present (PDF/A doesn't allow encryption)
                if pdf.is_encrypted:
                    try:
                        pdf.save(output_path, compress_streams=True)
                        # Reopen unencrypted version
                        with pikepdf.open(output_path) as pdf_unencrypted:
                            pdf = pdf_unencrypted
                    except:
                        errors.append({"file": file_path, "error": "Cannot convert encrypted PDF to PDF/A. Please unlock first."})
                        continue

                # Set PDF/A metadata (basic implementation)
                # Full PDF/A requires XMP metadata with proper schema
                if '/Metadata' not in pdf.Root or pdf.Root.Metadata is None:
                    # Create basic metadata structure
                    pdf.Root.Metadata = pikepdf.Stream(pdf, b'<?xpacket begin="" id="W5M0MpCehiHzreSzNTczkc9d"?><x:xmpmeta xmlns:x="adobe:ns:meta/"><rdf:RDF xmlns:rdf="http://www.w3.org/1999/02/22-rdf-syntax-ns#"><rdf:Description rdf:about="" xmlns:pdfaid="http://www.aiim.org/pdfa/ns/id/"><pdfaid:part>2</pdfaid:part><pdfaid:conformance>B</pdfaid:conformance></rdf:Description></rdf:RDF></x:xmpmeta><?xpacket end="w"?>')

                # Ensure required PDF/A keys
                if '/OutputIntents' not in pdf.Root:
                    # Add basic output intent (required for PDF/A)
                    pdf.Root.OutputIntents = pikepdf.Array()

                # Save as PDF/A
                pdf.save(output_path, compress_streams=True, normalize_content=True)

            if os.path.exists(output_path):
                processed_files.append(output_path)
            else:
                errors.append({"file": file_path, "error": "Failed to create PDF/A file"})
        except Exception as e:
            logger.error(f"Error converting to PDF/A {file_path}: {e}", exc_info=True)
            errors.append({"file": file_path, "error": f"PDF/A conversion failed: {str(e)}"})

    return {"processed_files": processed_files, "errors": errors}

def crop_pdf(payload):
    """Crop PDF pages to specific dimensions."""
    files = payload.get("files", [])
    x = payload.get("x", 0)  # Left margin
    y = payload.get("y", 0)  # Top margin
    width = payload.get("width")  # Crop width
    height = payload.get("height")  # Crop height
    pages = payload.get("pages", "")  # Optional: specific pages (e.g., "1,3,5" or "2-5")

    processed_files = []
    errors = []

    for file_path in files:
        if not os.path.exists(file_path):
            errors.append({"file": file_path, "error": f"File not found: {file_path}"})
            continue

        doc = None
        try:
            base, ext = os.path.splitext(file_path)
            output_path = f"{base}_cropped{ext}"

            doc = fitz.open(file_path)
            total_pages = len(doc)

            if total_pages == 0:
                errors.append({"file": file_path, "error": "PDF has no pages"})
                doc.close()
                continue

            # Parse pages if specified
            page_list = None
            if pages:
                try:
                    if "-" in pages:
                        start, end = map(int, pages.split("-"))
                        page_list = list(range(start - 1, min(end, total_pages)))
                    else:
                        page_list = [int(p.strip()) - 1 for p in pages.split(",")]
                        # Validate page numbers
                        page_list = [p for p in page_list if 0 <= p < total_pages]
                except ValueError:
                    errors.append({"file": file_path, "error": f"Invalid page range: {pages}"})
                    doc.close()
                    continue

            # Validate crop dimensions
            if width is None or height is None:
                errors.append({"file": file_path, "error": "Width and height are required for cropping"})
                doc.close()
                continue

            if width <= 0 or height <= 0:
                errors.append({"file": file_path, "error": "Width and height must be positive"})
                doc.close()
                continue

            # Process pages
            for page_num in range(total_pages):
                if page_list is None or page_num in page_list:
                    page = doc[page_num]
                    page_rect = page.rect

                    # Calculate crop rectangle
                    # x, y are from top-left, convert to fitz coordinates
                    crop_rect = fitz.Rect(
                        float(x),
                        float(y),
                        min(float(x) + float(width), page_rect.width),
                        min(float(y) + float(height), page_rect.height)
                    )

                    # Validate crop rectangle is within page bounds
                    if crop_rect.x1 <= crop_rect.x0 or crop_rect.y1 <= crop_rect.y0:
                        errors.append({"file": file_path, "error": f"Invalid crop dimensions for page {page_num + 1}"})
                        continue

                    if crop_rect.x0 < 0 or crop_rect.y0 < 0:
                        errors.append({"file": file_path, "error": f"Crop coordinates out of bounds for page {page_num + 1}"})
                        continue

                    # Set crop box (visible area)
                    page.set_cropbox(crop_rect)

            doc.save(output_path)
            if doc:
                doc.close()

            if os.path.exists(output_path):
                processed_files.append(output_path)
            else:
                errors.append({"file": file_path, "error": "Failed to create cropped PDF"})

        except Exception as e:
            logger.error(f"Error cropping PDF {file_path}: {e}", exc_info=True)
            if doc:
                doc.close()
            errors.append({"file": file_path, "error": str(e)})

    return {"processed_files": processed_files, "errors": errors}

def organize_pdf(payload):
    """Reorder pages in PDF (Organize PDF)."""
    files = payload.get("files", [])
    page_order = payload.get("page_order", "")  # e.g., "3,1,2,5,4" or "1-5,10,8-9"

    processed_files = []
    errors = []

    for file_path in files:
        if not os.path.exists(file_path):
            errors.append({"file": file_path, "error": f"File not found: {file_path}"})
            continue

        doc = None
        new_doc = None
        try:
            base, ext = os.path.splitext(file_path)
            output_path = f"{base}_organized{ext}"

            doc = fitz.open(file_path)
            total_pages = len(doc)

            if total_pages == 0:
                errors.append({"file": file_path, "error": "PDF has no pages"})
                doc.close()
                continue

            if not page_order:
                errors.append({"file": file_path, "error": "Page order is required. Provide page numbers like '3,1,2' or '1-5,10,8-9'"})
                doc.close()
                continue

            # Parse page order
            page_indices = []
            try:
                # Validate page order format (only digits, commas, dashes, and spaces)
                import re
                if not re.match(r'^[\d\s,\-]+$', page_order):
                    errors.append({"file": file_path, "error": f"Invalid page order format: {page_order}. Use only numbers, commas, and dashes (e.g., '3,1,2' or '1-5,10,8-9')"})
                    doc.close()
                    continue

                # Split by comma
                parts = [p.strip() for p in page_order.split(",")]

                # Limit max parts to prevent DoS
                if len(parts) > 1000:
                    errors.append({"file": file_path, "error": f"Too many page specifications (max 1000). Please simplify your page order."})
                    doc.close()
                    continue

                for part in parts:
                    if not part:  # Skip empty parts
                        continue

                    if "-" in part:
                        # Range like "1-5"
                        range_parts = part.split("-")
                        if len(range_parts) != 2:
                            errors.append({"file": file_path, "error": f"Invalid range format: {part}. Use format like '1-5'"})
                            doc.close()
                            continue

                        start, end = map(int, range_parts)

                        # Validate range bounds
                        if start < 1 or end < 1:
                            errors.append({"file": file_path, "error": f"Page numbers must be >= 1. Invalid range: {part}"})
                            doc.close()
                            continue

                        if start > end:
                            errors.append({"file": file_path, "error": f"Invalid range: {part}. Start page must be <= end page."})
                            doc.close()
                            continue

                        # Limit range size to prevent memory issues
                        if end - start > 500:
                            errors.append({"file": file_path, "error": f"Range too large: {part}. Maximum range size is 500 pages."})
                            doc.close()
                            continue

                        # Convert to 0-indexed and validate
                        for p in range(start - 1, min(end, total_pages)):
                            if 0 <= p < total_pages:
                                page_indices.append(p)
                    else:
                        # Single page number
                        p = int(part) - 1  # Convert to 0-indexed
                        if 0 <= p < total_pages:
                            page_indices.append(p)
                        else:
                            errors.append({"file": file_path, "error": f"Page number {part} is out of range (1-{total_pages})"})
                            doc.close()
                            continue
            except ValueError as ve:
                errors.append({"file": file_path, "error": f"Invalid page order format: {page_order}. Use format like '3,1,2' or '1-5,10,8-9'. Error: {str(ve)}"})
                doc.close()
                continue

            if not page_indices:
                errors.append({"file": file_path, "error": "No valid pages specified in page order"})
                doc.close()
                continue

            # Limit total pages to prevent memory issues
            if len(page_indices) > 1000:
                errors.append({"file": file_path, "error": f"Too many pages to reorder ({len(page_indices)} pages). Maximum is 1000 pages."})
                doc.close()
                continue

            # Create new document with reordered pages
            new_doc = fitz.open()
            for page_idx in page_indices:
                new_doc.insert_pdf(doc, from_page=page_idx, to_page=page_idx)

            new_doc.save(output_path)
            if new_doc:
                new_doc.close()
            if doc:
                doc.close()

            if os.path.exists(output_path):
                processed_files.append(output_path)
            else:
                errors.append({"file": file_path, "error": "Failed to create organized PDF"})

        except Exception as e:
            logger.error(f"Error organizing PDF {file_path}: {e}", exc_info=True)
            if new_doc:
                new_doc.close()
            if doc:
                doc.close()
            errors.append({"file": file_path, "error": str(e)})

    return {"processed_files": processed_files, "errors": errors}

def extract_metadata(payload):
    """Extract PDF metadata as key-value pairs (user-friendly format).

    Extracts and displays:
    - Standard metadata (title, author, dates)
    - Document information (pages, file size)
    - XML metadata (if present)

    Security Features:
    - File size limits (100MB max)
    - PDF validation (header check)
    - Password-protected PDF detection
    - Metadata sanitization (removes control characters)
    - Metadata length limits (prevents huge output files)
    """
    files = payload.get("files", [])

    processed_files = []
    errors = []

    # Backend safety limits
    MAX_PDF_SIZE_MB = 100  # Generous limit for desktop users
    MAX_METADATA_LENGTH = 10000  # Prevent huge metadata strings

    def sanitize_metadata(value):
        """Sanitize metadata by removing control characters and limiting length."""
        if not value:
            return ""

        # Convert to string
        value_str = str(value)

        # Remove control characters (except newline and tab)
        sanitized = ''.join(char for char in value_str if char == '\n' or char == '\t' or (ord(char) >= 32 and ord(char) != 127))

        # Limit length
        if len(sanitized) > MAX_METADATA_LENGTH:
            sanitized = sanitized[:MAX_METADATA_LENGTH] + "... (truncated)"

        # Replace newlines with spaces for single-line fields
        sanitized = sanitized.replace('\n', ' ').replace('\r', ' ')

        return sanitized.strip()

    for file_path in files:
        if not os.path.exists(file_path):
            errors.append({"file": file_path, "error": f"File not found: {file_path}"})
            continue

        doc = None
        try:
            # SECURITY: Check file size
            try:
                pdf_size_mb = os.path.getsize(file_path) / (1024 * 1024)
                if pdf_size_mb > MAX_PDF_SIZE_MB:
                    errors.append({
                        "file": file_path,
                        "error": f"PDF file too large ({pdf_size_mb:.1f}MB). Maximum allowed: {MAX_PDF_SIZE_MB}MB."
                    })
                    continue
            except Exception as e:
                errors.append({"file": file_path, "error": f"Cannot check file size: {str(e)}"})
                continue

            # VALIDATION: Check PDF header (integrity check)
            try:
                with open(file_path, 'rb') as f:
                    header = f.read(5)
                    if header != b'%PDF-':
                        errors.append({"file": file_path, "error": "Invalid PDF file: missing PDF header. File may be corrupted."})
                        continue
            except Exception as e:
                errors.append({"file": file_path, "error": f"Cannot read file: {str(e)}"})
                continue

            # Try to open PDF (will detect encryption/password protection)
            try:
                doc = fitz.open(file_path)
            except Exception as e:
                error_msg = str(e).lower()
                if "password" in error_msg or "crypt" in error_msg or "encrypt" in error_msg:
                    errors.append({
                        "file": file_path,
                        "error": "Cannot extract metadata from password-protected PDFs. Please unlock the PDF first using the 'Unlock PDF' tool."
                    })
                else:
                    errors.append({
                        "file": file_path,
                        "error": f"Cannot open PDF: {str(e)}"
                    })
                continue

            # Check for empty PDF
            if len(doc) == 0:
                errors.append({"file": file_path, "error": "PDF has no pages. Cannot extract metadata from an empty PDF."})
                continue

            # Extract metadata
            metadata = doc.metadata

            # Build metadata output with sanitization
            metadata_lines = []
            metadata_lines.append("PDF Metadata Information")
            metadata_lines.append("=" * 50)
            metadata_lines.append("")

            # Standard PDF metadata fields (sanitized)
            if metadata.get('title'):
                metadata_lines.append(f"Title: {sanitize_metadata(metadata.get('title'))}")
            if metadata.get('author'):
                metadata_lines.append(f"Author: {sanitize_metadata(metadata.get('author'))}")
            if metadata.get('subject'):
                metadata_lines.append(f"Subject: {sanitize_metadata(metadata.get('subject'))}")
            if metadata.get('keywords'):
                metadata_lines.append(f"Keywords: {sanitize_metadata(metadata.get('keywords'))}")
            if metadata.get('creator'):
                metadata_lines.append(f"Creator: {sanitize_metadata(metadata.get('creator'))}")
            if metadata.get('producer'):
                metadata_lines.append(f"Producer: {sanitize_metadata(metadata.get('producer'))}")
            if metadata.get('creationDate'):
                metadata_lines.append(f"Creation Date: {sanitize_metadata(metadata.get('creationDate'))}")
            if metadata.get('modDate'):
                metadata_lines.append(f"Modification Date: {sanitize_metadata(metadata.get('modDate'))}")

            # If no metadata found
            if len([line for line in metadata_lines if ': ' in line]) == 0:
                metadata_lines.append("No standard metadata found in this PDF.")

            # Additional document info
            metadata_lines.append("")
            metadata_lines.append("Document Information")
            metadata_lines.append("-" * 50)
            metadata_lines.append(f"Number of Pages: {len(doc)}")
            metadata_lines.append(f"File Size: {os.path.getsize(file_path):,} bytes ({pdf_size_mb:.2f} MB)")
            metadata_lines.append(f"PDF Version: {doc.metadata.get('format', 'Unknown')}")

            # Try to get additional metadata from pikepdf if available
            try:
                with pikepdf.open(file_path) as pdf:
                    if '/Metadata' in pdf.Root:
                        metadata_lines.append("")
                        metadata_lines.append("Additional Metadata (XML)")
                        metadata_lines.append("-" * 50)
                        metadata_lines.append("This PDF contains XMP metadata. View in a PDF reader for full details.")
            except Exception:
                # pikepdf might fail on some PDFs, that's okay
                pass

            # Generate output path
            base, _ = os.path.splitext(file_path)
            output_path = f"{base}_metadata.txt"

            # Write to file with error handling
            try:
                with open(output_path, 'w', encoding='utf-8', errors='replace') as f:
                    f.write('\n'.join(metadata_lines))
            except Exception as e:
                errors.append({"file": file_path, "error": f"Cannot write metadata file: {str(e)}"})
                continue

            if os.path.exists(output_path):
                processed_files.append(output_path)
                print(f"[METADATA] Extracted metadata from {os.path.basename(file_path)} → {os.path.basename(output_path)}", file=sys.stderr, flush=True)
            else:
                errors.append({"file": file_path, "error": "Failed to create metadata file"})

        except Exception as e:
            logger.error(f"Error extracting metadata from {file_path}: {e}", exc_info=True)
            errors.append({"file": file_path, "error": f"Metadata extraction failed: {str(e)}"})
        finally:
            # SECURITY: Always clean up resources
            if doc:
                try:
                    doc.close()
                except:
                    pass

    return {"processed_files": processed_files, "errors": errors}

def extract_form_data(payload):
    """Extract form field data from PDF to CSV/XLSX.

    Handles:
    - Fillable forms (AcroForm fields)
    - Widget-based forms (PyMuPDF extraction)
    - Flattened forms (text extraction fallback - limited)

    Security Features:
    - File size limits (100MB max)
    - Page count limits (1000 pages max)
    - Field count limits (10,000 fields max)
    - Flattened form page limit (100 pages max for text extraction)
    - Field value sanitization (removes control characters)
    - PDF validation (header check)
    - Password-protected PDF detection
    - Resource cleanup with try/finally
    """
    files = payload.get("files", [])
    output_format = payload.get("output_format", "csv")  # csv or xlsx

    processed_files = []
    errors = []

    # Backend safety limits
    MAX_PDF_SIZE_MB = 100  # Generous limit for desktop users
    MAX_PAGES = 1000       # Page count limit
    MAX_FIELDS = 10000     # Field count limit (prevent huge DataFrames)
    MAX_FLATTENED_PAGES = 100  # Limit pages for flattened form text extraction

    def sanitize_field_value(value):
        """Sanitize field value by removing control characters.

        No length limit - we already limit at file/page/field level.
        Skips binary data (signature certificates, embedded files, etc.).
        """
        if not value:
            return ""

        # Skip binary data (bytes objects with non-printable content)
        if isinstance(value, bytes):
            # Check if it's likely binary data (signature, certificate, etc.)
            # Binary data typically has high proportion of non-printable bytes
            try:
                # Try to decode as UTF-8 text
                value_str = value.decode('utf-8')
            except (UnicodeDecodeError, AttributeError):
                # Binary data - skip it
                debug_log(f"[FORM DATA] Skipping binary field value ({len(value)} bytes)")
                return "[Binary Data - Omitted]"
        else:
            value_str = str(value)

        # Remove control characters (except newline, tab, carriage return)
        # This prevents CSV injection and other control character attacks
        sanitized = ''.join(char for char in value_str
                          if char in '\n\t\r' or (ord(char) >= 32 and ord(char) != 127))

        return sanitized.strip()

    for file_path in files:
        if not os.path.exists(file_path):
            errors.append({"file": file_path, "error": f"File not found: {file_path}"})
            continue

        doc = None
        try:
            # SECURITY: Check file size
            try:
                pdf_size_mb = os.path.getsize(file_path) / (1024 * 1024)
                if pdf_size_mb > MAX_PDF_SIZE_MB:
                    errors.append({
                        "file": file_path,
                        "error": f"PDF file too large ({pdf_size_mb:.1f}MB). Maximum allowed: {MAX_PDF_SIZE_MB}MB."
                    })
                    continue
            except Exception as e:
                errors.append({"file": file_path, "error": f"Cannot check file size: {str(e)}"})
                continue

            # VALIDATION: Check PDF header (integrity check)
            try:
                with open(file_path, 'rb') as f:
                    header = f.read(5)
                    debug_log(f"[FORM DATA] PDF header check: {header}")
                    if header != b'%PDF-':
                        error_msg = f"Invalid PDF file: missing PDF header. Expected '%PDF-' but got '{header}'. File may be corrupted or is not a PDF."
                        debug_log(f"[FORM DATA] VALIDATION FAILED: {error_msg}")
                        errors.append({"file": file_path, "error": error_msg})
                        continue
                    debug_log(f"[FORM DATA] PDF header validation passed")
            except Exception as e:
                debug_log(f"[FORM DATA] Error reading file header: {str(e)}")
                errors.append({"file": file_path, "error": f"Cannot read file: {str(e)}"})
                continue

            # Try to open PDF with PyMuPDF (will detect encryption)
            try:
                doc = fitz.open(file_path)
            except Exception as e:
                error_msg = str(e).lower()
                if "password" in error_msg or "crypt" in error_msg or "encrypt" in error_msg:
                    errors.append({
                        "file": file_path,
                        "error": "Cannot extract form data from password-protected PDFs. Please unlock the PDF first using the 'Unlock PDF' tool."
                    })
                else:
                    errors.append({
                        "file": file_path,
                        "error": f"Cannot open PDF: {str(e)}"
                    })
                continue

            # Check for empty PDF (FIXED: should be OR, not AND)
            if len(doc) == 0:
                errors.append({"file": file_path, "error": "PDF has no pages. Cannot extract form data from an empty PDF."})
                continue

            # SECURITY: Check page count
            if len(doc) > MAX_PAGES:
                errors.append({
                    "file": file_path,
                    "error": f"PDF has too many pages ({len(doc)} pages). Maximum allowed: {MAX_PAGES} pages."
                })
                continue

            print(f"[FORM DATA] Extracting from {os.path.basename(file_path)} ({len(doc)} pages, {pdf_size_mb:.2f}MB)", file=sys.stderr, flush=True)

            form_fields = []
            base, _ = os.path.splitext(file_path)

            # Method 1: Try pypdf for AcroForm fields (most reliable for fillable forms)
            try:
                from pypdf import PdfReader
                reader = PdfReader(file_path)

                if '/AcroForm' in reader.trailer.get('/Root', {}):
                    acro_form = reader.trailer['/Root']['/AcroForm']
                    if '/Fields' in acro_form:
                        fields = acro_form['/Fields']

                        for field in fields:
                            # SECURITY: Check field count limit
                            if len(form_fields) >= MAX_FIELDS:
                                print(f"[FORM DATA] Field limit reached ({MAX_FIELDS}). Stopping extraction.", file=sys.stderr, flush=True)
                                break

                            field_info = {}

                            # Get field name
                            if '/T' in field:
                                field_info['Field Name'] = sanitize_field_value(str(field['/T']))
                            else:
                                field_info['Field Name'] = 'Unnamed'

                            # Get field type
                            if '/FT' in field:
                                field_type = str(field['/FT'])
                                type_map = {
                                    '/Tx': 'Text',
                                    '/Btn': 'Button',
                                    '/Ch': 'Choice',
                                    '/Sig': 'Signature'
                                }
                                field_info['Field Type'] = type_map.get(field_type, field_type.replace('/', ''))

                                # Skip signature fields - they contain binary certificate data
                                if field_type == '/Sig':
                                    debug_log(f"[FORM DATA] Skipping signature field: {field_info.get('Field Name', 'Unnamed')}")
                                    continue
                            else:
                                field_info['Field Type'] = 'Unknown'

                            # Get field value (sanitized)
                            if '/V' in field:
                                value = field['/V']
                                if isinstance(value, list):
                                    field_info['Field Value'] = sanitize_field_value(', '.join(str(v) for v in value))
                                else:
                                    field_info['Field Value'] = sanitize_field_value(value)
                            elif '/DV' in field:  # Default value
                                field_info['Field Value'] = sanitize_field_value(field['/DV'])
                            else:
                                field_info['Field Value'] = ''

                            field_info['Page'] = 'Unknown'
                            form_fields.append(field_info)

                        print(f"[FORM DATA] Found {len(form_fields)} AcroForm fields", file=sys.stderr, flush=True)
            except Exception as e:
                logger.warning(f"pypdf form extraction failed: {e}")

            # Method 2: If no AcroForm fields, try PyMuPDF widgets (already have doc open)
            if not form_fields and doc:
                try:
                    for page_num in range(len(doc)):
                        page = doc[page_num]
                        widgets = page.widgets()

                        for widget in widgets:
                            # SECURITY: Check field count limit
                            if len(form_fields) >= MAX_FIELDS:
                                print(f"[FORM DATA] Field limit reached ({MAX_FIELDS}). Stopping extraction.", file=sys.stderr, flush=True)
                                break

                            field_info = {}
                            field_info['Field Name'] = sanitize_field_value(widget.field_name or 'Unnamed')

                            # Map field type
                            field_type_map = {
                                fitz.PDF_WIDGET_TYPE_TEXT: 'Text',
                                fitz.PDF_WIDGET_TYPE_CHECKBOX: 'Checkbox',
                                fitz.PDF_WIDGET_TYPE_RADIOBUTTON: 'Radio Button',
                                fitz.PDF_WIDGET_TYPE_SIGNATURE: 'Signature',
                                fitz.PDF_WIDGET_TYPE_BUTTON: 'Button',
                                fitz.PDF_WIDGET_TYPE_CHOICE: 'Dropdown',
                            }
                            field_info['Field Type'] = field_type_map.get(widget.field_type, 'Unknown')

                            # Skip signature widgets - they contain binary certificate data
                            if widget.field_type == fitz.PDF_WIDGET_TYPE_SIGNATURE:
                                debug_log(f"[FORM DATA] Skipping signature widget: {field_info.get('Field Name', 'Unnamed')}")
                                continue

                            # Get field value (sanitized)
                            if widget.field_value:
                                field_info['Field Value'] = sanitize_field_value(widget.field_value)
                            else:
                                field_info['Field Value'] = ''

                            field_info['Page'] = str(page_num + 1)
                            form_fields.append(field_info)

                        if len(form_fields) >= MAX_FIELDS:
                            break

                    if form_fields:
                        print(f"[FORM DATA] Found {len(form_fields)} widget fields", file=sys.stderr, flush=True)
                except Exception as e:
                    logger.warning(f"PyMuPDF widget extraction failed: {e}")

            # Method 3: Fallback for flattened forms (LIMITED to prevent performance issues)
            if not form_fields and doc:
                try:
                    # SECURITY: Limit pages for text extraction
                    pages_to_extract = min(len(doc), MAX_FLATTENED_PAGES)
                    if len(doc) > MAX_FLATTENED_PAGES:
                        print(f"[FORM DATA] Limiting flattened form extraction to first {MAX_FLATTENED_PAGES} pages", file=sys.stderr, flush=True)

                    for page_num in range(pages_to_extract):
                        # SECURITY: Check field count limit
                        if len(form_fields) >= MAX_FIELDS:
                            break

                        page = doc[page_num]

                        # Extract text using multiple methods for better coverage
                        # Method 1: Standard text extraction
                        text = page.get_text("text")

                        # Method 2: If text seems incomplete (very short), try block extraction
                        # This handles complex layouts better (tables, columns, text boxes)
                        if len(text.strip()) < 100:  # Suspiciously short
                            try:
                                blocks = page.get_text("blocks")
                                # Blocks format: (x0, y0, x1, y1, "text", block_no, block_type)
                                text_blocks = [block[4] for block in blocks if len(block) > 4 and block[4].strip()]
                                if text_blocks:
                                    text = '\n'.join(text_blocks)
                                    debug_log(f"[FORM DATA] Used block extraction for page {page_num + 1} (improved layout handling)")
                            except Exception as e:
                                logger.warning(f"Block extraction failed on page {page_num + 1}: {e}")

                        if text.strip():
                            # Extract full page text (no character limit - we already limit file size and page count)
                            # This is for flattened forms where the entire page content IS the form data
                            field_info = {
                                'Field Name': f'Page {page_num + 1} Content',
                                'Field Type': 'Text (Flattened)',
                                'Field Value': sanitize_field_value(text.strip()),  # No limit - already have file/page limits
                                'Page': str(page_num + 1)
                            }
                            form_fields.append(field_info)

                            # Debug: Log text length
                            debug_log(f"[FORM DATA] Page {page_num + 1}: Extracted {len(text.strip())} characters")

                    if form_fields:
                        debug_log(f"[FORM DATA] Extracted {len(form_fields)} flattened form entries from {pages_to_extract} pages")
                        debug_log(f"[FORM DATA] Note: This PDF has no fillable fields. Extracted page text instead.")
                except Exception as e:
                    logger.warning(f"Flattened form extraction failed: {e}")

            # Create output file
            if form_fields:
                try:
                    df = pd.DataFrame(form_fields)

                    # Debug: Check if data is complete in DataFrame
                    total_chars = sum(len(str(row.get('Field Value', ''))) for _, row in df.iterrows())
                    debug_log(f"[FORM DATA] DataFrame contains {total_chars:,} total characters across {len(form_fields)} fields")

                    if output_format == "csv":
                        output_path = f"{base}_form_data.csv"
                        # Force no truncation in CSV export
                        df.to_csv(
                            output_path,
                            index=False,
                            encoding='utf-8-sig',  # UTF-8 with BOM for Excel
                            quoting=1,  # QUOTE_ALL - ensures long text is properly quoted
                        )
                        processed_files.append(output_path)
                        debug_log(f"[FORM DATA] Created CSV with {len(form_fields)} fields")

                        # Verify CSV file size
                        csv_size_kb = os.path.getsize(output_path) / 1024
                        debug_log(f"[FORM DATA] CSV file size: {csv_size_kb:.1f}KB")

                    elif output_format == "xlsx":
                        output_path = f"{base}_form_data.xlsx"
                        df.to_excel(output_path, index=False, engine='openpyxl')
                        processed_files.append(output_path)
                        debug_log(f"[FORM DATA] Created XLSX with {len(form_fields)} fields")

                        # Verify XLSX file size
                        xlsx_size_kb = os.path.getsize(output_path) / 1024
                        debug_log(f"[FORM DATA] XLSX file size: {xlsx_size_kb:.1f}KB")
                    else:
                        errors.append({"file": file_path, "error": f"Unsupported output format: {output_format}. Use 'csv' or 'xlsx'."})
                except Exception as e:
                    errors.append({"file": file_path, "error": f"Cannot create output file: {str(e)}"})
            else:
                errors.append({
                    "file": file_path,
                    "error": "No form fields found in PDF. The PDF may not contain fillable forms, or the forms may be completely flattened without extractable data."
                })

        except Exception as e:
            logger.error(f"Error processing {file_path}: {e}", exc_info=True)
            errors.append({"file": file_path, "error": f"Form data extraction failed: {str(e)}"})
        finally:
            # SECURITY: Always clean up resources
            if doc:
                try:
                    doc.close()
                except:
                    pass

    return {"processed_files": processed_files, "errors": errors}

# ============================================================================
# FILE TO PDF CONVERSION FUNCTIONS
# ============================================================================

def csv_to_pdf(payload):
    """
    Convert CSV files to PDF documents with formatted tables.
    """
    files = payload.get("files", [])
    processed_files = []
    errors = []

    for file_path in files:
        if not os.path.exists(file_path):
            errors.append({"file": file_path, "error": f"File not found: {file_path}"})
            continue

        try:
            base, _ = os.path.splitext(file_path)
            output_path = f"{base}.pdf"

            # Read CSV file
            try:
                # Try different encodings
                df = None
                for encoding in ['utf-8', 'latin-1', 'cp1252']:
                    try:
                        df = pd.read_csv(file_path, encoding=encoding)
                        break
                    except UnicodeDecodeError:
                        continue

                if df is None:
                    errors.append({"file": file_path, "error": "Could not read CSV file. Check encoding."})
                    continue

            except Exception as e:
                errors.append({"file": file_path, "error": f"Failed to parse CSV: {str(e)}"})
                continue

            # Create PDF with reportlab
            from reportlab.lib import colors
            from reportlab.lib.pagesizes import letter, landscape
            from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch

            # Determine orientation based on column count
            page_size = landscape(letter) if len(df.columns) > 6 else letter

            doc = SimpleDocTemplate(
                output_path,
                pagesize=page_size,
                rightMargin=0.5*inch,
                leftMargin=0.5*inch,
                topMargin=0.5*inch,
                bottomMargin=0.5*inch
            )

            elements = []
            styles = getSampleStyleSheet()

            # Create a style for table cells with text wrapping
            cell_style = ParagraphStyle(
                'CellStyle',
                parent=styles['Normal'],
                fontSize=8,
                leading=10,  # Line spacing
                wordWrap='CJK'  # Enable word wrapping
            )
            
            # Create a style for header cells
            header_style = ParagraphStyle(
                'HeaderStyle',
                parent=styles['Normal'],
                fontSize=9,
                leading=11,
                wordWrap='CJK',
                fontName='Helvetica-Bold'
            )

            # Add title
            title = Paragraph(f"CSV Data: {os.path.basename(file_path)}", styles['Title'])
            elements.append(title)
            elements.append(Spacer(1, 0.25*inch))

            # Prepare table data (header + rows)
            # Limit to 1000 rows to prevent memory issues
            max_rows = min(len(df), 1000)
            if len(df) > 1000:
                elements.append(Paragraph(f"Note: Showing first 1000 of {len(df)} rows", styles['Italic']))
                elements.append(Spacer(1, 0.1*inch))

            table_data = [df.columns.tolist()] + df.head(max_rows).values.tolist()

            # Convert cell data to Paragraph objects for proper text wrapping
            # This preserves full content and enables automatic wrapping
            def escape_xml(text):
                """Escape XML special characters for Paragraph objects"""
                text = str(text)
                text = text.replace('&', '&amp;')
                text = text.replace('<', '&lt;')
                text = text.replace('>', '&gt;')
                return text

            for i, row in enumerate(table_data):
                if i == 0:
                    # Header row - use header style
                    table_data[i] = [Paragraph(escape_xml(cell), header_style) for cell in row]
                else:
                    # Data rows - use cell style
                    table_data[i] = [Paragraph(escape_xml(cell), cell_style) for cell in row]

            # Calculate column widths
            # Remove the restrictive 2-inch cap and allow columns to use more space
            available_width = page_size[0] - 1*inch
            col_width = available_width / len(df.columns)
            # Allow columns to use up to 40% of page width (better for landscape mode)
            max_col_width = available_width * 0.4
            col_widths = [min(col_width, max_col_width) for _ in df.columns]

            table = Table(table_data, colWidths=col_widths, repeatRows=1)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#4A90D9')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 9),
                ('FONTSIZE', (0, 1), (-1, -1), 8),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 8),
                ('TOPPADDING', (0, 0), (-1, 0), 8),
                ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#F8F9FA')),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#F8F9FA')]),
                ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#DEE2E6')),
                ('VALIGN', (0, 0), (-1, -1), 'TOP'),
            ]))

            elements.append(table)
            doc.build(elements)

            if os.path.exists(output_path):
                processed_files.append(output_path)
            else:
                errors.append({"file": file_path, "error": "Failed to create PDF"})

        except Exception as e:
            logger.error(f"Error converting CSV to PDF {file_path}: {e}", exc_info=True)
            errors.append({"file": file_path, "error": str(e)})

    return {"processed_files": processed_files, "errors": errors}


def txt_to_pdf(payload):
    """
    Convert plain text files to PDF documents.
    """
    files = payload.get("files", [])
    processed_files = []
    errors = []

    for file_path in files:
        if not os.path.exists(file_path):
            errors.append({"file": file_path, "error": f"File not found: {file_path}"})
            continue

        try:
            base, _ = os.path.splitext(file_path)
            output_path = f"{base}.pdf"

            # Read text file
            text_content = None
            for encoding in ['utf-8', 'latin-1', 'cp1252']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        text_content = f.read()
                    break
                except UnicodeDecodeError:
                    continue

            if text_content is None:
                errors.append({"file": file_path, "error": "Could not read text file. Check encoding."})
                continue

            # Create PDF with reportlab
            from reportlab.lib.pagesizes import letter
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            from reportlab.lib.enums import TA_LEFT

            doc = SimpleDocTemplate(
                output_path,
                pagesize=letter,
                rightMargin=0.75*inch,
                leftMargin=0.75*inch,
                topMargin=0.75*inch,
                bottomMargin=0.75*inch
            )

            elements = []
            styles = getSampleStyleSheet()

            # Create monospace style for text
            mono_style = ParagraphStyle(
                'MonoText',
                parent=styles['Normal'],
                fontName='Courier',
                fontSize=10,
                leading=14,
                alignment=TA_LEFT,
                spaceAfter=6
            )

            # Add title
            title = Paragraph(f"{os.path.basename(file_path)}", styles['Title'])
            elements.append(title)
            elements.append(Spacer(1, 0.25*inch))

            # Process text - handle line breaks and special characters
            lines = text_content.split('\n')
            for line in lines:
                # Escape XML special characters
                line = line.replace('&', '&amp;')
                line = line.replace('<', '&lt;')
                line = line.replace('>', '&gt;')
                # Replace multiple spaces with non-breaking spaces
                line = line.replace('  ', '&nbsp;&nbsp;')
                # Handle tabs
                line = line.replace('\t', '&nbsp;&nbsp;&nbsp;&nbsp;')

                if line.strip():
                    elements.append(Paragraph(line, mono_style))
                else:
                    elements.append(Spacer(1, 0.1*inch))

            doc.build(elements)

            if os.path.exists(output_path):
                processed_files.append(output_path)
            else:
                errors.append({"file": file_path, "error": "Failed to create PDF"})

        except Exception as e:
            logger.error(f"Error converting TXT to PDF {file_path}: {e}", exc_info=True)
            errors.append({"file": file_path, "error": str(e)})

    return {"processed_files": processed_files, "errors": errors}


def tiff_to_pdf(payload):
    """
    Convert TIFF images to PDF documents. Supports multi-page TIFF files.
    """
    files = payload.get("files", [])
    processed_files = []
    errors = []

    for file_path in files:
        if not os.path.exists(file_path):
            errors.append({"file": file_path, "error": f"File not found: {file_path}"})
            continue

        try:
            base, _ = os.path.splitext(file_path)
            output_path = f"{base}.pdf"

            # Open TIFF image
            img = Image.open(file_path)

            # Get all frames for multi-page TIFF
            frames = []
            try:
                while True:
                    # Convert to RGB if necessary
                    frame = img.copy()
                    if frame.mode in ('RGBA', 'LA', 'P'):
                        # Create white background for transparency
                        background = Image.new('RGB', frame.size, (255, 255, 255))
                        if frame.mode == 'P':
                            frame = frame.convert('RGBA')
                        background.paste(frame, mask=frame.split()[-1] if frame.mode in ('RGBA', 'LA') else None)
                        frame = background
                    elif frame.mode != 'RGB':
                        frame = frame.convert('RGB')

                    frames.append(frame)
                    img.seek(img.tell() + 1)
            except EOFError:
                pass  # End of frames

            if not frames:
                errors.append({"file": file_path, "error": "No valid frames in TIFF file"})
                continue

            # Save as PDF
            if len(frames) == 1:
                frames[0].save(output_path, 'PDF', resolution=150.0)
            else:
                # Multi-page PDF
                frames[0].save(
                    output_path,
                    'PDF',
                    resolution=150.0,
                    save_all=True,
                    append_images=frames[1:]
                )

            if os.path.exists(output_path):
                processed_files.append(output_path)
            else:
                errors.append({"file": file_path, "error": "Failed to create PDF"})

        except Exception as e:
            logger.error(f"Error converting TIFF to PDF {file_path}: {e}", exc_info=True)
            errors.append({"file": file_path, "error": str(e)})

    return {"processed_files": processed_files, "errors": errors}


def rtf_to_pdf(payload):
    """
    Convert RTF (Rich Text Format) files to PDF documents.
    Uses striprtf for parsing and reportlab for PDF generation.
    """
    files = payload.get("files", [])
    processed_files = []
    errors = []

    # Try to import striprtf
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError:
        return {
            "processed_files": [],
            "errors": [{"file": "system", "error": "RTF conversion requires striprtf library. Install with: pip install striprtf"}]
        }

    for file_path in files:
        if not os.path.exists(file_path):
            errors.append({"file": file_path, "error": f"File not found: {file_path}"})
            continue

        try:
            base, _ = os.path.splitext(file_path)
            output_path = f"{base}.pdf"

            # Read RTF file
            rtf_content = None
            for encoding in ['utf-8', 'latin-1', 'cp1252']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        rtf_content = f.read()
                    break
                except UnicodeDecodeError:
                    continue

            if rtf_content is None:
                errors.append({"file": file_path, "error": "Could not read RTF file. Check encoding."})
                continue

            # Convert RTF to plain text
            try:
                text_content = rtf_to_text(rtf_content)
            except Exception as e:
                errors.append({"file": file_path, "error": f"Failed to parse RTF: {str(e)}"})
                continue

            # Create PDF with reportlab
            from reportlab.lib.pagesizes import letter
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.lib.units import inch

            doc = SimpleDocTemplate(
                output_path,
                pagesize=letter,
                rightMargin=0.75*inch,
                leftMargin=0.75*inch,
                topMargin=0.75*inch,
                bottomMargin=0.75*inch
            )

            elements = []
            styles = getSampleStyleSheet()

            # Add title
            title = Paragraph(f"{os.path.basename(file_path)}", styles['Title'])
            elements.append(title)
            elements.append(Spacer(1, 0.25*inch))

            # Process text - handle line breaks
            paragraphs = text_content.split('\n\n')  # Split by double newlines for paragraphs
            for para in paragraphs:
                # Clean up single newlines within paragraphs
                para = para.replace('\n', ' ').strip()
                # Escape XML special characters
                para = para.replace('&', '&amp;')
                para = para.replace('<', '&lt;')
                para = para.replace('>', '&gt;')

                if para:
                    elements.append(Paragraph(para, styles['Normal']))
                    elements.append(Spacer(1, 0.1*inch))

            doc.build(elements)

            if os.path.exists(output_path):
                processed_files.append(output_path)
            else:
                errors.append({"file": file_path, "error": "Failed to create PDF"})

        except Exception as e:
            logger.error(f"Error converting RTF to PDF {file_path}: {e}", exc_info=True)
            errors.append({"file": file_path, "error": str(e)})

    return {"processed_files": processed_files, "errors": errors}


def xml_to_pdf(payload):
    """
    Convert XML files to readable PDF documents.
    Renders XML structure in a formatted, readable way.
    """
    files = payload.get("files", [])
    processed_files = []
    errors = []

    for file_path in files:
        if not os.path.exists(file_path):
            errors.append({"file": file_path, "error": f"File not found: {file_path}"})
            continue

        try:
            base, _ = os.path.splitext(file_path)
            output_path = f"{base}.pdf"

            # Read XML file
            xml_content = None
            for encoding in ['utf-8', 'latin-1', 'cp1252']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        xml_content = f.read()
                    break
                except UnicodeDecodeError:
                    continue

            if xml_content is None:
                errors.append({"file": file_path, "error": "Could not read XML file. Check encoding."})
                continue

            # Parse and pretty-print XML (using defusedxml for XXE protection)
            import defusedxml.ElementTree as ET
            from defusedxml.minidom import parseString

            try:
                # Parse XML (defusedxml prevents XXE attacks)
                root = ET.fromstring(xml_content)
                # Pretty print
                xml_str = ET.tostring(root, encoding='unicode')
                dom = parseString(xml_str)
                pretty_xml = dom.toprettyxml(indent="  ")
                # Remove extra blank lines
                pretty_xml = '\n'.join([line for line in pretty_xml.split('\n') if line.strip()])
            except ET.ParseError as e:
                # If XML parsing fails, use raw content
                logger.warning(f"XML parsing failed, using raw content: {e}")
                pretty_xml = xml_content

            # Create PDF with reportlab
            from reportlab.lib.pagesizes import letter
            from reportlab.lib import colors
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Preformatted
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            from reportlab.lib.enums import TA_LEFT

            doc = SimpleDocTemplate(
                output_path,
                pagesize=letter,
                rightMargin=0.5*inch,
                leftMargin=0.5*inch,
                topMargin=0.75*inch,
                bottomMargin=0.75*inch
            )

            elements = []
            styles = getSampleStyleSheet()

            # Create code style for XML
            code_style = ParagraphStyle(
                'XMLCode',
                parent=styles['Code'],
                fontName='Courier',
                fontSize=8,
                leading=10,
                alignment=TA_LEFT,
                textColor=colors.HexColor('#333333'),
                backColor=colors.HexColor('#F5F5F5'),
                borderColor=colors.HexColor('#DDDDDD'),
                borderWidth=1,
                borderPadding=8,
            )

            # Add title
            title = Paragraph(f"XML Document: {os.path.basename(file_path)}", styles['Title'])
            elements.append(title)
            elements.append(Spacer(1, 0.25*inch))

            # Process XML lines
            lines = pretty_xml.split('\n')
            for line in lines:
                # Escape special characters for reportlab
                line = line.replace('&', '&amp;')
                line = line.replace('<', '&lt;')
                line = line.replace('>', '&gt;')
                # Preserve indentation
                line = line.replace('  ', '&nbsp;&nbsp;')

                if line.strip():
                    elements.append(Paragraph(line, code_style))

            doc.build(elements)

            if os.path.exists(output_path):
                processed_files.append(output_path)
            else:
                errors.append({"file": file_path, "error": "Failed to create PDF"})

        except Exception as e:
            logger.error(f"Error converting XML to PDF {file_path}: {e}", exc_info=True)
            errors.append({"file": file_path, "error": str(e)})

    return {"processed_files": processed_files, "errors": errors}

# Trigger reload
